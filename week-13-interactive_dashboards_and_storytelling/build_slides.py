"""Build Week 13 lecture slides — Interactive Dashboards & Data Storytelling.

Visual style: ACL@NCU — white BG, navy + teal accents.

設計準則 (本次修訂)：
- 受眾是行為科學研究生，非資訊背景 → 每個英文術語首次出現都附中文解釋
- 用日常類比 (Google 下一頁、Russian doll、體檢) 取代 buzzword
- 環境是 Windows PowerShell + ipython + VS Code → demo callouts 給 PS 指令與 VS Code 操作
- 每個概念有 .py (整支跑) 與 .ipynb (cell-by-cell 跑) 兩種使用方式

Output: week-13-slides.pptx (16:9).
Run:    python build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# ACL@NCU palette
# ============================================================
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE    = RGBColor(0xF7, 0xF8, 0xFA)
BG_SECTION     = RGBColor(0x14, 0x32, 0x5C)
BG_BREAK       = RGBColor(0xFB, 0xEA, 0xC0)
BG_PRACTICE    = RGBColor(0xEC, 0xF7, 0xF6)
TEXT_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED     = RGBColor(0x5F, 0x6B, 0x83)
ACCENT_PRIMARY = RGBColor(0x14, 0x32, 0x5C)
ACCENT_TEAL    = RGBColor(0x0D, 0x9B, 0x9B)
ACCENT_AMBER   = RGBColor(0xE8, 0xA1, 0x2A)
ACCENT_RED     = RGBColor(0xD3, 0x4F, 0x4F)
ACCENT_GREEN   = RGBColor(0x2E, 0x8B, 0x57)
ACCENT_ORANGE  = RGBColor(0xF9, 0x71, 0x16)
CODE_BG        = RGBColor(0x1E, 0x29, 0x3B)
CODE_COMMENT   = RGBColor(0x8B, 0x9D, 0xB8)
HAIRLINE       = RGBColor(0xE2, 0xE6, 0xEC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"

FOOTER_STR = "Week 13 — Interactive Dashboards & Data Storytelling  ·  ACL@NCU"


# ============================================================
# Primitives
# ============================================================
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_top_band(slide, color=ACCENT_TEAL, height=Inches(0.14)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=20, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_section_label(slide, text, color=ACCENT_TEAL):
    add_text(slide, text,
             Inches(0.7), Inches(0.55), Inches(11), Inches(0.35),
             size=14, bold=True, color=color, font=FONT_SANS)


def add_title(slide, title_zh, *, size=32, top=0.95, color=TEXT_DARK):
    add_text(slide, title_zh,
             Inches(0.7), Inches(top), Inches(12), Inches(0.85),
             size=size, bold=True, color=color, font=FONT_CJK)


def add_bullets(slide, items, left, top, width, height, *,
                size=18, color=TEXT_DARK, line_spacing=1.25,
                bullet_color=ACCENT_TEAL, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    use_font = font or FONT_CJK
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = "▸  "
        b.font.name = FONT_SANS
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = use_font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i > 0:
            p.space_before = Pt(size * 0.45)
    return tb


def _color_for_code_line(line, lang):
    s = line.lstrip()
    if not s:
        return None
    if lang in ("bash", "powershell", "yaml", "python", "markdown") and s.startswith("#"):
        return CODE_COMMENT
    if s.startswith(">"):
        return ACCENT_AMBER
    if s.startswith("$") or s.startswith("PS "):
        return ACCENT_GREEN
    return None


def add_code(slide, code, left, top, width, height, *,
             size=14, lang="python", line_spacing=1.15, padding=0.18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(padding); tf.margin_bottom = Inches(padding)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        c = _color_for_code_line(line, lang) or TEXT_LIGHT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(size)
        r.font.color.rgb = c
    return box


def add_card(slide, left, top, width, height, *, fill=BG_OFFWHITE,
             accent=None, accent_w=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = HAIRLINE
    card.line.width = Pt(0.5)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(accent_w), height)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return card


def add_callout(slide, text, left, top, width, height, *,
                size=24, color=ACCENT_PRIMARY, fill=BG_OFFWHITE,
                accent=ACCENT_TEAL, font=FONT_CJK, align=PP_ALIGN.LEFT):
    add_card(slide, left, top, width, height, fill=fill, accent=accent,
             accent_w=0.10)
    add_text(slide, text,
             left + Inches(0.4), top + Inches(0.2),
             width - Inches(0.6), height - Inches(0.4),
             size=size, bold=True, color=color, font=font,
             align=align, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page, total, *, theme="light"):
    color = TEXT_MUTED if theme == "light" else ACCENT_TEAL
    add_text(slide, FOOTER_STR,
             Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
             size=11, color=color)
    add_text(slide, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=color, align=PP_ALIGN.RIGHT)


def add_demo_callout(slide, label, path, *, left=0.85, top=6.40,
                     width=11.6, height=0.5, accent=ACCENT_AMBER):
    add_card(slide, Inches(left), Inches(top), Inches(width), Inches(height),
             fill=BG_OFFWHITE, accent=accent, accent_w=0.08)
    add_text(slide, "▶ DEMO",
             Inches(left + 0.18), Inches(top + 0.10),
             Inches(1.05), Inches(height - 0.15),
             size=13, bold=True, color=accent)
    add_text(slide, label,
             Inches(left + 1.30), Inches(top + 0.10),
             Inches(5.2), Inches(height - 0.15),
             size=13, color=TEXT_DARK, font=FONT_CJK)
    add_text(slide, path,
             Inches(left + 6.55), Inches(top + 0.10),
             Inches(width - 6.7), Inches(height - 0.15),
             size=12, color=ACCENT_PRIMARY, font=FONT_MONO)


def add_table(slide, headers, rows, left, top, width, height, *,
              header_fill=ACCENT_PRIMARY, header_color=TEXT_LIGHT,
              col_widths=None, header_size=13, body_size=12,
              first_col_bold=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w

    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_color

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT_CJK
            r.font.size = Pt(body_size)
            r.font.color.rgb = TEXT_DARK
            if first_col_bold and ci == 0:
                r.font.bold = True
    return table_shape


def add_slide_base(prs, *, bg=BG_WHITE, band=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, bg)
    if band:
        add_top_band(s)
    return s


# ============================================================
# Mini chart primitives
# ============================================================
def add_mini_line_chart(slide, left, top, width, height, values,
                        *, color=ACCENT_TEAL, marker=True, annotation=None):
    n = len(values)
    pad = Inches(0.10)
    inner_l = left + pad
    inner_b = top + height - pad
    plot_w = width - 2 * pad
    plot_h = height - 2 * pad
    ax = slide.shapes.add_connector(1, inner_l, inner_b,
                                    inner_l + plot_w, inner_b)
    ax.line.color.rgb = TEXT_MUTED
    ax.line.width = Pt(0.75)
    xs = [inner_l + i * (plot_w / max(1, n - 1)) for i in range(n)]
    ys = [inner_b - plot_h * max(0.04, v) for v in values]
    for i in range(n - 1):
        c = slide.shapes.add_connector(1, xs[i], ys[i], xs[i + 1], ys[i + 1])
        c.line.color.rgb = color
        c.line.width = Pt(2.0)
    if marker:
        r = Inches(0.07)
        for x, y in zip(xs, ys):
            m = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       x - r/2, y - r/2, r, r)
            m.fill.solid(); m.fill.fore_color.rgb = color
            m.line.fill.background()


# ============================================================
# Slide builders
# ============================================================
def slide_title(prs):
    s = add_slide_base(prs, band=False)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))

    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(2.5),
                               Inches(0.18), Inches(2.6))
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT_TEAL
    block.line.fill.background()

    add_text(s, "WEEK 13",
             Inches(1.1), Inches(2.4), Inches(11), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, "互動式儀表板",
             Inches(1.1), Inches(2.95), Inches(12), Inches(0.95),
             size=42, bold=True, color=TEXT_DARK, font=FONT_CJK)
    add_text(s, "與資料敘事",
             Inches(1.1), Inches(3.85), Inches(12), Inches(0.85),
             size=38, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, "用兩個真實大型資料集走完整 pipeline",
             Inches(1.1), Inches(4.85), Inches(12), Inches(0.6),
             size=20, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "NS5116 · Programming & AI Applications in Behavioral Science",
             Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
             size=14, color=TEXT_MUTED)
    add_text(s, "ACL@NCU  ·  Spring 2026  ·  2026-05-21  ·  張智宏",
             Inches(1.1), Inches(6.4), Inches(11), Inches(0.4),
             size=13, color=TEXT_MUTED)
    return s


def slide_why(prs, page, total):
    """為什麼今天要學互動圖表 — 直白版。"""
    s = add_slide_base(prs)
    add_section_label(s, "本週動機")
    add_title(s, "好的儀表板不是把所有圖塞進一頁")

    items = [
        "一張好的圖應該回答一個明確的問題",
        "讀者既要能自由探索，也要被引導看到結論",
        "互動性 (可 hover、可放大、可 filter) 是網頁圖表的關鍵",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.0),
                size=20)

    cards = [
        ("✗  雜亂的圖表牆",
         "所有變數都塞一張圖；沒有先後順序；讀者迷路",
         ACCENT_RED),
        ("✓  有順序的儀表板",
         "從整體到細節；一張圖一個結論；annotation 引導視線",
         ACCENT_GREEN),
    ]
    x0, y0 = Inches(0.9), Inches(4.3)
    w, h = Inches(5.7), Inches(2.0)
    gap = Inches(0.35)
    for i, (title, body, color) in enumerate(cards):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=color)
        add_text(s, title, left + Inches(0.25), y0 + Inches(0.2),
                 w - Inches(0.4), Inches(0.5),
                 size=20, bold=True, color=color, font=FONT_CJK)
        add_text(s, body, left + Inches(0.25), y0 + Inches(0.85),
                 w - Inches(0.4), h - Inches(1.0),
                 size=15, color=TEXT_DARK, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_two_datasets(prs, page, total):
    """兩個 dataset，同一條流程。"""
    s = add_slide_base(prs)
    add_section_label(s, "今天的兩個資料集")
    add_title(s, "差別只在「資料怎麼取得」")

    add_text(s,
             "Week 12 學的「載入→清理→描述→分析」流程，今天兩份真實資料各跑一遍",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
             size=16, color=TEXT_MUTED, font=FONT_CJK)

    cards = [
        ("資料 A · PsyArXiv 心理學論文",
         "OSF API · 即時 REST 介面",
         "約 1,000 篇最新心理學 preprint",
         "技巧：分頁抓取 + 把 JSON 攤平成表格",
         ACCENT_TEAL),
        ("資料 B · 教育部高教統計",
         "stats.moe.gov.tw · 公開 CSV",
         "105–113 學年度共約 7,200 列",
         "技巧：跨年度合併 + 中文欄位清理",
         ACCENT_ORANGE),
    ]
    x0, y0 = Inches(0.9), Inches(2.7)
    w, h = Inches(5.7), Inches(3.5)
    gap = Inches(0.35)
    for i, (title, src, scale, skill, color) in enumerate(cards):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=color, accent_w=0.10)
        add_text(s, title, left + Inches(0.3), y0 + Inches(0.25),
                 w - Inches(0.5), Inches(0.55),
                 size=20, bold=True, color=color, font=FONT_CJK)
        add_text(s, src, left + Inches(0.3), y0 + Inches(0.95),
                 w - Inches(0.5), Inches(0.4),
                 size=14, color=TEXT_MUTED, font=FONT_MONO)
        add_text(s, scale, left + Inches(0.3), y0 + Inches(1.55),
                 w - Inches(0.5), Inches(0.55),
                 size=17, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, "重點技巧", left + Inches(0.3), y0 + Inches(2.25),
                 w - Inches(0.5), Inches(0.35),
                 size=12, bold=True, color=TEXT_MUTED, font=FONT_CJK)
        add_text(s, skill, left + Inches(0.3), y0 + Inches(2.65),
                 w - Inches(0.5), Inches(0.6),
                 size=15, bold=True, color=color, font=FONT_CJK)

    add_callout(s,
        "資料載進來之後，清理、分析、畫圖的程式碼幾乎一模一樣。",
        Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.45),
        size=15, color=ACCENT_PRIMARY)

    add_footer(s, page, total)
    return s


def slide_objectives(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "今天的學習目標")
    add_title(s, "下課後你會⋯")

    items = [
        "從網路 API 用「分頁」方式抓取上千筆資料",
        "把巢狀的 JSON 攤平成一張表格 (DataFrame)",
        "用 pd.concat 合併多年度的 CSV，處理欄位不一致",
        "用 Plotly 畫互動式的長條圖、折線圖、散布圖",
        "把圖表嵌入 Streamlit 網頁",
        "用 annotation 與圖表說明，告訴讀者「為什麼這張圖重要」",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5),
                size=19)

    add_footer(s, page, total)
    return s


def slide_dev_env(prs, page, total):
    """新增：今天在 Windows + PowerShell + VS Code 怎麼用 demo。"""
    s = add_slide_base(prs)
    add_section_label(s, "今天的開發環境")
    add_title(s, "在 Windows 怎麼跑今天的 demo？")

    # 兩個 column: PowerShell vs VS Code
    add_text(s, "方式 1 — PowerShell 一次跑完整支",
             Inches(0.9), Inches(2.0), Inches(6.0), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    code1 = (
        "# 切到 code 資料夾\n"
        "PS> cd .\\week-13-...\\code\n"
        "\n"
        "# 跑整支腳本，看完整結果\n"
        "PS> python .\\osf_psyarxiv_pipeline.py\n"
        "PS> python .\\moe_higher_ed_pipeline.py"
    )
    add_code(s, code1, Inches(0.9), Inches(2.5), Inches(5.9), Inches(2.4),
             size=12, lang="powershell")

    add_text(s, "方式 2 — VS Code 一格一格跑",
             Inches(7.1), Inches(2.0), Inches(6.0), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_bullets(s, [
        "打開 .ipynb 檔（同個資料夾）",
        "右上角選 Python 直譯器 (Python 3.10+)",
        "Shift+Enter 一格一格跑，看每段輸出",
        "適合上課跟著老師一步一步走",
    ], Inches(7.1), Inches(2.5), Inches(5.9), Inches(2.4), size=14)

    # 第三 column - 終端機選擇
    add_text(s, "方式 3 — IPython 互動式 (REPL)",
             Inches(0.9), Inches(5.05), Inches(11.5), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    code3 = (
        "PS> ipython                       # 啟動 IPython\n"
        "In [1]: import pandas as pd       # 一邊輸入一邊看結果\n"
        "In [2]: df = pd.read_csv(\"moe_higher_ed.csv\")\n"
        "In [3]: df.head()                 # 立即看表格"
    )
    add_code(s, code3, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.5),
             size=12, lang="powershell")

    add_footer(s, page, total)
    return s


def slide_plotly_role(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 認識 PLOTLY")
    add_title(s, "為什麼這週要學 Plotly？")

    add_text(s, "Matplotlib 我們已經會用了 — 但它畫的是 *靜態* 圖（給論文 PDF 用）",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "Plotly 畫的是「會動的圖」：滑鼠移過去看數值、可放大、按 legend 篩選",
             Inches(0.9), Inches(2.45), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)

    rows = [
        ("論文 figure",                "推薦",        "PDF 流程不友善"),
        ("資料探索 (邊看邊調)",         "沒有 hover",  "移過去就有數值"),
        ("放到 Streamlit 網頁",        "只能存圖",    "直接互動"),
        ("Jupyter / VS Code notebook", "可以",        "可以"),
    ]
    add_table(s, ["要做什麼", "Matplotlib", "Plotly"], rows,
              Inches(0.9), Inches(3.1), Inches(11.5), Inches(2.4),
              col_widths=[Inches(5.0), Inches(3.2), Inches(3.3)],
              header_size=14, body_size=14, first_col_bold=True)

    add_callout(s,
        "Plotly 不取代 Matplotlib。論文還是 Matplotlib，網頁儀表板用 Plotly。",
        Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.7),
        size=16, color=ACCENT_PRIMARY)

    add_footer(s, page, total)
    return s


def slide_plotly_basics(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · PLOTLY EXPRESS")
    add_title(s, "Plotly Express：一行函式畫一張圖")

    code = (
        "import plotly.express as px\n"
        "import streamlit as st\n"
        "\n"
        "# 一個 px.bar 呼叫 = 一張完整 bar chart\n"
        "fig = px.bar(df, x=\"city\", y=\"n_students\",\n"
        "             color=\"sector\", barmode=\"group\")\n"
        "\n"
        "# 把這個 figure 物件直接交給 streamlit\n"
        "st.plotly_chart(fig, use_container_width=True)"
    )
    add_code(s, code, Inches(0.85), Inches(2.0), Inches(7.4), Inches(3.5),
             size=14)

    add_text(s, "今天會用到的三個函式",
             Inches(8.55), Inches(2.0), Inches(4.5), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_bullets(s, [
        "px.bar()    類別比較",
        "px.line()   時間趨勢",
        "px.scatter() 兩變數關係",
    ], Inches(8.55), Inches(2.55), Inches(4.5), Inches(1.6),
       size=15, font=FONT_MONO)

    add_text(s, "嵌進 Streamlit 的要點",
             Inches(8.55), Inches(4.3), Inches(4.5), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_bullets(s, [
        "不用先呼叫 .show()",
        "use_container_width=True 讓圖隨欄寬縮放",
    ], Inches(8.55), Inches(4.85), Inches(4.5), Inches(1.5),
       size=14)

    add_demo_callout(s, "在 IPython 試一下", "import plotly.express as px",
                     accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


# ----------------------------------------------------------------
# Dataset A
# ----------------------------------------------------------------
def slide_psyarxiv_intro(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · PSYARXIV")
    add_title(s, "心理學界最近在討論什麼？")

    add_text(s, "PsyArXiv — 心理學界主要的論文預印本網站",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
             size=18, bold=True, color=ACCENT_PRIMARY, font=FONT_CJK)
    add_text(s, "Preprint (預印本) = 還沒正式審查的論文，但作者把它公開讓大家看",
             Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.4),
             size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)

    add_bullets(s, [
        "由 COS (開放科學中心) 維運，每天有數十篇新論文上傳",
        "用 REST API 取得 — 等於「網址 = 問題，回傳 JSON = 答案」",
        "每篇都有標題、tags (作者自填關鍵字)、subjects (主題分類)、日期",
        "正好用來練：API 抓取、分頁、把 JSON 攤平成表格",
    ], Inches(0.9), Inches(3.0), Inches(11.5), Inches(2.0), size=15)

    # sample metadata
    add_card(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.45),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "拿到一筆 preprint 大概長這樣", Inches(1.1), Inches(5.15),
             Inches(8), Inches(0.4),
             size=13, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s,
             "title          : \"Future self-continuity & intergenerational concern\"\n"
             "date_published : 2026-05-13\n"
             "primary_subject: Cognitive Neuroscience\n"
             "n_tags         : 5",
             Inches(1.1), Inches(5.5), Inches(11.0), Inches(0.95),
             size=12, color=TEXT_DARK, font=FONT_MONO)

    add_demo_callout(s, "VS Code 打開 .ipynb 跟著做",
                     "code/osf_psyarxiv_pipeline.ipynb")
    add_footer(s, page, total)
    return s


def slide_pagination(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 分頁 (PAGINATION)")
    add_title(s, "為什麼不能只抓第一頁就好？")

    add_text(s,
             "「分頁」(pagination) — 想像成 Google 搜尋每頁只給你 10 個結果，要看更多得按下一頁。",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_bullets(s, [
        "OSF API 一次最多回 100 筆 — 抓 1,000 筆要做 10 次呼叫",
        "for 迴圈跑 page=1,2,...,10，把每次的結果累積起來",
        "每次呼叫之間 time.sleep(0.3) — 對 server 禮貌，避免被封鎖",
        "r.raise_for_status() — 如果伺服器回錯誤，立刻停下來，不要繼續",
    ], Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.1), size=15)

    code = (
        "for page in range(1, n_pages + 1):\n"
        "    params = {\"page\": page, \"page[size]\": 100,\n"
        "              \"sort\": \"-date_published\"}\n"
        "    r = requests.get(OSF_ENDPOINT, params=params, timeout=30)\n"
        "    r.raise_for_status()\n"
        "    all_items.extend(r.json().get(\"data\", []))\n"
        "    time.sleep(0.3)"
    )
    add_code(s, code, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.65),
             size=12)

    add_demo_callout(s, "在 .ipynb 第 3 格看 pagination 跑動",
                     "fetch_psyarxiv(n_pages=10)")
    add_footer(s, page, total)
    return s


def slide_json_to_df(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 把 JSON 變成表格")
    add_title(s, "巢狀 JSON → 一列一篇 preprint")

    add_text(s,
             "JSON 是巢狀的 — 像 Russian doll 套娃娃。要分析就得攤平成「一列代表一筆」。",
             Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s, "原始巢狀 JSON",
             Inches(0.9), Inches(2.55), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=ACCENT_RED, font=FONT_CJK)
    json_text = (
        "{\n"
        "  \"attributes\": {\n"
        "    \"title\": \"...\",\n"
        "    \"date_published\": \"2026-05-13\",\n"
        "    \"tags\": [\"longitudinal\", ...],\n"
        "    \"subjects\": [\n"
        "      [\n"
        "        {\"text\": \"Social Sciences\"},\n"
        "        {\"text\": \"Cog. Neuroscience\"}\n"
        "      ]\n"
        "    ]\n"
        "  }\n"
        "}"
    )
    add_code(s, json_text, Inches(0.9), Inches(3.0), Inches(5.7), Inches(3.3),
             size=12, lang="markdown")

    add_text(s, "整理過的 DataFrame (一列一篇)",
             Inches(6.95), Inches(2.55), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_table(s,
              ["title", "date", "subject", "tags"],
              [
                  ("Future self-continuity ...", "2026-05-13", "Cog. Neurosci.", "5"),
                  ("Mind-wandering and ...",     "2026-05-12", "Cog. Psy.",      "3"),
                  ("Meta-analytic methods ...",  "2026-05-12", "Meta-science",   "7"),
              ],
              Inches(6.95), Inches(3.0), Inches(5.7), Inches(2.5),
              col_widths=[Inches(2.0), Inches(1.0), Inches(1.9), Inches(0.8)],
              header_size=11, body_size=10)
    add_text(s, "重點：subjects 是分類路徑，我們只留最具體那一層",
             Inches(6.95), Inches(5.7), Inches(5.7), Inches(0.6),
             size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_cleaning_decisions(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 清理資料的紀律")
    add_title(s, "每個動作都要說得出代價")

    add_text(s,
             "Week 12 學的紀律：清理資料時，每個動作對應一個觀察，並交代代價。",
             Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.45),
             size=15, color=TEXT_MUTED, font=FONT_CJK)

    rows = [
        ("date_published 是字串",
         "用 pd.to_datetime 轉成時間型別",
         "無法解析者變 NaT（plot 會略過）"),
        ("有些 title 是空字串",
         "篩掉 len(title)==0 的列",
         "失去 <1% 資料，避免污染統計"),
        ("primary_subject 有缺值",
         "fillna('Unspecified')",
         "可能掩蓋編碼問題，但保留 n"),
    ]
    add_table(s,
              ["觀察到什麼", "做了什麼動作", "動作的代價"],
              rows,
              Inches(0.85), Inches(2.6), Inches(11.6), Inches(2.6),
              col_widths=[Inches(3.7), Inches(3.6), Inches(4.3)],
              header_size=13, body_size=13)

    add_callout(s,
        "清理函式寫成獨立的 clean(df) — 沒有副作用，可以單獨測試。",
        Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.8),
        size=16, color=ACCENT_PRIMARY)

    add_demo_callout(s, "在 .ipynb 第 5 格看清理過程",
                     "code/osf_psyarxiv_pipeline.ipynb")
    add_footer(s, page, total)
    return s


def slide_psyarxiv_bar(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 圖 1 · 長條圖")
    add_title(s, "最近一千篇，最熱門的主題是哪些？")

    add_bullets(s, [
        "用水平 bar (orientation='h') — subject 名稱很長，放 y 軸不會擠",
        "color='n' 把數字本身對應到顏色，一眼看出高低",
        "categoryorder='total ascending' 自動由大到小排序",
    ], Inches(0.9), Inches(1.95), Inches(6.5), Inches(2.0), size=15)

    code = (
        "fig = px.bar(\n"
        "    counts, x=\"n\", y=\"subject\",\n"
        "    orientation=\"h\",\n"
        "    color=\"n\",\n"
        "    color_continuous_scale=\"Blues\",\n"
        "    title=\"Top 15 PsyArXiv subjects\",\n"
        ")"
    )
    add_code(s, code, Inches(0.9), Inches(4.0), Inches(6.4), Inches(2.2),
             size=13)

    add_card(s, Inches(7.7), Inches(1.95), Inches(4.8), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "預覽 · 前 8 大 subject",
             Inches(7.9), Inches(2.10), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    subjects = [
        ("Social & Behavioral", 23),
        ("Psychiatry",          22),
        ("Meta-science",        18),
        ("Cognitive Neurosci.", 14),
        ("Clinical Psychology", 14),
        ("Cognitive Psy.",      13),
        ("Developmental",       12),
        ("Neuroscience",        11),
    ]
    max_v = max(v for _, v in subjects)
    y = 2.55
    for label, v in subjects:
        add_text(s, label, Inches(7.9), Inches(y),
                 Inches(2.0), Inches(0.40),
                 size=11, color=TEXT_DARK, font=FONT_CJK)
        bar_w = Inches(2.4 * v / max_v)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(9.95), Inches(y + 0.08),
                                 bar_w, Inches(0.22))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_TEAL
        bar.line.fill.background()
        add_text(s, str(v), Inches(9.95) + bar_w + Inches(0.05),
                 Inches(y), Inches(0.6), Inches(0.40),
                 size=11, color=TEXT_DARK)
        y += 0.45

    add_footer(s, page, total)
    return s


def slide_psyarxiv_line(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 圖 2 · 折線 + 標注")
    add_title(s, "發表數量是不是有月份差異？")

    add_bullets(s, [
        "markers=True 同時畫點與線",
        "hovermode='x unified' — 滑鼠移過去，同一月份所有數字一起顯示",
        "add_annotation — 在最高點畫箭頭並加文字解釋為什麼這點重要",
    ], Inches(0.9), Inches(1.95), Inches(6.5), Inches(2.0), size=15)

    code = (
        "monthly = (df.groupby(\"year_month\").size()\n"
        "             .reset_index(name=\"n_preprints\"))\n"
        "fig = px.line(monthly, x=\"year_month\",\n"
        "              y=\"n_preprints\", markers=True)\n"
        "\n"
        "peak = monthly.loc[monthly.n_preprints.idxmax()]\n"
        "fig.add_annotation(\n"
        "    x=peak.year_month, y=peak.n_preprints,\n"
        "    text=\"Peak ...\", arrowhead=2)"
    )
    add_code(s, code, Inches(0.9), Inches(4.05), Inches(6.4), Inches(2.25),
             size=12)

    add_card(s, Inches(7.7), Inches(1.95), Inches(4.8), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "預覽 · 每月發表數量",
             Inches(7.9), Inches(2.10), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    values = [0.35, 0.55, 0.90, 0.65, 0.75, 0.60]
    add_mini_line_chart(s, Inches(7.9), Inches(2.55),
                        Inches(4.4), Inches(3.6),
                        values)
    # peak annotation 位置 (對應 values[2] 的最高點)
    add_text(s, "◀ peak",
             Inches(9.5), Inches(2.6), Inches(1.2), Inches(0.3),
             size=10, bold=True, color=ACCENT_RED)

    add_footer(s, page, total)
    return s


def slide_psyarxiv_scatter(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 A · 圖 3 · 散布圖")
    add_title(s, "標題越長的論文，是不是 tag 也越多？")

    add_bullets(s, [
        "散布圖 (scatter) 是探索資料時最常用的圖 — 看出 outlier 與群聚",
        "hover_name='title' 滑鼠移過去就看到具體論文標題",
        "opacity=0.6 半透明 — 點堆疊時可以看出密度",
        "subject 太多時 showlegend=False，避免 legend 擋畫面",
    ], Inches(0.9), Inches(1.95), Inches(6.5), Inches(2.3), size=15)

    code = (
        "fig = px.scatter(\n"
        "    df, x=\"n_tags\", y=\"title_len\",\n"
        "    color=\"primary_subject\",\n"
        "    hover_name=\"title\",\n"
        "    opacity=0.6,\n"
        ")\n"
        "fig.update_layout(showlegend=False)"
    )
    add_code(s, code, Inches(0.9), Inches(4.4), Inches(6.4), Inches(1.95),
             size=13)

    # Right scatter sketch
    add_card(s, Inches(7.7), Inches(1.95), Inches(4.8), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "預覽 · title 長度 vs. tag 數",
             Inches(7.9), Inches(2.10), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    left, top, w, h = Inches(8.0), Inches(2.55), Inches(4.2), Inches(3.4)
    pad = Inches(0.20)
    x0, y0 = left + pad, top + h - pad
    ax_x = s.shapes.add_connector(1, x0, y0, left + w - pad, y0)
    ax_x.line.color.rgb = TEXT_MUTED
    ax_x.line.width = Pt(0.75)
    ax_y = s.shapes.add_connector(1, x0, y0, x0, top + pad)
    ax_y.line.color.rgb = TEXT_MUTED
    ax_y.line.width = Pt(0.75)
    import random
    rng = random.Random(13)
    plot_w = w - 2 * pad
    plot_h = h - 2 * pad
    palette = [ACCENT_TEAL, ACCENT_PRIMARY, ACCENT_AMBER, ACCENT_GREEN, ACCENT_RED]
    for _ in range(80):
        dx = rng.random() ** 1.2
        dy = rng.random() ** 0.8
        cx = x0 + Inches(0.05) + dx * (plot_w - Inches(0.1))
        cy = y0 - Inches(0.05) - dy * (plot_h - Inches(0.1))
        r = Inches(0.07)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 cx - r/2, cy - r/2, r, r)
        dot.fill.solid()
        dot.fill.fore_color.rgb = palette[rng.randint(0, 4)]
        dot.line.fill.background()
        dot.fill.transparency = 0.4
    add_text(s, "n_tags →", x0, y0 + Inches(0.05), Inches(2), Inches(0.3),
             size=10, color=TEXT_MUTED)
    add_text(s, "title_len ↑", x0 - Inches(0.05), top + pad - Inches(0.05),
             Inches(2), Inches(0.3),
             size=10, color=TEXT_MUTED)

    add_footer(s, page, total)
    return s


# ----------------------------------------------------------------
# BREAK
# ----------------------------------------------------------------
def slide_break(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_BREAK)
    add_text(s, "BREAK",
             Inches(0), Inches(2.7), Inches(13.333), Inches(1.3),
             size=72, bold=True, color=ACCENT_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, "10 分鐘休息  ·  回來繼續 教育部高教統計",
             Inches(0), Inches(4.1), Inches(13.333), Inches(0.7),
             size=22, color=TEXT_DARK, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    return s


# ----------------------------------------------------------------
# Dataset B
# ----------------------------------------------------------------
def slide_moe_intro(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 B · 教育部高教統計")
    add_title(s, "少子化怎麼反映在高教資料？")

    add_bullets(s, [
        "教育部統計處每年都會公開「大專校院校別學生數」CSV",
        "本週合併 105–113 學年度 9 年的資料 → 約 7,200 列",
        "169 所大專、21 個縣市 — 涵蓋全國",
        "可同時練：跨年度合併、中文欄位清理、用規則 (rule-based) 推導分類",
    ], Inches(0.9), Inches(1.95), Inches(11.5), Inches(2.1), size=16)

    kpis = [
        ("18.0%",  "總學生數 9 年累計跌幅",     ACCENT_RED),
        ("-23.5 萬", "9 年累計減少人數",        ACCENT_AMBER),
        ("-29%",   "私立部門 9 年跌幅",        ACCENT_ORANGE),
    ]
    x0, y0 = Inches(0.9), Inches(4.4)
    w, h = Inches(3.8), Inches(1.9)
    gap = Inches(0.15)
    for i, (num, label, color) in enumerate(kpis):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=color, accent_w=0.10)
        add_text(s, num, left + Inches(0.3), y0 + Inches(0.25),
                 w - Inches(0.5), Inches(0.85),
                 size=44, bold=True, color=color)
        add_text(s, label, left + Inches(0.3), y0 + Inches(1.15),
                 w - Inches(0.5), Inches(0.6),
                 size=14, color=TEXT_DARK, font=FONT_CJK)

    add_demo_callout(s, "VS Code 打開 .ipynb 跟著做",
                     "code/moe_higher_ed_pipeline.ipynb")
    add_footer(s, page, total)
    return s


def slide_schema_align(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 B · 跨年度合併")
    add_title(s, "不同年度的欄位不一樣怎麼辦？")

    add_text(s,
             "Schema (資料表結構) = 「這張表有哪些欄位」。教育部的 schema 多年來有變動。",
             Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_bullets(s, [
        "105–106 學年度：23 欄，沒有「總計」、「縣市」、「體系別」",
        "107 學年度起：26 欄，加入「縣市名稱」、「體系別」",
        "113 學年度：28 欄，新增「華語先修生男/女」",
        "pd.concat(ignore_index=True, sort=False) — 缺欄位自動補成 NaN",
    ], Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.0), size=15)

    code = (
        "def fetch_all(years=range(105, 114)):\n"
        "    parts = [fetch_year(y) for y in years]\n"
        "    return pd.concat(parts, ignore_index=True, sort=False)\n"
        "\n"
        "# 缺值 (NaN) 不一定是真的缺 — 也可能是「那年沒這個欄位」"
    )
    add_code(s, code, Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.5),
             size=13)

    add_callout(s,
        "拿到合併好的表格，第一件事永遠是 df.isna().mean() — 看缺值在哪。",
        Inches(0.9), Inches(6.30), Inches(11.5), Inches(0.55),
        size=15, color=ACCENT_PRIMARY)

    add_footer(s, page, total)
    return s


def slide_chinese_cleaning(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 B · 中文欄位清理")
    add_title(s, "真實資料常常藏著隱形規則")

    rows = [
        ("「縣市名稱」=「30 臺北市」",
         "用 regex 拆出純名稱",
         "→ city_name"),
        ("「體系別」=「1 一般」/「2 技職」",
         "用 str.replace 去掉前面數字",
         "→ system"),
        ("「總計」型別變成 object",
         "to_numeric + 千分號處理",
         "→ NaN 化髒值"),
        ("公私立沒有專屬欄位",
         "看「國立/市立」開頭推導",
         "→ sector"),
    ]
    add_table(s,
              ["觀察", "操作", "結果"],
              rows,
              Inches(0.85), Inches(2.05), Inches(11.6), Inches(3.4),
              col_widths=[Inches(3.7), Inches(4.8), Inches(3.1)],
              header_size=13, body_size=12)

    add_text(s, "⚠ 用規則推導 (rule-based) 一定要人工抽查 — 「國防大學」、「警察大學」需要特殊規則",
             Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.45),
             size=14, italic=True, color=ACCENT_RED, font=FONT_CJK)

    add_demo_callout(s, "在 .ipynb 第 5 格看清理過程",
                     "code/moe_higher_ed_pipeline.ipynb")
    add_footer(s, page, total)
    return s


def slide_moe_total(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 B · 圖 1 · 總量趨勢")
    add_title(s, "整體真的在下降嗎？")

    add_bullets(s, [
        "政策報告的第一張圖：先回答「整體是不是真的在下降」",
        "用 add_annotation 標出 105 學年度起點 — 提醒讀者這是少子化轉折點",
        "Annotation 是告訴讀者「為什麼這點重要」，不只是標數值",
    ], Inches(0.9), Inches(1.95), Inches(6.5), Inches(2.0), size=15)

    code = (
        "agg = df.groupby(\"學年度\")[\"總計\"].sum().reset_index()\n"
        "fig = px.line(agg, x=\"學年度\", y=\"總計\",\n"
        "              markers=True)\n"
        "fig.add_annotation(\n"
        "    x=agg.iloc[0][\"學年度\"],\n"
        "    y=agg.iloc[0][\"總計\"],\n"
        "    text=\"少子化骨牌效應起點\",\n"
        "    showarrow=True, arrowhead=2,\n"
        ")"
    )
    add_code(s, code, Inches(0.9), Inches(4.1), Inches(6.4), Inches(2.3),
             size=12)

    add_card(s, Inches(7.7), Inches(1.95), Inches(4.8), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_ORANGE)
    add_text(s, "預覽 · 總學生數 105–113",
             Inches(7.9), Inches(2.10), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=ACCENT_ORANGE, font=FONT_CJK)
    real = [1.309, 1.274, 1.245, 1.213, 1.203, 1.186, 1.140, 1.095, 1.074]
    vmin, vmax = min(real), max(real)
    values = [(v - vmin) / (vmax - vmin) * 0.80 + 0.10 for v in real]
    add_mini_line_chart(s, Inches(7.9), Inches(2.65),
                        Inches(4.4), Inches(3.2),
                        values, color=ACCENT_ORANGE)
    add_text(s, "起點 1.31M",
             Inches(8.2), Inches(2.50), Inches(1.5), Inches(0.3),
             size=10, bold=True, color=ACCENT_RED)
    add_text(s, "終點 1.07M",
             Inches(11.05), Inches(5.55), Inches(1.4), Inches(0.3),
             size=10, bold=True, color=ACCENT_RED)
    add_text(s, "105   106   107   108   109   110   111   112   113",
             Inches(7.95), Inches(5.92), Inches(4.4), Inches(0.3),
             size=9, color=TEXT_MUTED)
    add_text(s, "9 年累計 −18%",
             Inches(8.55), Inches(6.10), Inches(3.5), Inches(0.25),
             size=10, bold=True, italic=True, color=ACCENT_ORANGE, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_moe_sector(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "資料 B · 圖 2 · 公立 vs. 私立")
    add_title(s, "下降集中在公立還是私立？")

    add_bullets(s, [
        "用 stacked bar — 一根長條同時呈現兩個 sector 的數量與總和",
        "結論：私立 9 年減 25.6 萬，公立反而微增 2.1 萬",
        "Caption 應該寫結論（「私立首當其衝」），不要描述圖表類型",
    ], Inches(0.9), Inches(1.95), Inches(6.5), Inches(2.0), size=15)

    code = (
        "agg = (df.groupby([\"學年度\", \"sector\"])\n"
        "         [\"總計\"].sum().reset_index())\n"
        "fig = px.bar(\n"
        "    agg, x=\"學年度\", y=\"總計\",\n"
        "    color=\"sector\", barmode=\"stack\",\n"
        "    color_discrete_map={\n"
        "        \"公立\": \"#1d4ed8\",\n"
        "        \"私立\": \"#f97316\"},\n"
        ")"
    )
    add_code(s, code, Inches(0.9), Inches(4.0), Inches(6.4), Inches(2.3),
             size=13)

    add_card(s, Inches(7.7), Inches(1.95), Inches(4.8), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_ORANGE)
    add_text(s, "預覽 · 公立 (藍) / 私立 (橘)",
             Inches(7.9), Inches(2.10), Inches(4.5), Inches(0.35),
             size=12, bold=True, color=ACCENT_ORANGE, font=FONT_CJK)

    pub_pri = [
        (430, 879), (431, 843), (430, 815), (431, 782),
        (439, 764), (445, 741), (449, 691), (451, 643), (451, 623),
    ]
    max_total = max(a + b for a, b in pub_pri)
    bar_area_l = Inches(7.95)
    bar_area_t = Inches(2.6)
    bar_area_w = Inches(4.45)
    bar_area_h = Inches(3.5)
    bar_count = len(pub_pri)
    bar_w_each = Inches(0.40)
    gap = (bar_area_w - bar_w_each * bar_count) / (bar_count + 1)
    for i, (pub, pri) in enumerate(pub_pri):
        total_h = bar_area_h * ((pub + pri) / max_total)
        pub_h = bar_area_h * (pub / max_total)
        pri_h = total_h - pub_h
        x = bar_area_l + gap + i * (bar_w_each + gap)
        y_bot = bar_area_t + bar_area_h
        b1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x, y_bot - total_h, bar_w_each, pri_h)
        b1.fill.solid(); b1.fill.fore_color.rgb = ACCENT_ORANGE
        b1.line.fill.background()
        b2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                x, y_bot - pub_h, bar_w_each, pub_h)
        b2.fill.solid(); b2.fill.fore_color.rgb = ACCENT_PRIMARY
        b2.line.fill.background()
    years = ["105", "106", "107", "108", "109", "110", "111", "112", "113"]
    for i, yr in enumerate(years):
        x = bar_area_l + gap + i * (bar_w_each + gap)
        add_text(s, yr, x - Inches(0.1),
                 bar_area_t + bar_area_h + Inches(0.05),
                 bar_w_each + Inches(0.2), Inches(0.3),
                 size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_footer(s, page, total)
    return s


# ----------------------------------------------------------------
# Synthesis
# ----------------------------------------------------------------
def slide_storytelling(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 4 · 資料敘事")
    add_title(s, "讓圖表會說故事的五個原則")

    rows = [
        ("一張圖只回答一個問題",
         "別把太多訊息塞進同一張圖",
         "Top subjects bar 只回答「哪個主題最多」"),
        ("先看大局，再看細節",
         "從整體開始，再拆分類別",
         "MOE: 先總量下降，再拆公私立"),
        ("把重點標出來",
         "用 annotation + 文字解釋",
         "「少子化骨牌效應起點」"),
        ("呈現不確定性",
         "點密度、誤差棒都是訊號",
         "scatter 用 opacity=0.6 暗示密度"),
        ("複雜的圖要有理由",
         "能用 bar 就不用 map",
         "城市排名用 bar 比 choropleth 更清楚"),
    ]
    add_table(s,
              ["原則", "意思", "本週對應例子"],
              rows,
              Inches(0.85), Inches(2.0), Inches(11.6), Inches(4.4),
              col_widths=[Inches(3.4), Inches(3.4), Inches(4.8)],
              header_size=13, body_size=12, first_col_bold=True)

    add_footer(s, page, total)
    return s


def slide_pitfalls(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "COMMON PITFALLS")
    add_title(s, "互動圖表常見錯誤")

    pitfalls = [
        ("✗", "抓 API 只拿第一頁",
         "誤以為 PsyArXiv 只有 100 篇心理學論文"),
        ("✗", "concat 後不檢查欄位",
         "某些欄位全 NaN 卻沒人發現"),
        ("✗", "to_numeric 不加 errors='coerce'",
         "一個髒字串就讓整支 script 炸掉"),
        ("✗", "塞 200 種 subject 進 legend",
         "圖完全沒有訊息量；應該 showlegend=False"),
        ("✗", "Annotation 只寫值不寫意義",
         "「Peak: 87」遠不如「Peak — conference deadline」"),
    ]
    y = 2.05
    for mark, head, body in pitfalls:
        add_text(s, mark, Inches(0.9), Inches(y),
                 Inches(0.6), Inches(0.5),
                 size=22, bold=True, color=ACCENT_RED)
        add_text(s, head, Inches(1.6), Inches(y),
                 Inches(4.5), Inches(0.5),
                 size=17, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, Inches(6.2), Inches(y),
                 Inches(6.5), Inches(0.5),
                 size=15, color=TEXT_MUTED, font=FONT_CJK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += 0.85

    add_footer(s, page, total)
    return s


def slide_homework(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "HOMEWORK · 下週上課前繳交")
    add_title(s, "把兩個資料集整合進你的儀表板")

    add_bullets(s, [
        "新增 PsyArXiv 頁面：fetch_psyarxiv() + clean()，至少 2 張 Plotly 圖",
        "新增 MOE 頁面：fetch_all() + clean()，至少 2 張 Plotly 圖",
        "至少一張圖要含 annotation；每張圖下方寫 1–2 句 takeaway",
        "每個資料集至少一個 widget（slider / selectbox / multiselect）",
        "Docstring 寫一段：哪個清理決定你刻意與範例不同？理由？",
    ], Inches(0.9), Inches(1.95), Inches(11.5), Inches(2.6), size=15)

    rows = [
        ("Pipeline 完整 (load → clean → describe → plot)", "30%"),
        ("Plotly 圖數量與多樣性 (≥ 4 張，含 annotation)",   "25%"),
        ("資料敘事 — caption + annotation",                  "20%"),
        ("清理決定的反思（與範例不同 + 理由）",              "15%"),
        ("Streamlit 互動 (widget + 排版)",                   "10%"),
    ]
    add_table(s, ["評分項目", "比重"], rows,
              Inches(0.9), Inches(4.8), Inches(11.5), Inches(2.0),
              col_widths=[Inches(9.0), Inches(2.5)],
              header_size=13, body_size=12)

    add_footer(s, page, total)
    return s


def slide_takeaway(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_SECTION)
    add_text(s, "TAKEAWAY",
             Inches(0.7), Inches(1.0), Inches(12), Inches(0.55),
             size=20, bold=True, color=ACCENT_TEAL)
    add_text(s, "資料可信，圖表才有說服力",
             Inches(0.7), Inches(1.65), Inches(12.5), Inches(1.0),
             size=42, bold=True, color=TEXT_LIGHT, font=FONT_CJK)

    pillars = [
        ("Pipeline",
         "讓資料可信",
         "load → clean → describe"),
        ("Interactivity",
         "讓資料可探索",
         "Plotly + Streamlit"),
        ("Storytelling",
         "讓分析有方向",
         "一張圖 / 一個問題"),
    ]
    x0, y0 = Inches(0.9), Inches(3.5)
    w, h = Inches(3.85), Inches(2.5)
    gap = Inches(0.18)
    for i, (title, sub, body) in enumerate(pillars):
        left = x0 + i * (w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y0, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1F, 0x40, 0x6F)
        card.line.color.rgb = ACCENT_TEAL
        card.line.width = Pt(1)
        add_text(s, title, left + Inches(0.25), y0 + Inches(0.25),
                 w - Inches(0.5), Inches(0.55),
                 size=22, bold=True, color=ACCENT_TEAL)
        add_text(s, sub, left + Inches(0.25), y0 + Inches(0.95),
                 w - Inches(0.5), Inches(0.5),
                 size=17, color=TEXT_LIGHT, font=FONT_CJK)
        add_text(s, body, left + Inches(0.25), y0 + Inches(1.55),
                 w - Inches(0.5), Inches(0.85),
                 size=13, color=ACCENT_TEAL, font=FONT_MONO)

    add_text(s, "儀表板的目標不是展示所有資料，而是幫使用者做判斷。",
             Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
             size=18, italic=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)
    return s


# ============================================================
# Build the deck
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 24

    slide_title(prs)                                       # 1
    slide_why(prs, 2, TOTAL)                               # 2
    slide_two_datasets(prs, 3, TOTAL)                      # 3
    slide_objectives(prs, 4, TOTAL)                        # 4
    slide_dev_env(prs, 5, TOTAL)                           # 5 (NEW — PS / VS Code / ipython)
    slide_plotly_role(prs, 6, TOTAL)                       # 6
    slide_plotly_basics(prs, 7, TOTAL)                     # 7

    # Dataset A
    slide_psyarxiv_intro(prs, 8, TOTAL)                    # 8
    slide_pagination(prs, 9, TOTAL)                        # 9
    slide_json_to_df(prs, 10, TOTAL)                       # 10
    slide_cleaning_decisions(prs, 11, TOTAL)               # 11
    slide_psyarxiv_bar(prs, 12, TOTAL)                     # 12
    slide_psyarxiv_line(prs, 13, TOTAL)                    # 13
    slide_psyarxiv_scatter(prs, 14, TOTAL)                 # 14

    slide_break(prs, 15, TOTAL)                            # 15

    # Dataset B
    slide_moe_intro(prs, 16, TOTAL)                        # 16
    slide_schema_align(prs, 17, TOTAL)                     # 17
    slide_chinese_cleaning(prs, 18, TOTAL)                 # 18
    slide_moe_total(prs, 19, TOTAL)                        # 19
    slide_moe_sector(prs, 20, TOTAL)                       # 20

    # Synthesis
    slide_storytelling(prs, 21, TOTAL)                     # 21
    slide_pitfalls(prs, 22, TOTAL)                         # 22
    slide_homework(prs, 23, TOTAL)                         # 23
    slide_takeaway(prs, 24, TOTAL)                         # 24

    out = "week-13-slides.pptx"
    prs.save(out)
    print(f"Saved {out}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
