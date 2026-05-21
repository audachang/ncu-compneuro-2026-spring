"""Build Week 14 (Supplement) lecture slides — Pingouin Statistics.

Output : week-14-pingouin-slides.pptx (16:9)
Visual : ACL@NCU palette — white BG, navy + teal accents, JhengHei + Calibri.
Source : week-14-pingouin-slide-outline.md

Run    : python build_pingouin_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# Palette — keep consistent with Week 13
# ============================================================
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE    = RGBColor(0xF7, 0xF8, 0xFA)
BG_SECTION     = RGBColor(0x14, 0x32, 0x5C)
BG_PRACTICE    = RGBColor(0xEC, 0xF7, 0xF6)
BG_PITFALL     = RGBColor(0xFE, 0xF1, 0xF1)
TEXT_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED     = RGBColor(0x5F, 0x6B, 0x83)
ACCENT_NAVY    = RGBColor(0x14, 0x32, 0x5C)
ACCENT_TEAL    = RGBColor(0x0D, 0x9B, 0x9B)
ACCENT_AMBER   = RGBColor(0xE8, 0xA1, 0x2A)
ACCENT_RED     = RGBColor(0xD3, 0x4F, 0x4F)
ACCENT_GREEN   = RGBColor(0x2E, 0x8B, 0x57)
CODE_BG        = RGBColor(0x1E, 0x29, 0x3B)
HAIRLINE       = RGBColor(0xE2, 0xE6, 0xEC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"
FOOTER_STR = "Week 14 — Pingouin Statistics  ·  ACL@NCU"


# ============================================================
# Primitives
# ============================================================
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_top_band(slide, color=ACCENT_TEAL, height=Inches(0.14)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # CJK font for runs that include CJK characters
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.makeelement(qn('a:ea'),
                                {'typeface': FONT_CJK})
    rPr.append(eastAsia)
    return tb


def add_section_label(slide, text):
    add_text(slide, text, Inches(0.85), Inches(0.45), Inches(8), Inches(0.32),
             size=12, bold=True, color=ACCENT_TEAL,
             font=FONT_SANS)


def add_title(slide, title_zh, *, size=30, top=0.85, color=TEXT_DARK):
    return add_text(slide, title_zh, Inches(0.85), Inches(top), Inches(11.6),
                    Inches(0.85), size=size, bold=True, color=color,
                    font=FONT_CJK)


def add_bullets(slide, items, left, top, width, height, *,
                size=16, color=TEXT_DARK, font=FONT_SANS,
                line_spacing=1.25, bullet_char="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True

    from pptx.oxml.ns import qn

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass

        # bullet glyph
        bp = p.add_run()
        bp.text = f"{bullet_char}  "
        bp.font.name = font
        bp.font.size = Pt(size)
        bp.font.color.rgb = ACCENT_TEAL
        bp.font.bold = True

        # body
        body = p.add_run()
        body.text = str(item)
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = color
        rPr = body._r.get_or_add_rPr()
        eastAsia = rPr.makeelement(qn('a:ea'),
                                    {'typeface': FONT_CJK})
        rPr.append(eastAsia)
    return tb


def add_callout(slide, text, left, top, width, height, *,
                fill=BG_PRACTICE, border=ACCENT_TEAL,
                color=TEXT_DARK, size=14, bold=False):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1.0)
    tf = card.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT_SANS
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    eastAsia = rPr.makeelement(qn('a:ea'),
                                {'typeface': FONT_CJK})
    rPr.append(eastAsia)
    return card


def add_code(slide, code, left, top, width, height, *, size=12):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = CODE_BG
    card.line.fill.background()
    tf = card.text_frame
    tf.margin_left = Inches(0.20); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.14)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = (RGBColor(0xA8, 0xC0, 0xFF) if line.lstrip().startswith("#")
                              else RGBColor(0xEC, 0xEF, 0xF4))
    return card


def add_footer(slide, page, total):
    add_text(slide, FOOTER_STR, Inches(0.85), Inches(7.10),
             Inches(9), Inches(0.30),
             size=10, color=TEXT_MUTED, font=FONT_SANS)
    add_text(slide, f"{page} / {total}", Inches(11.85), Inches(7.10),
             Inches(1.30), Inches(0.30),
             size=10, color=TEXT_MUTED, font=FONT_SANS,
             align=PP_ALIGN.RIGHT)


def add_table(slide, headers, rows, left, top, width, height, *,
              header_fill=ACCENT_NAVY, alt_fill=BG_OFFWHITE,
              header_color=TEXT_LIGHT, body_color=TEXT_DARK, size=13):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    from pptx.oxml.ns import qn

    def cell_text(cell, text, *, color, bold=False, fill=None):
        cell.text = ""
        if fill is not None:
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(text)
        run.font.name = FONT_SANS
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        rPr = run._r.get_or_add_rPr()
        eastAsia = rPr.makeelement(qn('a:ea'),
                                    {'typeface': FONT_CJK})
        rPr.append(eastAsia)

    for j, h in enumerate(headers):
        cell_text(tbl.cell(0, j), h, color=header_color, bold=True,
                  fill=header_fill)
    for i, row in enumerate(rows, start=1):
        fill = alt_fill if i % 2 == 0 else BG_WHITE
        for j, v in enumerate(row):
            cell_text(tbl.cell(i, j), v, color=body_color, fill=fill)
    return tbl_shape


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ============================================================
# Slide base
# ============================================================
def new_slide(prs, *, bg=BG_WHITE, band=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg)
    if band:
        add_top_band(slide, ACCENT_TEAL)
    return slide


# ============================================================
# Build individual slides
# ============================================================
TOTAL = 25  # 將在 build() 結束時動態確認


def slide_01_title(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_SECTION)
    # accent dot top-left
    accent = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.85), Inches(0.85),
                                    Inches(0.32), Inches(0.32))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT_TEAL
    accent.line.fill.background()

    add_text(slide, "Week 14 (Supplement)", Inches(0.85), Inches(1.45),
             Inches(11), Inches(0.5), size=18, color=ACCENT_TEAL,
             bold=True, font=FONT_SANS)
    add_text(slide, "用 Pingouin 做行為實驗統計", Inches(0.85), Inches(2.05),
             Inches(11.5), Inches(1.4), size=44, bold=True,
             color=TEXT_LIGHT, font=FONT_CJK)
    add_text(slide, "Behavioral Statistics with Pingouin",
             Inches(0.85), Inches(3.55), Inches(11.5), Inches(0.6),
             size=22, color=RGBColor(0xCA, 0xDC, 0xFC), italic=True,
             font=FONT_SANS)
    add_text(slide, "從 t-test 到 Streamlit dashboard 的完整 pipeline",
             Inches(0.85), Inches(4.20), Inches(11.5), Inches(0.6),
             size=18, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_CJK)

    add_text(slide, "NS5116 · 2026 Spring · ACL@NCU · 張智宏 Erik Chang",
             Inches(0.85), Inches(6.60), Inches(11), Inches(0.4),
             size=12, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_SANS)
    set_notes(slide,
              "本週把 ML 主題暫停一格，回到統計推論。學生在 Week 7 學了 NumPy、"
              "Week 12 學了 pandas、Week 13 學了 Plotly、Week 11–12 學了 Streamlit。"
              "今天把這些全串起來，做出 publication-ready 的統計分析。")


def slide_02_scipy_vs_pg(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§1  WHY PINGOUIN")
    add_title(slide, "scipy.stats 給你「結果」，pingouin 給你「報告」")

    add_bullets(slide,
                ["scipy.stats.ttest_rel() 只回傳 (t, p) namedtuple",
                 "論文要求：effect size, df, 95% CI, Bayes Factor, power",
                 "自己從 scipy 組起來 ≈ 寫 20 行 boilerplate",
                 "pg.ttest() 一行就有 8 欄位的 DataFrame"],
                Inches(0.85), Inches(1.85), Inches(6.2), Inches(3),
                size=16)

    add_table(slide,
              ["", "scipy.stats", "pingouin"],
              [["輸出型別", "namedtuple", "DataFrame"],
               ["欄位數",   "2–3",         "7–8"],
               ["Effect size", "自己算", "內建"],
               ["Bayes Factor", "無",     "內建"],
               ["API 一致性", "散亂",     "統一"]],
              Inches(7.4), Inches(1.85), Inches(5.2), Inches(3.6),
              size=13)
    add_footer(slide, page, total)
    set_notes(slide,
              "強調對研究生的痛點：論文 reviewer 會要求 effect size。"
              "pingouin 把這些都打包進一個 DataFrame，省下整理表格的時間。")


def slide_03_objectives(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "OBJECTIVES")
    add_title(slide, "今天你會學會什麼")

    items = [
        ("①", "安裝 pingouin 並理解 API 哲學"),
        ("②", "t-test 家族 + effect size + Bayes Factor"),
        ("③", "One-way / Repeated-measures / Mixed ANOVA"),
        ("④", "Correlation 與 partial correlation"),
        ("⑤", "Streamlit 部署互動式統計報告"),
    ]
    for i, (num, label) in enumerate(items):
        top = 2.0 + i * 0.85
        # circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(1.10), Inches(top),
                                        Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_TEAL
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = num
        run.font.name = FONT_SANS; run.font.size = Pt(18)
        run.font.bold = True; run.font.color.rgb = TEXT_LIGHT

        add_text(slide, label, Inches(1.85), Inches(top + 0.08),
                 Inches(10), Inches(0.5), size=18, color=TEXT_DARK,
                 font=FONT_CJK)
    add_footer(slide, page, total)
    set_notes(slide,
              "100 分鐘課程 + 5 個 hands-on practice。"
              "前 4 個是純 Python，第 5 個整合到 Streamlit。")


def slide_04_install(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§1  INSTALL & FIRST TEST")
    add_title(slide, "一行安裝，一行 test")

    add_code(slide, "pip install pingouin",
             Inches(0.85), Inches(1.90), Inches(11.6), Inches(0.65),
             size=16)

    add_code(slide,
             "import pingouin as pg\n\n"
             "# paired t-test — Stroop within-subject\n"
             "result = pg.ttest(congruent, incongruent, paired=True)\n"
             "print(result)",
             Inches(0.85), Inches(2.75), Inches(7.0), Inches(2.4),
             size=14)

    add_callout(slide,
                "輸出欄位 (pingouin ≥ 0.6):\n"
                "T  ·  dof  ·  p_val  ·  CI95  ·\n"
                "cohen_d  ·  BF10  ·  power  ·  alternative",
                Inches(8.10), Inches(2.75), Inches(4.5), Inches(2.4),
                fill=BG_PRACTICE, border=ACCENT_TEAL, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "提醒學生 pingouin 0.6 之後欄位用 underscore (p_val, cohen_d, CI95)。"
              "舊版是 hyphen (p-val, cohen-d, CI95%)。如果他們 import 時遇到 KeyError，"
              "先用 pip show pingouin 確認版本。")


def slide_05_philosophy(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§1  API PHILOSOPHY")
    add_title(slide, "One test = One DataFrame")

    add_bullets(slide,
                ["所有 function 回傳統一格式：pd.DataFrame",
                 "可直接 to_latex() / to_csv() / st.dataframe()",
                 "欄位命名與論文用語一致 (cohen_d, CI95)",
                 "與 pandas / statsmodels 生態系完美銜接"],
                Inches(0.85), Inches(1.85), Inches(6.5), Inches(3),
                size=16)

    # Diagram: DataFrame -> 4 outputs
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(8.0), Inches(2.3),
                                 Inches(3.5), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT_NAVY
    box.line.fill.background()
    add_text(slide, "pingouin\nDataFrame", Inches(8.0), Inches(2.42),
             Inches(3.5), Inches(0.6), size=15, bold=True,
             color=TEXT_LIGHT, font=FONT_SANS, align=PP_ALIGN.CENTER)

    targets = [("→ LaTeX (paper)", 3.50),
               ("→ CSV (download)", 4.10),
               ("→ Streamlit dashboard", 4.70),
               ("→ Markdown (report)", 5.30)]
    for txt, top in targets:
        add_text(slide, txt, Inches(8.20), Inches(top),
                 Inches(4), Inches(0.4),
                 size=13, color=TEXT_DARK, font=FONT_SANS)
    add_footer(slide, page, total)
    set_notes(slide,
              "重點：DataFrame 是 publication 的單位，不是 namedtuple。"
              "學生在 Week 12-13 已學過 pandas/Streamlit，能無痛接上。")


def slide_06_handson1(prs, page, total):
    slide = new_slide(prs, bg=BG_PRACTICE)
    add_section_label(slide, "HANDS-ON 1")
    add_title(slide, "比較 scipy 與 pingouin 輸出")

    add_callout(slide, "5 min", Inches(11.0), Inches(0.55),
                Inches(1.5), Inches(0.45),
                fill=ACCENT_AMBER, border=ACCENT_AMBER,
                color=TEXT_LIGHT, size=14, bold=True)

    add_bullets(slide,
                ["任務：對 young vs. old RT 跑 independent t-test",
                 "印出 Cohen's d (cohen_d) 與 Bayes Factor (BF10)",
                 "解讀：BF10 > 10 代表什麼？"],
                Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.0),
                size=17)

    add_code(slide,
             "import numpy as np, pingouin as pg\n"
             "np.random.seed(0)\n"
             "young = np.random.normal(400, 50, 25)\n"
             "old   = np.random.normal(480, 70, 25)\n"
             "# 你的程式碼從這裡開始",
             Inches(0.85), Inches(4.05), Inches(11.6), Inches(2.4),
             size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "BF10 > 10 = strong evidence for H1 (alternative)。"
              "若 BF10 < 1/10，反而支持 H0 (null) — 這是頻率派 t-test 做不到的。")


def slide_07_ttest_family(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§2  T-TEST FAMILY")
    add_title(slide, "Paired / Independent / One-sample")

    rows = [
        ("Paired", "同一受試者多次測量", "pg.ttest(a, b, paired=True)"),
        ("Independent", "兩組不同受試者", "pg.ttest(a, b, paired=False)"),
        ("One-sample", "與固定值比較", "pg.ttest(a, popmean)"),
    ]
    add_table(slide,
              ["Test 類型", "適用情境", "API 語法"],
              rows,
              Inches(0.85), Inches(1.95), Inches(11.6), Inches(2.5),
              size=14)

    add_callout(slide,
                "決策原則：先問「同一個 subject 嗎？」\n"
                "  Yes → paired=True\n"
                "  No  → paired=False (independent)\n"
                "  與固定值比 → 第二個 argument 給 scalar",
                Inches(0.85), Inches(4.85), Inches(11.6), Inches(1.55),
                fill=BG_PRACTICE, border=ACCENT_TEAL, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "最常見的錯誤：within-subject 設計用了 paired=False，"
              "失去 within-subject 的 statistical power。"
              "另一個錯誤：one-sample 時誤把 baseline 變成等長 array。")


def slide_08_nonparametric(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§2  NON-PARAMETRIC")
    add_title(slide, "當 normality 假設違反時")

    add_bullets(slide,
                ["pg.normality() 先檢查 — Shapiro-Wilk test",
                 "Wilcoxon signed-rank：pg.wilcoxon() (paired)",
                 "Mann-Whitney U：pg.mwu() (independent)",
                 "何時用：RT 嚴重 skewed、小樣本、Likert scale"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(3),
                size=16)

    add_callout(slide,
                "📌 老人受試者的 RT 常見 right-skew distribution\n"
                "📌 兒童的 accuracy data 不適合用 parametric test\n"
                "📌 跑 pg.normality() 看 normal 欄位是 True / False",
                Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.55),
                fill=BG_OFFWHITE, border=ACCENT_NAVY, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "normality 違反不等於不能跑 parametric — t-test 在 n > 30 時"
              "對 normality 違反相對 robust。但 reviewer 會問，所以最好附上檢驗結果。")


def slide_09_effsize(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§2  EFFECT SIZE")
    add_title(slide, "p < .05 不夠 — 為什麼要報 effect size")

    add_bullets(slide,
                ["p-value 受 sample size 影響很大 (n 大就容易顯著)",
                 "Cohen's d 反映實際差異大小，獨立於 n",
                 "慣例：0.2 small, 0.5 medium, 0.8 large",
                 "pingouin 自動回報；也可用 pg.compute_effsize()"],
                Inches(0.85), Inches(1.85), Inches(7.0), Inches(3),
                size=16)

    # mini visual: two scenarios
    add_callout(slide,
                "情境 A\nn = 1000\np = 0.01 ✓\nd = 0.10  ✗",
                Inches(8.30), Inches(2.0), Inches(2.0), Inches(2.0),
                fill=BG_PITFALL, border=ACCENT_RED, size=14, bold=True)
    add_callout(slide,
                "情境 B\nn = 30\np = 0.04 ✓\nd = 0.85  ✓",
                Inches(10.50), Inches(2.0), Inches(2.0), Inches(2.0),
                fill=BG_PRACTICE, border=ACCENT_GREEN, size=14, bold=True)
    add_text(slide, "同樣 p < .05，但 effect size 天差地遠",
             Inches(8.30), Inches(4.20), Inches(4.5), Inches(0.4),
             size=12, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(slide, page, total)
    set_notes(slide,
              "Lakens (2013) 是經典 reference — effect size reporting 已是論文標配。"
              "Hedges' g 是 small-sample bias-corrected 版本，n < 50 時建議用 g 而非 d。")


def slide_10_handson2(prs, page, total):
    slide = new_slide(prs, bg=BG_PRACTICE)
    add_section_label(slide, "HANDS-ON 2")
    add_title(slide, "完整的 Stroop effect 報告")

    add_callout(slide, "10 min", Inches(11.0), Inches(0.55),
                Inches(1.5), Inches(0.45),
                fill=ACCENT_AMBER, border=ACCENT_AMBER,
                color=TEXT_LIGHT, size=14, bold=True)

    add_bullets(slide,
                ["任務：產出 paired-t 與 Wilcoxon 的合併表格",
                 "用 pd.concat([t_res, w_res], keys=...) 串接",
                 "寫一句 APA-style 結果報告"],
                Inches(0.85), Inches(1.85), Inches(11.5), Inches(1.6),
                size=17)

    add_callout(slide,
                "APA 範例：\n"
                "Paired-samples t-test revealed a significant Stroop effect, "
                "t(27) = -4.21, p < .001, Cohen's d = 0.77, 95% CI [-94, -33].",
                Inches(0.85), Inches(3.85), Inches(11.6), Inches(2.5),
                fill=BG_WHITE, border=ACCENT_TEAL, size=14, bold=False)
    add_footer(slide, page, total)
    set_notes(slide,
              "APA 7th 的 reporting standard — 必須包含 test statistic, df, p, "
              "effect size, CI。pingouin 把這些都一次給你。")


def slide_11_anova_oneway(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§3  ANOVA")
    add_title(slide, "One-way ANOVA — 3 個 condition 以上")

    add_bullets(slide,
                ["多次 t-test → inflated Type I error (alpha 累積)",
                 "pg.anova(data=df, dv='rt', between='load')",
                 "輸出：F, DF1, DF2, p_unc, np2 (partial η²)",
                 "顯著後務必做 post-hoc 釐清哪兩組有差"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(3),
                size=16)

    add_callout(slide,
                "情境：N-back working memory load (1-back, 2-back, 3-back)\n"
                "問題：load 對 RT 的整體效應是否顯著？\n"
                "若 F 顯著 (p < .05)，再跑 pairwise tests 看哪兩個 load 不同。",
                Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.55),
                fill=BG_OFFWHITE, border=ACCENT_NAVY, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "為什麼 partial η² 重要：它是 between-factor variance 占總 variance 的比例，"
              "比 R² 更能呈現單一因子的解釋力。慣例：0.01 small, 0.06 medium, 0.14 large。")


def slide_12_rm_anova(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§3  REPEATED MEASURES")
    add_title(slide, "rm-ANOVA — 同一受試者多次測量")

    add_code(slide,
             "pg.rm_anova(\n"
             "    data=df, dv='rt',\n"
             "    within='load', subject='subject',\n"
             "    detailed=True\n"
             ")",
             Inches(0.85), Inches(1.95), Inches(6.0), Inches(2.4),
             size=14)

    add_bullets(slide,
                ["必傳 subject= 參數 (告訴 pg 哪些 row 是同一人)",
                 "資料須為 long format (用 df.melt() 轉)",
                 "自動回報 Greenhouse-Geisser ε",
                 "若 ε < 0.75 → 引用 GG 校正版本 p-value"],
                Inches(7.20), Inches(1.95), Inches(5.5), Inches(3),
                size=15)

    add_callout(slide,
                "⚠️ Sphericity 違反是 within-subject design 最容易踩的雷。",
                Inches(0.85), Inches(5.05), Inches(11.6), Inches(0.85),
                fill=BG_PITFALL, border=ACCENT_RED, size=14, bold=True)
    add_footer(slide, page, total)
    set_notes(slide,
              "Sphericity = 各 condition 兩兩差的 variance 相等。"
              "Mauchly test 顯著 → 違反 sphericity → 用 GG 校正。"
              "pingouin 把 ε 直接放進輸出表格，不用另外算。")


def slide_13_posthoc_mixed(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§3  POST-HOC & MIXED")
    add_title(slide, "哪兩組有差？兩個 factor 怎麼跑？")

    add_bullets(slide,
                ["pg.pairwise_tests(padjust='bonf') — pairwise + Bonferroni",
                 "pg.pairwise_tukey() — 經典 Tukey HSD",
                 "pg.mixed_anova(within=, between=, subject=) — 2 × 2 factorial",
                 "重點欄位：F, p_unc, np2, p_corr"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(3),
                size=16)

    add_callout(slide,
                "經典設計：age (between, young/old) × condition (within, cong/incong)\n"
                "→ Mixed ANOVA 同時檢驗 main effect 與 interaction\n"
                "→ Interaction 顯著 → 做 simple effects analysis",
                Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.55),
                fill=BG_OFFWHITE, border=ACCENT_NAVY, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "Bonferroni 是 conservative 的校正；如果 family-wise comparison 很多，"
              "可改用 FDR (padjust='fdr_bh') 較不保守。")


def slide_14_handson3(prs, page, total):
    slide = new_slide(prs, bg=BG_PRACTICE)
    add_section_label(slide, "HANDS-ON 3")
    add_title(slide, "三組老化研究 — ANOVA + Post-hoc")

    add_callout(slide, "15 min", Inches(11.0), Inches(0.55),
                Inches(1.5), Inches(0.45),
                fill=ACCENT_AMBER, border=ACCENT_AMBER,
                color=TEXT_LIGHT, size=14, bold=True)

    add_bullets(slide,
                ["模擬 young / middle / old 三組，各 n=20",
                 "跑 one-way ANOVA",
                 "加做 pg.pairwise_tukey()",
                 "指出哪兩組顯著不同"],
                Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.0),
                size=17)

    add_code(slide,
             "rows = []\n"
             "for grp, mu in [('young', 420), ('middle', 480), ('old', 560)]:\n"
             "    for _ in range(20):\n"
             "        rows.append({'group': grp, 'rt': np.random.normal(mu, 70)})\n"
             "df = pd.DataFrame(rows)\n"
             "# 你的程式碼從這裡開始",
             Inches(0.85), Inches(4.05), Inches(11.6), Inches(2.4),
             size=13)
    add_footer(slide, page, total)
    set_notes(slide,
              "預期：F 顯著、young-old 差最大、young-middle 與 middle-old 都顯著。"
              "讓學生看 Tukey HSD 自動完成 family-wise 多重比較校正。")


def slide_15_correlation(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§4  CORRELATION")
    add_title(slide, "Pearson、Spearman、Partial correlation")

    rows = [
        ("pg.corr(x, y, method='pearson')", "兩變項，含 CI 與 BF10"),
        ("pg.corr(x, y, method='spearman')", "Rank-based，robust to outliers"),
        ("pg.partial_corr(data, x, y, covar=)", "控制 covariate 後的相關"),
        ("df.rcorr(padjust='fdr_bh', stars=True)", "Correlation matrix + FDR"),
    ]
    add_table(slide, ["API", "用途"], rows,
              Inches(0.85), Inches(1.95), Inches(11.6), Inches(2.7),
              size=14)

    add_callout(slide,
                "應用：individual differences 研究 — "
                "working memory (WM) 與 fluid intelligence (Gf) 的關係，"
                "控制 age / education 等 covariates。",
                Inches(0.85), Inches(5.05), Inches(11.6), Inches(1.50),
                fill=BG_OFFWHITE, border=ACCENT_NAVY, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "Partial correlation 在發展心理學、認知老化研究中超常用 —"
              "幾乎所有跨年齡層的相關研究都要控制 age 作為 covariate。")


def slide_16_handson4(prs, page, total):
    slide = new_slide(prs, bg=BG_PRACTICE)
    add_section_label(slide, "HANDS-ON 4")
    add_title(slide, "WM-Gf 控制不同 covariates")

    add_callout(slide, "10 min", Inches(11.0), Inches(0.55),
                Inches(1.5), Inches(0.45),
                fill=ACCENT_AMBER, border=ACCENT_AMBER,
                color=TEXT_LIGHT, size=14, bold=True)

    add_bullets(slide,
                ["任務：分別控制 (a) age, (b) edu, (c) 兩者一起",
                 "比較 partial r 在三種設定下的變化",
                 "解讀：哪個變項對 WM-Gf 關係影響最大？"],
                Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.0),
                size=17)

    add_code(slide,
             "pg.partial_corr(data=df, x='wm', y='gf', covar='age')\n"
             "pg.partial_corr(data=df, x='wm', y='gf', covar='edu')\n"
             "pg.partial_corr(data=df, x='wm', y='gf', covar=['age', 'edu'])",
             Inches(0.85), Inches(4.20), Inches(11.6), Inches(2.0),
             size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "若控制某 covariate 後 r 大幅下降，代表原本的相關有一部分是 spurious。"
              "若 r 基本不變，代表該 covariate 不是混淆變項。")


def slide_17_streamlit_why(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§5  STREAMLIT INTEGRATION")
    add_title(slide, "分析跑完了 — 然後呢？")

    add_bullets(slide,
                ["Notebook 只能自己看，PI / collaborator 無法即時探索",
                 "Streamlit：把 DataFrame 直接渲染、加 sidebar filter",
                 "pingouin 回傳 DataFrame ⇒ 直接 st.dataframe()",
                 "部署到 Streamlit Cloud：分享 URL 給合作者"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(3.0),
                size=16)

    # arrow flow — use shape's own text frame so labels are centred
    blocks = [("Jupyter\nNotebook", 0.85),
              ("pingouin\nDataFrame", 5.15),
              ("Streamlit\nApp + URL", 9.45)]
    from pptx.oxml.ns import qn
    for (txt, left) in blocks:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(5.15),
                                   Inches(3.0), Inches(1.30))
        b.fill.solid(); b.fill.fore_color.rgb = ACCENT_NAVY
        b.line.fill.background()
        tf = b.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.10); tf.margin_bottom = Inches(0.10)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, line in enumerate(txt.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.name = FONT_SANS
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = TEXT_LIGHT
            rPr = run._r.get_or_add_rPr()
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_CJK})
            rPr.append(ea)
    # arrows
    for x in [3.95, 8.25]:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     Inches(x), Inches(5.65),
                                     Inches(1.15), Inches(0.30))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT_TEAL
        arr.line.fill.background()
    add_footer(slide, page, total)
    set_notes(slide,
              "強調可重現性 / 開放科學 — 把分析變成可分享的 URL，"
              "符合 preregistration & open analysis 的精神。")


def slide_18_dashboard_layout(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§5  DASHBOARD STRUCTURE")
    add_title(slide, "四個區塊一個 app")

    add_bullets(slide,
                ["Sidebar：資料上傳、RT 過濾 slider、accuracy filter",
                 "主畫面 metrics：n subjects, n trials, n conditions",
                 "Plotly box plot：condition × group",
                 "st.tabs(['t-test', 'ANOVA', 'Correlation']) 切換分析"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(3.2),
                size=16)

    # wireframe — sidebar + main
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.85), Inches(5.15),
                                  Inches(2.5), Inches(1.50))
    side.fill.solid(); side.fill.fore_color.rgb = BG_OFFWHITE
    side.line.color.rgb = HAIRLINE
    add_text(slide, "sidebar\n(filters)", Inches(0.85), Inches(5.40),
             Inches(2.5), Inches(1.0),
             size=12, color=TEXT_MUTED, font=FONT_SANS,
             align=PP_ALIGN.CENTER)
    main = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(3.55), Inches(5.15),
                                  Inches(9.0), Inches(1.50))
    main.fill.solid(); main.fill.fore_color.rgb = BG_WHITE
    main.line.color.rgb = HAIRLINE
    add_text(slide, "metrics  |  plotly chart  |  tabs (t-test · ANOVA · corr)",
             Inches(3.55), Inches(5.40), Inches(9.0), Inches(1.0),
             size=12, color=TEXT_MUTED, font=FONT_SANS,
             align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    set_notes(slide,
              "參考 05_streamlit_stats_app.py。學生只需照範例改 condition 名稱、"
              "新增 tab 就能改成自己的研究。")


def slide_19_code_walkthrough(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§5  CODE WALKTHROUGH")
    add_title(slide, "100 行串接整個 pipeline")

    add_code(slide,
             "import pingouin as pg, plotly.express as px, streamlit as st\n\n"
             "@st.cache_data\n"
             "def simulate_data(n_subj=30, seed=42): ...\n\n"
             "with st.sidebar:\n"
             "    uploaded = st.file_uploader('upload CSV', type=['csv'])\n"
             "    rt_lower = st.slider('RT lower', 100, 400, 200)\n\n"
             "tab1, tab2, tab3 = st.tabs(['t-test', 'ANOVA', 'Correlation'])\n"
             "with tab2:\n"
             "    aov = pg.mixed_anova(data=subj_mean, dv='rt',\n"
             "                          within='condition', between='group',\n"
             "                          subject='subject')\n"
             "    st.dataframe(aov.round(4))\n"
             "    st.download_button('下載 CSV', aov.to_csv().encode(),\n"
             "                       file_name='anova.csv')",
             Inches(0.85), Inches(1.85), Inches(11.6), Inches(4.85),
             size=13)
    add_footer(slide, page, total)
    set_notes(slide,
              "重點：pingouin 的 DataFrame 直接餵給 st.dataframe，"
              "再用 to_csv().encode() 餵給 st.download_button。100 行可以做出一個完整 app。")


def slide_20_handson5(prs, page, total):
    slide = new_slide(prs, bg=BG_PRACTICE)
    add_section_label(slide, "HANDS-ON 5")
    add_title(slide, "加入 Non-parametric tab")

    add_callout(slide, "15 min", Inches(11.0), Inches(0.55),
                Inches(1.5), Inches(0.45),
                fill=ACCENT_AMBER, border=ACCENT_AMBER,
                color=TEXT_LIGHT, size=14, bold=True)

    add_bullets(slide,
                ["在現有 05_streamlit_stats_app.py 加第四個 tab",
                 "跑 Wilcoxon signed-rank test",
                 "改 st.tabs([...]) 列表，新增 'Non-parametric'"],
                Inches(0.85), Inches(1.85), Inches(11.5), Inches(2.0),
                size=17)

    add_code(slide,
             "tab1, tab2, tab3, tab4 = st.tabs(\n"
             "    ['Paired t-test', 'Mixed ANOVA', 'Correlation', 'Non-parametric']\n"
             ")\n\n"
             "with tab4:\n"
             "    w = pg.wilcoxon(wide['congruent'], wide['incongruent'])\n"
             "    st.dataframe(w.round(4))",
             Inches(0.85), Inches(4.10), Inches(11.6), Inches(2.3),
             size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "讓學生實際感受 Streamlit 的「改一行就更新」開發體驗。"
              "存檔後瀏覽器自動 rerun，不用重啟 server。")


def slide_21_deploy(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "§5  DEPLOY")
    add_title(slide, "三步驟讓全世界看到")

    steps = [
        ("1", "Push to GitHub", "把 .py 與 requirements.txt 推上 repo"),
        ("2", "share.streamlit.io", "New app → 連結 repo → 選 main file"),
        ("3", "Get a public URL", "貼到 supplementary 或寄給 collaborator"),
    ]
    for i, (num, title, sub) in enumerate(steps):
        left = 0.85 + i * 4.20
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(left), Inches(2.20),
                                      Inches(3.90), Inches(3.50))
        card.fill.solid(); card.fill.fore_color.rgb = BG_OFFWHITE
        card.line.color.rgb = HAIRLINE
        # number circle
        ncircle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(left + 0.30),
                                         Inches(2.50),
                                         Inches(0.70), Inches(0.70))
        ncircle.fill.solid(); ncircle.fill.fore_color.rgb = ACCENT_TEAL
        ncircle.line.fill.background()
        add_text(slide, num, Inches(left + 0.30), Inches(2.55),
                 Inches(0.70), Inches(0.60), size=22, bold=True,
                 color=TEXT_LIGHT, align=PP_ALIGN.CENTER, font=FONT_SANS)
        add_text(slide, title, Inches(left + 0.30), Inches(3.45),
                 Inches(3.30), Inches(0.50),
                 size=17, bold=True, color=TEXT_DARK, font=FONT_SANS)
        add_text(slide, sub, Inches(left + 0.30), Inches(3.95),
                 Inches(3.30), Inches(1.5),
                 size=13, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(slide, page, total)
    set_notes(slide,
              "Streamlit Cloud 免費版 1GB RAM，跑這種統計 dashboard 完全夠用。"
              "Build 約 2-3 分鐘 (要裝 pingouin)，之後就能秒開。")


def slide_22_pitfalls(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "PITFALLS")
    add_title(slide, "5 個容易踩的雷")

    pitfalls = [
        ("Wide format 餵給 rm_anova",
         "用 df.melt() 轉成 long format"),
        ("只報 p，沒報 effect size",
         "務必附上 cohen_d / np2 / hedges"),
        ("Within-subject 用 paired=False",
         "確認 design 後選正確的 paired 參數"),
        ("Sphericity 違反卻引 p_unc",
         "用 GG 校正後的 p-value (eps < 0.75 時)"),
        ("Pairwise 沒做多重比較校正",
         "加 padjust='bonf' 或 'fdr_bh'"),
    ]
    for i, (bad, good) in enumerate(pitfalls):
        top = 1.85 + i * 0.95
        # red X
        add_text(slide, "✗", Inches(0.85), Inches(top),
                 Inches(0.40), Inches(0.50),
                 size=22, bold=True, color=ACCENT_RED, font=FONT_SANS)
        add_text(slide, bad, Inches(1.30), Inches(top + 0.05),
                 Inches(5.5), Inches(0.50),
                 size=15, color=TEXT_DARK, font=FONT_CJK)
        # green check
        add_text(slide, "✓", Inches(7.05), Inches(top),
                 Inches(0.40), Inches(0.50),
                 size=22, bold=True, color=ACCENT_GREEN, font=FONT_SANS)
        add_text(slide, good, Inches(7.50), Inches(top + 0.05),
                 Inches(5.3), Inches(0.50),
                 size=15, color=TEXT_DARK, font=FONT_CJK)
    add_footer(slide, page, total)
    set_notes(slide,
              "每個錯誤都對應講義 §Recap & Common Pitfalls 區塊。"
              "建議學生把這頁印出來貼在電腦旁邊。")


def slide_23_homework(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "HOMEWORK")
    add_title(slide, "完整的 Flanker 分析 pipeline")

    add_bullets(slide,
                ["資料準備：trial-level CSV (subject, group, condition, rt, accuracy)",
                 "描述性：pg.normality() + pg.homoscedasticity() 檢查假設",
                 "推論：mixed ANOVA (group × condition) + simple effects",
                 "Correlation：Flanker effect × accuracy partial corr (控制 age)",
                 "Streamlit：擴充 dashboard + 部署到 Streamlit Cloud",
                 "APA-style writeup (≤ 200 字)"],
                Inches(0.85), Inches(1.85), Inches(11.6), Inches(4.0),
                size=15)

    add_callout(slide,
                "繳交：week-14-pingouin-hw.ipynb  +  Streamlit Cloud URL\n"
                "評分：可執行 (40%) / 分析正確 (40%) / 報告清晰 (20%)",
                Inches(0.85), Inches(5.85), Inches(11.6), Inches(0.95),
                fill=BG_OFFWHITE, border=ACCENT_NAVY, size=14)
    add_footer(slide, page, total)
    set_notes(slide,
              "Homework 設計成 final project 的暖身 — 同樣的 pipeline 學生可以替換成自己的研究資料。"
              "鼓勵學生把 dashboard URL 放進履歷或 GitHub README。")


def slide_24_refs(prs, page, total):
    slide = new_slide(prs)
    add_section_label(slide, "REFERENCES")
    add_title(slide, "學完還想多看一點")

    refs = [
        ("Vallat, R. (2018). Pingouin: statistics in Python.",
         "JOSS, 3(31), 1026.  doi.org/10.21105/joss.01026"),
        ("Pingouin documentation",
         "https://pingouin-stats.org/"),
        ("Lakens, D. (2013). Calculating and reporting effect sizes.",
         "Frontiers in Psychology, 4, 863."),
        ("Streamlit documentation",
         "https://docs.streamlit.io/"),
    ]
    for i, (title, sub) in enumerate(refs):
        top = 1.95 + i * 1.10
        add_text(slide, title, Inches(0.85), Inches(top),
                 Inches(11.5), Inches(0.55),
                 size=16, bold=True, color=TEXT_DARK, font=FONT_SANS)
        add_text(slide, sub, Inches(0.85), Inches(top + 0.55),
                 Inches(11.5), Inches(0.45),
                 size=13, italic=True, color=TEXT_MUTED, font=FONT_SANS)
    add_footer(slide, page, total)
    set_notes(slide,
              "Lakens (2013) 是必讀；pingouin 官網的 examples gallery 也是很好的入門。")


def slide_25_next_week(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_SECTION)

    add_text(slide, "NEXT WEEK", Inches(0.85), Inches(0.85),
             Inches(11), Inches(0.5), size=14, color=ACCENT_TEAL,
             bold=True, font=FONT_SANS)
    add_text(slide, "Week 15: Final Project Workshop", Inches(0.85), Inches(1.50),
             Inches(11), Inches(1.0), size=36, bold=True,
             color=TEXT_LIGHT, font=FONT_SANS)
    add_text(slide, "把 W11–W14 學的工具整合",
             Inches(0.85), Inches(2.85), Inches(11), Inches(0.6),
             size=22, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC),
             font=FONT_CJK)

    items = [
        "提案題目：行為實驗 + 統計報告 + Streamlit dashboard",
        "整合 Week 7 (NumPy) + 12 (pandas) + 13 (Plotly) + 14 (pingouin)",
        "評分重點：pipeline 完整性、可重現性、視覺化品質",
        "Deliverables：.ipynb + Streamlit URL + 5-min lightning talk",
    ]
    for i, txt in enumerate(items):
        add_text(slide, "•  " + txt, Inches(1.20), Inches(3.85 + i * 0.55),
                 Inches(11), Inches(0.50), size=16,
                 color=TEXT_LIGHT, font=FONT_CJK)

    add_text(slide, FOOTER_STR, Inches(0.85), Inches(7.10),
             Inches(10), Inches(0.30),
             size=10, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_SANS)
    add_text(slide, f"{page} / {total}", Inches(11.85), Inches(7.10),
             Inches(1.30), Inches(0.30),
             size=10, color=RGBColor(0xCA, 0xDC, 0xFC), font=FONT_SANS,
             align=PP_ALIGN.RIGHT)
    set_notes(slide,
              "Final project 是 8 週工具的整合驗收。"
              "建議學生提案時就考慮 Streamlit 部署，這是與一般 ML 課程的差異化亮點。")


# ============================================================
# Build
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_scipy_vs_pg,
        slide_03_objectives,
        slide_04_install,
        slide_05_philosophy,
        slide_06_handson1,
        slide_07_ttest_family,
        slide_08_nonparametric,
        slide_09_effsize,
        slide_10_handson2,
        slide_11_anova_oneway,
        slide_12_rm_anova,
        slide_13_posthoc_mixed,
        slide_14_handson3,
        slide_15_correlation,
        slide_16_handson4,
        slide_17_streamlit_why,
        slide_18_dashboard_layout,
        slide_19_code_walkthrough,
        slide_20_handson5,
        slide_21_deploy,
        slide_22_pitfalls,
        slide_23_homework,
        slide_24_refs,
        slide_25_next_week,
    ]
    total_slides = len(builders)
    for page, fn in enumerate(builders, start=1):
        fn(prs, page, total_slides)

    out = "week-14-pingouin-slides.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({total_slides} slides)")


if __name__ == "__main__":
    build()
