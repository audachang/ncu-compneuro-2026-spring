"""Build week-10 lecture slides as PPTX for students.

Source content: week-10-lecture_plan_150min.md (the 150-minute classroom plan).
Output: student-facing deck — no teacher notes, no risk/fallback tables,
no formative assessment, no pre-class checklists.

Design (from scientific-slides skill):
- Visual-first hierarchy with large fonts (titles 40-60pt, body 24-28pt, code 18-22pt)
- Two-column layouts for OS-specific commands and CLI/web alternatives
- Modern dark-slate + cyan palette (developer/CS aesthetic)
- One topic per slide, generous whitespace

Local dependency (NOT in git — see archive/ in .gitignore):
    archive/extracted_images/slide_{02,04,05,14}.png

These are concept visuals extracted from
archive/Git_Visual_Workflow_Mastery.pptx. To regenerate after a fresh clone:

    python -c "
    import os
    from pptx import Presentation
    p = Presentation('week-10-agentic_workflows_and_github/archive/Git_Visual_Workflow_Mastery.pptx')
    out = 'week-10-agentic_workflows_and_github/archive/extracted_images'
    os.makedirs(out, exist_ok=True)
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if hasattr(sh, 'image'):
                with open(f'{out}/slide_{i:02d}.{sh.image.ext}', 'wb') as f:
                    f.write(sh.image.blob)
    "

The committed week-10-slides.pptx already has these images embedded, so
you only need to re-run build_slides.py if you change content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Color palette ----
BG_LIGHT      = RGBColor(0xF7, 0xF9, 0xFB)
BG_DARK       = RGBColor(0x1E, 0x29, 0x3B)
BG_DIVIDER    = RGBColor(0x0F, 0x76, 0x8F)
BG_BREAK      = RGBColor(0x4A, 0x14, 0x8C)
TEXT_DARK     = RGBColor(0x1E, 0x29, 0x3B)
TEXT_LIGHT    = RGBColor(0xF7, 0xF9, 0xFB)
TEXT_MUTED    = RGBColor(0x55, 0x65, 0x7B)
ACCENT_CYAN   = RGBColor(0x00, 0xB8, 0xD4)
ACCENT_AMBER  = RGBColor(0xFF, 0xB7, 0x00)
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)
CODE_COMMENT  = RGBColor(0x8B, 0x9D, 0xB8)

# ---- Slide size: 16:9 ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- Fonts ----
FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"


# ============================================================
# Primitives
# ============================================================

def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_accent_bar(slide, color=ACCENT_CYAN, top=Inches(0), height=Inches(0.18)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=28, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, items, left, top, width, height, *,
                size=26, color=TEXT_DARK, line_spacing=1.2,
                bullet_char="•", bullet_color=ACCENT_CYAN, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    use_font = font or FONT_CJK
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = bullet_char + "  "
        b.font.name = FONT_SANS
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = use_font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i > 0:
            p.space_before = Pt(size * 0.45)
    return tb


def _color_for_line(line, lang):
    stripped = line.lstrip()
    if not stripped:
        return None  # blank
    if lang in ("bash", "yaml", "python", "html", "markdown") and stripped.startswith("#"):
        return CODE_COMMENT
    if stripped.startswith(">"):
        return ACCENT_AMBER
    if stripped.startswith("$"):
        return ACCENT_GREEN
    return None  # default


def add_code_block(slide, code, left, top, width, height, *,
                   size=20, bg=BG_DARK, fg=TEXT_LIGHT, lang="bash",
                   line_spacing=1.10, padding=0.22):
    """Render a code block in IDE-style dark box.

    Use line_spacing=1.05 and smaller padding for tight code; size=18 for
    long code blocks to avoid overflow.
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.fill.background()
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.margin_left = Inches(0.28)
    tf.margin_right = Inches(0.28)
    tf.margin_top = Inches(padding)
    tf.margin_bottom = Inches(padding)
    tf.word_wrap = True

    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        color = _color_for_line(line, lang) or fg
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_footer(slide, page_num, total, dark=False):
    fg = TEXT_LIGHT if dark else TEXT_MUTED
    add_text(slide, "Week 10 — Agentic Workflows & GitHub",
             Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
             size=14, color=fg)
    add_text(slide, f"{page_num} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.6), Inches(0.35),
             size=14, color=fg, align=PP_ALIGN.RIGHT)


def add_section_label(slide, label_text, color=ACCENT_CYAN):
    add_text(slide, label_text,
             Inches(0.7), Inches(0.5), Inches(8), Inches(0.5),
             size=18, bold=True, color=color)


def slide_full_image(prs, image_path, page, total, *, label=None):
    """Full-bleed image slide. Used to embed pre-rendered concept visuals
    from the archived Git_Visual_Workflow_Mastery deck. Source images are
    16:9 (17.78 x 10), same aspect ratio as our deck, so they fit edge-to-edge.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_DARK)
    s.shapes.add_picture(image_path, 0, 0, SLIDE_W, SLIDE_H)
    # Tiny page indicator in corner — keeps deck navigation consistent
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.6), Inches(0.35),
             size=14, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)
    return s


def add_title(slide, title_zh, *, size=40, top=0.95):
    add_text(slide, title_zh,
             Inches(0.7), Inches(top), Inches(12), Inches(0.9),
             size=size, bold=True, color=TEXT_DARK, font=FONT_CJK)


def add_subtitle(slide, text, *, top=1.55):
    add_text(slide, text,
             Inches(0.7), Inches(top), Inches(12), Inches(0.55),
             size=20, color=TEXT_MUTED, font=FONT_CJK)


# ============================================================
# Slide builders
# ============================================================

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(2.2),
                             Inches(0.3), Inches(2.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    add_text(s, "WEEK 10",
             Inches(1.3), Inches(2.1), Inches(11), Inches(0.7),
             size=28, bold=True, color=ACCENT_CYAN)
    add_text(s, "Agentic Workflows × GitHub",
             Inches(1.3), Inches(2.7), Inches(11.5), Inches(1.1),
             size=52, bold=True, color=TEXT_LIGHT)
    add_text(s, "用 Claude Code 建立資料分析專案環境 SOP",
             Inches(1.3), Inches(3.95), Inches(11.5), Inches(0.7),
             size=28, color=ACCENT_CYAN, font=FONT_CJK)
    add_text(s, "× GitHub Repo & GitHub Pages",
             Inches(1.3), Inches(4.55), Inches(11.5), Inches(0.7),
             size=28, color=ACCENT_CYAN, font=FONT_CJK)

    add_text(s, "NS5116  Programming & AI Applications in Behavioral Science",
             Inches(1.3), Inches(5.85), Inches(11), Inches(0.5),
             size=18, color=TEXT_LIGHT)
    add_text(s, "Spring 2026  •  2026-04-30  •  150 min",
             Inches(1.3), Inches(6.3), Inches(11), Inches(0.5),
             size=16, color=TEXT_MUTED)
    return s


def slide_objectives(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "LEARNING OUTCOMES")
    add_title(s, "本週結束後，你會...", size=42)
    items = [
        "依 SOP 從零建立可重現 (reproducible) 的 Python 資料分析專案",
        "用 Claude Code 自動產生標準專案結構（venv / requirements / src / data）",
        "把本機專案推上 GitHub，理解 repo 是專案的「正式對外形象」",
        "設定 GitHub Pages，把分析結果發布為 https://<user>.github.io/<repo>/",
        "解釋 GitHub repo 與 Pages 在資料科學作品集中的角色",
    ]
    add_bullets(s, items,
                Inches(0.9), Inches(2.1), Inches(12), Inches(4.8),
                size=26, color=TEXT_DARK)
    add_footer(s, page, total)
    return s


def slide_agenda(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "TODAY'S 150 MINUTES")
    add_title(s, "今天的時間表", size=42)

    # 8-row table
    rows = [
        ("0", "為什麼要 SOP & 公開 repo", "10 min", "講授"),
        ("1", "SOP 全景圖（6 步驟）", "10 min", "講授"),
        ("2", "實作 #1：資料夾 / venv / git init", "25 min", "Hands-on"),
        ("3", "Claude Code 產生專案骨架", "10 min", "講授+示範"),
        ("4", "實作 #2：scaffold + 第一份分析", "20 min", "Hands-on"),
        ("—", "休息", "10 min", "—"),
        ("5", "GitHub Repo & README 寫法", "15 min", "講授+示範"),
        ("6", "實作 #3：建遠端 repo & push", "15 min", "Hands-on"),
        ("7", "GitHub Pages 把分析變網頁", "15 min", "講授+示範"),
        ("8", "實作 #4：啟用 Pages、發布 index.html", "15 min", "Hands-on"),
        ("9", "Recap、作業說明、Q&A", "5 min", "講授"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 4,
        Inches(0.7), Inches(1.95),
        Inches(11.9), Inches(4.8),
    )
    t = table_shape.table
    t.columns[0].width = Inches(0.9)
    t.columns[1].width = Inches(6.6)
    t.columns[2].width = Inches(1.6)
    t.columns[3].width = Inches(2.8)

    headers = ["#", "內容", "時間", "形式"]
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = BG_DARK
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        is_break = row[1] == "休息"
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            if is_break:
                cell.fill.fore_color.rgb = RGBColor(0xFB, 0xEA, 0xC0)
            else:
                cell.fill.fore_color.rgb = BG_LIGHT if ri % 2 else RGBColor(0xE8, 0xEE, 0xF4)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_CJK
            r.font.size = Pt(18)
            r.font.color.rgb = TEXT_DARK
    add_footer(s, page, total)
    return s


def slide_divider(prs, page, total, number, title_en, title_zh, label="PART"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_DIVIDER)
    add_text(s, f"{label} {number}",
             Inches(0.7), Inches(2.4), Inches(12), Inches(0.7),
             size=28, bold=True, color=ACCENT_AMBER)
    add_text(s, title_en,
             Inches(0.7), Inches(3.0), Inches(12), Inches(1.2),
             size=50, bold=True, color=TEXT_LIGHT)
    add_text(s, title_zh,
             Inches(0.7), Inches(4.5), Inches(12), Inches(0.8),
             size=28, color=TEXT_LIGHT, font=FONT_CJK)
    add_footer(s, page, total, dark=True)
    return s


def slide_break(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_BREAK)
    add_text(s, "BREAK",
             Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
             size=80, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
    add_text(s, "10 分鐘休息",
             Inches(0.7), Inches(3.7), Inches(12), Inches(0.9),
             size=40, color=TEXT_LIGHT, font=FONT_CJK, align=PP_ALIGN.CENTER)
    add_text(s, "回來後：GitHub Repo & GitHub Pages",
             Inches(0.7), Inches(4.7), Inches(12), Inches(0.6),
             size=22, color=ACCENT_CYAN, font=FONT_CJK, align=PP_ALIGN.CENTER)
    return s


# ---- Section 0 ----

def slide_why_sop(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 0 · 開場")
    add_title(s, "一份分析專案 = 4 個元素", size=40)

    # Four pillars as colored boxes
    pillars = [
        ("程式碼",  "code",         ACCENT_CYAN),
        ("環境",    "venv + deps",  ACCENT_GREEN),
        ("資料",    "raw / processed", ACCENT_AMBER),
        ("對外展示", "repo + Pages", RGBColor(0xE0, 0x4A, 0xCC)),
    ]
    for i, (zh, en, color) in enumerate(pillars):
        x = 0.7 + i * 3.05
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(1.95),
                                 Inches(2.85), Inches(2.4))
        box.fill.solid(); box.fill.fore_color.rgb = color
        box.line.fill.background()
        box.adjustments[0] = 0.08
        tf = box.text_frame
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.4); tf.margin_bottom = Inches(0.2)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = zh
        r.font.name = FONT_CJK; r.font.size = Pt(34); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)
        r2 = p2.add_run(); r2.text = en
        r2.font.name = FONT_MONO; r2.font.size = Pt(20)
        r2.font.color.rgb = TEXT_LIGHT

    add_text(s, "四者缺一 → 別人（包括三個月後的你）就重現不了",
             Inches(0.7), Inches(4.7), Inches(12), Inches(0.7),
             size=26, bold=True, color=TEXT_DARK, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_three_tragedies(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s, color=ACCENT_RED)
    add_section_label(s, "SECTION 0 · 痛點")
    add_title(s, "三個你一定遇過的悲劇", size=40)

    tragedies = [
        ("🐛", "「我電腦上可以跑」",
         "→ 沒 venv、沒 requirements.txt，到別人電腦就壞了"),
        ("🤔", "「這 code 是我寫的嗎？」",
         "→ 沒 README、沒 notebook 註解，三個月後自己都看不懂"),
        ("📦", "「老師，作品給你 zip 檔可以嗎？」",
         "→ 沒公開 repo、沒 live demo URL，履歷上沒東西可貼"),
    ]
    y = 2.0
    for icon, problem, cause in tragedies:
        # Icon
        add_text(s, icon,
                 Inches(0.7), Inches(y), Inches(1.0), Inches(1.0),
                 size=44, color=ACCENT_RED, align=PP_ALIGN.CENTER)
        # Quote + cause
        add_text(s, problem,
                 Inches(1.8), Inches(y), Inches(11.0), Inches(0.6),
                 size=26, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, cause,
                 Inches(1.8), Inches(y + 0.55), Inches(11.0), Inches(0.55),
                 size=20, color=TEXT_MUTED, font=FONT_CJK)
        y += 1.55

    add_text(s, "本週解法：用 SOP 把每次新專案都做對 + 用 GitHub 把成果端出去",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             size=20, bold=True, color=ACCENT_CYAN, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


# ---- Section 1: SOP 全景 ----

def slide_sop_overview(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 1 · SOP 全景圖")
    add_title(s, "資料分析專案環境設定 SOP — 6 步驟", size=36)

    steps = [
        ("1", "建立專案資料夾並進入"),
        ("2", "建立並啟用 Python venv"),
        ("3", "git init + 第一個 commit"),
        ("4", "用 Claude Code 產生專案骨架"),
        ("5", "寫第一份分析（notebook 或 script）"),
        ("6", "push 到 GitHub + 啟用 GitHub Pages"),
    ]
    # 2-column 3-row grid
    for i, (num, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.3
        y = 1.85 + row * 1.55
        # number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(x), Inches(y),
                                    Inches(1.0), Inches(1.0))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_CYAN
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = FONT_SANS; r.font.size = Pt(36); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
        # description
        add_text(s, desc,
                 Inches(x + 1.2), Inches(y + 0.18), Inches(5.0), Inches(0.7),
                 size=22, color=TEXT_DARK, font=FONT_CJK)
    add_text(s, "每次新專案都從這 6 步開始 — 不省略",
             Inches(0.7), Inches(6.65), Inches(12), Inches(0.5),
             size=22, bold=True, color=ACCENT_AMBER, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_sop_table(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 1 · SOP 對照表")
    add_title(s, "每一步的「目的 → 工具 → 產出」", size=36)

    rows = [
        ("1", "隔離專案空間",      "mkdir, cd",                    "專案根目錄"),
        ("2", "鎖定套件版本",      "python -m venv, pip",          ".venv/, requirements.txt"),
        ("3", "開始版本控制",      "git init, .gitignore",         ".git/, .gitignore"),
        ("4", "建立標準目錄",      "Claude Code",                  "src/, data/, notebooks/, README"),
        ("5", "產出可重現分析",    "Jupyter / Python",             "analysis.ipynb / report.html"),
        ("6", "對外發布",          "GitHub + Pages",               "repo URL + *.github.io URL"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 4,
        Inches(0.5), Inches(1.95),
        Inches(12.3), Inches(4.7),
    )
    t = table_shape.table
    t.columns[0].width = Inches(0.7)
    t.columns[1].width = Inches(3.0)
    t.columns[2].width = Inches(3.8)
    t.columns[3].width = Inches(4.8)
    headers = ["#", "目的", "工具", "產出"]
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = BG_DARK
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if ri % 2 else RGBColor(0xE8, 0xEE, 0xF4)
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            is_code = ci in (2, 3)
            r.font.name = FONT_MONO if is_code else FONT_CJK
            r.font.size = Pt(18 if is_code else 19)
            r.font.color.rgb = TEXT_DARK
    add_footer(s, page, total)
    return s


# ---- Section 2: 實作 #1 ----

def slide_step1(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 2 · STEP 1")
    add_title(s, "建立專案資料夾", size=40)
    add_subtitle(s, "範例專案名稱：tw-airquality-mini")
    code = """# 進入你習慣的工作目錄
cd ~/projects

# 建立並進入新資料夾
mkdir tw-airquality-mini
cd tw-airquality-mini

# 確認位置
pwd"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.3),
                   Inches(11.9), Inches(4.3),
                   size=24, lang="bash")
    add_footer(s, page, total)
    return s


def slide_step2_venv(prs, page, total):
    """Two-column: Windows | macOS/Linux."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 2 · STEP 2")
    add_title(s, "建立並啟用 Python venv", size=40)
    add_subtitle(s, "Windows 與 macOS/Linux 啟用指令不同 — 認清自己的系統")

    # Left column: Windows
    add_text(s, "Windows (PowerShell / cmd)",
             Inches(0.6), Inches(2.15), Inches(6), Inches(0.5),
             size=22, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    win_code = """# 建立 venv
python -m venv .venv

# 啟用
.venv\\Scripts\\activate

# 提示符會出現 (.venv)
(.venv) C:\\>"""
    add_code_block(s, win_code,
                   Inches(0.6), Inches(2.65),
                   Inches(6.0), Inches(3.6),
                   size=20, lang="bash")

    # Right column: macOS/Linux
    add_text(s, "macOS / Linux (bash / zsh)",
             Inches(6.85), Inches(2.15), Inches(6), Inches(0.5),
             size=22, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    mac_code = """# 建立 venv
python -m venv .venv

# 啟用
source .venv/bin/activate

# 提示符會出現 (.venv)
(.venv) $"""
    add_code_block(s, mac_code,
                   Inches(6.85), Inches(2.65),
                   Inches(6.0), Inches(3.6),
                   size=20, lang="bash")

    add_text(s, "驗收：執行 which python（macOS）或 where python（Win）→ 路徑應指向 .venv",
             Inches(0.7), Inches(6.45), Inches(12), Inches(0.45),
             size=18, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_step2_install(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 2 · STEP 2 (續)")
    add_title(s, "安裝套件 + 凍結版本", size=40)
    code = """# 升級 pip
pip install --upgrade pip

# 安裝最小套件集合
pip install pandas matplotlib jupyter

# 把目前環境凍結成 requirements.txt
pip freeze > requirements.txt

# 確認
cat requirements.txt   # macOS/Linux
type requirements.txt  # Windows"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.05),
                   Inches(11.9), Inches(4.5),
                   size=22, lang="bash")
    add_text(s, "為什麼要 freeze？→ 別人 clone 後可以 pip install -r requirements.txt 復刻你的環境",
             Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_step3(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 2 · STEP 3")
    add_title(s, "git init + 第一個 commit", size=40)
    code = """# 初始化 git
git init
git branch -M main

# 建立 .gitignore（先暫時手寫，之後交給 Claude Code 完整版）
echo ".venv/"               > .gitignore
echo "__pycache__/"         >> .gitignore
echo ".ipynb_checkpoints/"  >> .gitignore

# 加入並提交
git add .gitignore requirements.txt
git commit -m "Initial commit: venv and gitignore"

# 確認
git log --oneline"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.05),
                   Inches(11.9), Inches(4.6),
                   size=20, lang="bash")
    add_text(s, ".venv/ 永遠不入 git — 它太大，而且每個人的系統會自己重建",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_section2_check(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s, color=ACCENT_GREEN)
    add_section_label(s, "SECTION 2 · 驗收", color=ACCENT_GREEN)
    add_title(s, "Step 1–3 完成 — 三個檢查點", size=38)

    items = [
        ("which python  /  where python", "顯示路徑包含 .venv"),
        ("git log --oneline",              "看到第一個 commit"),
        ("git status",                     ".venv/ 不在追蹤清單中"),
    ]
    y = 2.1
    for cmd, expect in items:
        # check icon
        add_text(s, "✓",
                 Inches(0.7), Inches(y), Inches(0.7), Inches(0.8),
                 size=40, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
        # cmd
        add_text(s, cmd,
                 Inches(1.5), Inches(y - 0.05), Inches(7), Inches(0.55),
                 size=24, color=TEXT_DARK, font=FONT_MONO, bold=True)
        # expectation
        add_text(s, "→ " + expect,
                 Inches(1.5), Inches(y + 0.55), Inches(11), Inches(0.5),
                 size=20, color=TEXT_MUTED, font=FONT_CJK)
        y += 1.4
    add_footer(s, page, total)
    return s


# ---- Section 3: Claude Code scaffold ----

def slide_why_scaffold(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 3 · 為什麼用 Claude Code")
    add_title(s, "不要再手工建 10 個資料夾", size=40)
    add_subtitle(s, "把這件事交給 Claude Code，每次都用同一份 prompt")

    add_text(s, "❌ 手工建立的問題",
             Inches(0.7), Inches(2.15), Inches(6), Inches(0.5),
             size=24, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_bullets(s,
                ["每次都要記得建哪些資料夾", "README 樣板要重寫一次",
                 ".gitignore 漏東漏西", "新成員照抄你的舊專案"],
                Inches(0.7), Inches(2.7), Inches(6), Inches(3.5),
                size=22, color=TEXT_DARK)

    add_text(s, "✅ Claude Code scaffold",
             Inches(7.0), Inches(2.15), Inches(6), Inches(0.5),
             size=24, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_bullets(s,
                ["一份 prompt 解決所有事", "每次都得到一致的結構",
                 ".gitignore / README 直接生成", "新專案 30 秒搞定"],
                Inches(7.0), Inches(2.7), Inches(6), Inches(3.5),
                size=22, color=TEXT_DARK)

    add_text(s, "Scaffold prompt = 專案啟動的「儀式」 — 每個新專案都從這份開始",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.55),
             size=20, bold=True, color=ACCENT_CYAN, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_scaffold_prompt(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 3 · Scaffold Prompt")
    add_title(s, "標準骨架 prompt（直接複製）", size=36)

    code = """> Scaffold a data analysis project in the current directory:
    data/
      raw/         (read-only original data; add .gitkeep)
      processed/   (cleaned outputs; add .gitkeep)
    notebooks/
      01_explore.ipynb   (starter notebook with title + pandas import)
    src/
      __init__.py
      load_data.py       (stub: load_csv(path) -> pd.DataFrame)
      clean_data.py      (stub: clean(df) -> pd.DataFrame)
    reports/             (HTML / PDF outputs; add .gitkeep)
    README.md            (title, one-paragraph description, How to run)
    .gitignore           (Python, Jupyter, venv, OS, data/raw/*.csv)
    requirements.txt     (keep existing entries, just confirm)

  Do not overwrite requirements.txt. Use placeholder content."""
    add_code_block(s, code,
                   Inches(0.5), Inches(1.95),
                   Inches(12.3), Inches(4.9),
                   size=17, lang="bash", line_spacing=1.12)
    add_footer(s, page, total)
    return s


# ---- Section 3.5: Prompt 寫法對比（分析空氣品質）----

def slide_prompt_two_styles(prs, page, total):
    """Concept overview: two prompt styles for the same analysis task."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 3 · Prompt 寫法 vs 產出")
    add_title(s, "同樣分析空氣品質,prompt 寫法決定品質", size=34)
    add_subtitle(s, "1 句話 vs 結構化指令 — 結果差 10 倍,後面兩頁示範")

    # Left card: Vague
    card_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.6), Inches(2.3),
                                Inches(6.05), Inches(4.25))
    card_l.fill.solid(); card_l.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEC)
    card_l.line.fill.background()
    card_l.adjustments[0] = 0.04

    add_text(s, "❌  Vague",
             Inches(0.85), Inches(2.45), Inches(5.5), Inches(0.55),
             size=26, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_text(s, "「幫我分析這份空氣品質資料」",
             Inches(0.85), Inches(3.05), Inches(5.5), Inches(0.55),
             size=18, color=TEXT_DARK, font=FONT_CJK)
    add_bullets(s,
                ["1 句話、無明確產出",
                 "Claude 自由發揮、你猜對方猜",
                 "沒有檢查點 → 跑完才發現方向錯",
                 "重跑結果不一樣、無法驗收"],
                Inches(0.85), Inches(3.75), Inches(5.5), Inches(2.6),
                size=17, color=TEXT_DARK, bullet_color=ACCENT_RED,
                line_spacing=1.18)

    # Right card: Structured
    card_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.85), Inches(2.3),
                                Inches(6.05), Inches(4.25))
    card_r.fill.solid(); card_r.fill.fore_color.rgb = RGBColor(0xE6, 0xF7, 0xEF)
    card_r.line.fill.background()
    card_r.adjustments[0] = 0.04

    add_text(s, "✅  Structured",
             Inches(7.10), Inches(2.45), Inches(5.5), Inches(0.55),
             size=26, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_text(s, "Goal · Inputs · Steps · Checkpoints · Constraints",
             Inches(7.10), Inches(3.05), Inches(5.7), Inches(0.55),
             size=15, color=TEXT_DARK, font=FONT_SANS)
    add_bullets(s,
                ["明確 deliverable + 明確輸入欄位",
                 "分步驟、在關鍵點要 Claude 暫停",
                 "限制套件、隨機種子 → 可重現",
                 "整段 prompt 可以存進 README 重用"],
                Inches(7.10), Inches(3.75), Inches(5.5), Inches(2.6),
                size=17, color=TEXT_DARK, bullet_color=ACCENT_GREEN,
                line_spacing=1.18)

    add_text(s, "原則:Prompt 像「上司交辦」,不是「亂丟一句話」",
             Inches(0.7), Inches(6.65), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_CYAN, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_vague_prompt(prs, page, total):
    """Vague prompt example with what-you-get outcomes."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s, color=ACCENT_RED)
    add_section_label(s, "SECTION 3 · ❌ Vague prompt", color=ACCENT_RED)
    add_title(s, "「幫我看一下這份空氣品質資料」", size=34)
    add_subtitle(s, "1 句話 → Claude 自由發揮 → 通常不會是你想要的")

    # Left: prompt
    add_text(s, "你打的 prompt",
             Inches(0.6), Inches(2.05), Inches(6.0), Inches(0.5),
             size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
    prompt = """> Analyze the air quality data
  in data/raw/ and make a plot."""
    add_code_block(s, prompt,
                   Inches(0.6), Inches(2.55),
                   Inches(6.1), Inches(1.6),
                   size=18, lang="bash", line_spacing=1.15)

    # Right: outcomes
    add_text(s, "Claude 給你什麼",
             Inches(6.85), Inches(2.05), Inches(6.0), Inches(0.5),
             size=20, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_bullets(s,
                ["隨機挑 1 個測站(沒問你哪一個)",
                 "可能用 pandas、seaborn、plotly 任一套",
                 "Y 軸沒單位、X 軸日期亂跳",
                 "圖檔存哪? 不知道、下次再跑也找不到",
                 "你看不懂程式邏輯、無法 commit"],
                Inches(6.85), Inches(2.55), Inches(6.0), Inches(3.85),
                size=17, color=TEXT_DARK, bullet_color=ACCENT_RED,
                line_spacing=1.20)

    # Bottom: typical conversation cost
    add_text(s, "💸  代價",
             Inches(0.6), Inches(4.45), Inches(6.0), Inches(0.5),
             size=20, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    cost_code = """# 實際對話通常長這樣
> 不是這個測站,我要新店
> 改成 PM2.5 不要 PM10
> 圖太醜了再改一次
> 為什麼每次顏色不一樣..."""
    add_code_block(s, cost_code,
                   Inches(0.6), Inches(4.95),
                   Inches(6.1), Inches(1.7),
                   size=15, lang="bash", line_spacing=1.15, padding=0.16)

    add_text(s, "結果:來回 5 次澄清,還是改不到你要的圖",
             Inches(0.7), Inches(6.75), Inches(12), Inches(0.4),
             size=17, color=ACCENT_RED, font=FONT_CJK, align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_structured_prompt(prs, page, total):
    """Structured prompt example with what-you-get outcomes."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s, color=ACCENT_GREEN)
    add_section_label(s, "SECTION 3 · ✅ Structured prompt", color=ACCENT_GREEN)
    add_title(s, "目標 / 輸入 / 步驟 / 檢查點 / 限制", size=34)
    add_subtitle(s, "同一筆 PM2.5 data,改寫成可驗收的 prompt")

    # Left: structured prompt (smaller font, more lines)
    prompt = """> Goal: 4-station PM2.5 daily trend (2025).

  Inputs:
    - data/raw/{中山,古亭,新店,板橋}_2025.csv
    - long format: 測站, 日期, 測項, 00..23

  Steps (PAUSE after step 3):
    1. Load the 4 CSVs, keep 測項 == 'PM2.5'.
    2. Melt 00..23 → long, compute daily mean
       per station.
    3. Print mean / p95 / days_above_35 per
       station — wait for OK.
    4. Plot 4 lines on one axes, save to
       reports/pm25_compare.png (DPI 150).

  Constraints:
    - pandas + matplotlib only (no seaborn).
    - Add 35 µg/m³ horizontal reference line.
    - np.random.seed(42) for any synthetic."""
    add_code_block(s, prompt,
                   Inches(0.5), Inches(2.05),
                   Inches(7.4), Inches(4.65),
                   size=13, lang="bash", line_spacing=1.10, padding=0.16)

    # Right: outcomes
    add_text(s, "Claude 給你什麼",
             Inches(8.05), Inches(2.05), Inches(5.0), Inches(0.5),
             size=20, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_bullets(s,
                ["明確 4 個測站 — 不會自己選錯",
                 "Step 3 暫停 → 你看到表格才決定要不要畫圖",
                 "圖檔路徑固定 → 直接 git add 進 commit",
                 "pandas + seed 固定 → 同學重跑結果一致",
                 "整段 prompt 可貼進 notebook 開頭當紀錄"],
                Inches(8.05), Inches(2.55), Inches(5.0), Inches(4.0),
                size=15, color=TEXT_DARK, bullet_color=ACCENT_GREEN,
                line_spacing=1.22)

    add_text(s, "5 個錨點:Goal · Inputs · Steps · Checkpoints · Constraints",
             Inches(0.7), Inches(6.75), Inches(12), Inches(0.4),
             size=17, bold=True, color=ACCENT_GREEN, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


# ---- Section 4: 實作 #2 ----

def slide_task_a_scaffold(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 4 · 任務 A")
    add_title(s, "跑 scaffold + 驗收結構", size=40)
    add_subtitle(s, "對自己的 tw-airquality-mini 跑 scaffold prompt")

    code = """# 1. 啟動 Claude Code
claude

# 2. 貼上 scaffold prompt（前一頁那段）
#    觀察 Claude 列出將建立的檔案 → 暫停 review

# 3. 確認後執行，檢查目錄
tree -L 2          # macOS/Linux
ls -R              # Windows
git status         # 應該看到一堆新檔案"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.2),
                   Inches(11.9), Inches(4.4),
                   size=22, lang="bash")
    add_footer(s, page, total)
    return s


def slide_task_b_first_analysis(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 4 · 任務 B")
    add_title(s, "第一份分析 + commit", size=40)

    # Left: prompt
    add_text(s, "1. 請 Claude 寫 notebook",
             Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             size=22, bold=True, color=ACCENT_CYAN, font=FONT_CJK)
    prompt_code = """> In notebooks/01_explore.ipynb,
  add cells that:
  1. Import pandas + matplotlib.
  2. Create a fake DataFrame with columns:
     date (10 days), station ('Taipei'),
     pm25 (random integers 10-80).
  3. Plot a line chart of pm25 over date.
  4. Save the plot to
     reports/pm25_demo.png."""
    add_code_block(s, prompt_code,
                   Inches(0.6), Inches(2.5),
                   Inches(6.1), Inches(3.9),
                   size=17, lang="bash", line_spacing=1.10)

    # Right: commit
    add_text(s, "2. 跑 notebook + commit",
             Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.5),
             size=22, bold=True, color=ACCENT_CYAN, font=FONT_CJK)
    commit_code = """# 在 Jupyter 執行所有 cell
# 確認 reports/pm25_demo.png 產生

git add .
git status   # 確認要進的檔案

git commit -m "Scaffold project + first \\
              demo notebook"

git log --oneline"""
    add_code_block(s, commit_code,
                   Inches(6.85), Inches(2.5),
                   Inches(6.0), Inches(3.9),
                   size=18, lang="bash", line_spacing=1.10)

    add_text(s, "Jupyter kernel 接不到 .venv？→ pip install ipykernel 然後重啟",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ---- Section 5: GitHub Repo ----

def slide_repo_is_namecard(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 5 · GitHub Repo")
    add_title(s, "Repo 是你的對外名片", size=42)
    add_subtitle(s, "招生委員、PI、雇主第一個會點的就是 repo URL")

    items = [
        "Repo 的「外觀分數」來自 → README 第一屏 + commit 訊息 + 檔案命名",
        "沒有 README 的 repo = 沒有名片，連網址都沒人想點",
        "Commit 訊息寫成「update」「fix」= 自己也找不到改了什麼",
        "從本週起，每個作業都繳交「GitHub repo URL + live demo URL」",
    ]
    add_bullets(s, items,
                Inches(0.9), Inches(2.4), Inches(12), Inches(4.5),
                size=24, color=TEXT_DARK)
    add_footer(s, page, total)
    return s


def slide_create_repo_two_ways(prs, page, total):
    """Two-column: Web vs CLI."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 5 · 建立遠端 Repo")
    add_title(s, "兩種方式建 repo", size=40)
    add_subtitle(s, "本機已有專案 → 不要勾「Add a README」")

    # Left: Web
    add_text(s, "🌐 網頁建立",
             Inches(0.6), Inches(2.15), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_CYAN, font=FONT_CJK)
    web_code = """# 1. GitHub → New repository
# 2. 名稱填 tw-airquality-mini
# 3. Public，不勾 Add README

# 之後在本機加遠端：
git remote add origin \\
  https://github.com/<user>/\\
  tw-airquality-mini.git

git push -u origin main"""
    add_code_block(s, web_code,
                   Inches(0.6), Inches(2.65),
                   Inches(6.0), Inches(3.7),
                   size=18, lang="bash", line_spacing=1.10)

    # Right: gh CLI (recommended)
    add_text(s, "⚡ gh CLI (推薦)",
             Inches(6.85), Inches(2.15), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    cli_code = """# 一行完成：建 repo + 設遠端 + push
gh repo create \\
  tw-airquality-mini \\
  --public \\
  --source=. \\
  --remote=origin \\
  --push

# 前提：先 gh auth login 過"""
    add_code_block(s, cli_code,
                   Inches(6.85), Inches(2.65),
                   Inches(6.0), Inches(3.7),
                   size=18, lang="bash", line_spacing=1.10)

    add_text(s, "Push 被拒（remote 有 README）→ git pull --rebase origin main 後再 push",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_readme_template(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 5 · README 寫法")
    add_title(s, "最小可行 README 架構", size=38)

    code = """# TW Air Quality Mini

一個示範用的台灣空氣品質迷你分析專案，作為 NS5116 Week 10 練習。

## How to run

git clone https://github.com/<user>/tw-airquality-mini.git
cd tw-airquality-mini
python -m venv .venv && source .venv/bin/activate
pip install -r requirements.txt
jupyter notebook notebooks/01_explore.ipynb

## Project structure

- data/raw/    — 原始 CSV（不入 git）
- notebooks/   — 探索性分析
- src/         — 共用函數
- reports/     — 產生的圖表與 HTML

## Author

Erik Chang, NCU 認知神經科學研究所"""
    add_code_block(s, code,
                   Inches(0.5), Inches(1.95),
                   Inches(12.3), Inches(4.95),
                   size=17, lang="markdown", line_spacing=1.10)
    add_footer(s, page, total)
    return s


def slide_readme_principles(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 5 · README 原則")
    add_title(s, "三個 README 寫作原則", size=42)
    items = [
        ("1", "第一段一句話講清楚這是什麼",
         "讀者 3 秒內要知道：這個 repo 在做什麼"),
        ("2", "一定要有 How to run",
         "包含 clone / venv / install / 執行的完整步驟"),
        ("3", "圖表 / 截圖 放最上面",
         "讓滑過的人 3 秒內看到成果，不用 clone 才知道好不好"),
    ]
    y = 2.0
    for num, title, desc in items:
        # number
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.7), Inches(y),
                                    Inches(0.95), Inches(0.95))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_AMBER
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = FONT_SANS; r.font.size = Pt(36); r.font.bold = True
        r.font.color.rgb = TEXT_DARK
        # title + desc
        add_text(s, title,
                 Inches(2.0), Inches(y), Inches(11), Inches(0.55),
                 size=24, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, desc,
                 Inches(2.0), Inches(y + 0.55), Inches(11), Inches(0.5),
                 size=18, color=TEXT_MUTED, font=FONT_CJK)
        y += 1.45
    add_footer(s, page, total)
    return s


# ---- Section 6: 實作 #3 ----

def slide_push_repo(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 6 · 實作 #3")
    add_title(s, "建遠端 repo & push", size=40)
    code = """# 推薦：gh CLI 一條龍
gh repo create tw-airquality-mini \\
  --public --source=. --remote=origin --push

# 備案：網頁建空 repo 後手動接遠端
git remote add origin https://github.com/<user>/tw-airquality-mini.git
git push -u origin main

# 確認
git remote -v
git log --oneline"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.05),
                   Inches(11.9), Inches(4.5),
                   size=20, lang="bash", line_spacing=1.10)
    add_text(s, "驗收：開啟 https://github.com/<user>/tw-airquality-mini，README 在首頁完整呈現",
             Inches(0.7), Inches(6.65), Inches(12), Inches(0.45),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_improve_readme(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 6 · 用 Claude Code 改 README")
    add_title(s, "讓 AI 把 README 寫漂亮", size=40)

    # Left: prompt
    add_text(s, "Prompt",
             Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    prompt_code = """> Improve README.md:
  - add a one-paragraph
    description
  - add a "How to run" section
    with venv + pip + jupyter
    commands
  - list the project structure
  - add a placeholder image link
    for reports/pm25_demo.png"""
    add_code_block(s, prompt_code,
                   Inches(0.6), Inches(2.5),
                   Inches(6.1), Inches(3.9),
                   size=18, lang="bash", line_spacing=1.10)

    # Right: commit + push
    add_text(s, "Commit + Push",
             Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    push_code = """# 檢查 Claude 改了什麼
git diff README.md

# 滿意就 commit + push
git add README.md
git commit -m "Improve README \\
              with how-to-run \\
              and structure"
git push"""
    add_code_block(s, push_code,
                   Inches(6.85), Inches(2.5),
                   Inches(6.0), Inches(3.9),
                   size=20, lang="bash", line_spacing=1.10)
    add_footer(s, page, total)
    return s


# ---- Section 7: GitHub Pages ----

def slide_pages_concept(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 7 · GitHub Pages")
    add_title(s, "免費、零部署成本的靜態網站", size=38)
    add_subtitle(s, "把 reports/ 裡的 HTML 直接公開出去")

    add_text(s, "適用場景",
             Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             size=22, bold=True, color=ACCENT_CYAN, font=FONT_CJK)
    items = [
        "個人作品集    →  https://<user>.github.io",
        "專案展示頁    →  https://<user>.github.io/<repo>/",
        "課堂作業 live demo（取代「附件 zip」）",
    ]
    add_bullets(s, items,
                Inches(0.9), Inches(2.7), Inches(12), Inches(2.0),
                size=23, color=TEXT_DARK, font=FONT_MONO,
                line_spacing=1.4)

    # Three publish methods
    add_text(s, "三種發布方式（本課用 docs/，最乾淨）",
             Inches(0.7), Inches(5.0), Inches(12), Inches(0.5),
             size=22, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    methods = [
        ("docs/ 子目錄",     "main 分支放 docs/index.html",  "✓ 推薦"),
        ("gh-pages branch", "另開分支放網頁檔",              "本週不教"),
        ("main root",       "整個 main 當網站根目錄",        "亂"),
    ]
    y = 5.55
    for name, desc, tag in methods:
        add_text(s, f"• {name}",
                 Inches(0.9), Inches(y), Inches(3.2), Inches(0.45),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, desc,
                 Inches(4.2), Inches(y), Inches(6.0), Inches(0.45),
                 size=17, color=TEXT_MUTED, font=FONT_CJK)
        add_text(s, tag,
                 Inches(10.4), Inches(y), Inches(2.5), Inches(0.45),
                 size=17, bold=True,
                 color=ACCENT_GREEN if "✓" in tag else TEXT_MUTED,
                 font=FONT_CJK)
        y += 0.4
    add_footer(s, page, total)
    return s


def slide_pages_steps(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 7 · 發布步驟")
    add_title(s, "用 docs/ 發布 — 4 個步驟", size=40)

    steps = [
        ("1", "在 repo 根目錄建 docs/ 資料夾"),
        ("2", "把 reports/analysis.html 複製成 docs/index.html"),
        ("3", "GitHub → Settings → Pages → Source: 'Deploy from a branch'"),
        ("4", "Branch: main，Folder: /docs → Save"),
    ]
    y = 1.95
    for num, desc in steps:
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.7), Inches(y),
                                    Inches(0.85), Inches(0.85))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_CYAN
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.name = FONT_SANS; r.font.size = Pt(28); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
        add_text(s, desc,
                 Inches(1.85), Inches(y + 0.13), Inches(11.4), Inches(0.6),
                 size=22, color=TEXT_DARK, font=FONT_CJK)
        y += 1.05

    add_text(s, "等 1–2 分鐘 → 拜訪 https://<user>.github.io/<repo>/",
             Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
             size=22, bold=True, color=ACCENT_AMBER, font=FONT_MONO,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


# ---- Section 8: 實作 #4 ----

def slide_nbconvert(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 8 · 任務 1")
    add_title(s, "把 notebook 轉成 HTML", size=40)

    code = """# 在專案根目錄
mkdir -p docs

# 用 jupyter nbconvert 把 notebook 轉成 HTML
# 並輸出到 docs/index.html
jupyter nbconvert --to html notebooks/01_explore.ipynb \\
                  --output index.html \\
                  --output-dir docs/

# 確認
ls docs/                           # macOS/Linux
dir docs\\                          # Windows"""
    add_code_block(s, code,
                   Inches(0.7), Inches(2.05),
                   Inches(11.9), Inches(4.4),
                   size=20, lang="bash", line_spacing=1.10)
    add_text(s, "找不到 nbconvert？→ pip install nbconvert 然後加進 requirements.txt",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             size=18, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_pages_landing(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 8 · 任務 2")
    add_title(s, "或讓 Claude Code 寫一個簡單首頁", size=38)

    # Left: prompt
    add_text(s, "Prompt",
             Inches(0.6), Inches(2.0), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    prompt_code = """> Create docs/index.html:
  - minimal landing page
  - project title
  - one-paragraph description
  - embedded image
    (../reports/pm25_demo.png
    copied to docs/)
  - link to notebook on GitHub
  - plain HTML + tiny inline CSS
    for readable typography"""
    add_code_block(s, prompt_code,
                   Inches(0.6), Inches(2.5),
                   Inches(6.1), Inches(3.9),
                   size=17, lang="bash", line_spacing=1.10)

    # Right: commit + push
    add_text(s, "Commit + Push",
             Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.5),
             size=24, bold=True, color=ACCENT_CYAN, font=FONT_SANS)
    push_code = """# 把 docs/ 整個目錄加進來
git add docs/

git commit -m "Add docs/index.html \\
              for GitHub Pages"

git push

# 然後到 Settings → Pages
# 啟用後等 1-2 分鐘"""
    add_code_block(s, push_code,
                   Inches(6.85), Inches(2.5),
                   Inches(6.0), Inches(3.9),
                   size=18, lang="bash", line_spacing=1.10)
    add_footer(s, page, total)
    return s


def slide_pages_enable(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s)
    add_section_label(s, "SECTION 8 · 啟用 Pages")
    add_title(s, "Settings → Pages", size=42)
    add_subtitle(s, "在 GitHub 網頁設定，存檔後等 1–2 分鐘")

    settings = [
        ("Source",  "Deploy from a branch"),
        ("Branch",  "main"),
        ("Folder",  "/docs"),
        ("Action",  "Save"),
    ]
    y = 2.5
    for label, value in settings:
        # label box
        lb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(2.0), Inches(y),
                                Inches(2.5), Inches(0.7))
        lb.fill.solid(); lb.fill.fore_color.rgb = BG_DARK
        lb.line.fill.background()
        tf = lb.text_frame
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        r = p.add_run(); r.text = label
        r.font.name = FONT_SANS; r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
        # value box
        vb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(4.5), Inches(y),
                                Inches(7.0), Inches(0.7))
        vb.fill.solid(); vb.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
        vb.line.color.rgb = RGBColor(0xCB, 0xD3, 0xDD)
        tf = vb.text_frame
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        r = p.add_run(); r.text = value
        r.font.name = FONT_MONO; r.font.size = Pt(22)
        r.font.color.rgb = TEXT_DARK
        y += 0.85

    add_text(s, "✅  Your site is live at https://<user>.github.io/<repo>/",
             Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
             size=22, bold=True, color=ACCENT_GREEN, font=FONT_MONO,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


# ---- Section 9: Recap & 作業 ----

def slide_recap(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_DARK)
    add_text(s, "RECAP",
             Inches(0.7), Inches(0.6), Inches(12), Inches(0.7),
             size=22, bold=True, color=ACCENT_AMBER)
    add_text(s, "本週你走過的完整 SOP",
             Inches(0.7), Inches(1.1), Inches(12), Inches(0.9),
             size=40, bold=True, color=TEXT_LIGHT, font=FONT_CJK)

    flow = [
        "資料夾",
        "venv",
        "git",
        "Claude Code\nscaffold",
        "分析\n(notebook)",
        "GitHub\nrepo",
        "GitHub\nPages",
    ]
    n = len(flow)
    box_w = 1.55
    gap = 0.18
    total_w = n * box_w + (n - 1) * gap
    start_x = (13.333 - total_w) / 2
    y = 3.0
    for i, label in enumerate(flow):
        x = start_x + i * (box_w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y),
                                 Inches(box_w), Inches(1.3))
        box.fill.solid(); box.fill.fore_color.rgb = ACCENT_CYAN
        box.line.fill.background()
        box.adjustments[0] = 0.12
        tf = box.text_frame
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
        for j, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.05
            r = p.add_run(); r.text = line
            r.font.name = FONT_CJK; r.font.size = Pt(16); r.font.bold = True
            r.font.color.rgb = BG_DARK
        # arrow
        if i < n - 1:
            arrow_x = x + box_w + 0.01
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    Inches(arrow_x), Inches(y + 0.5),
                                    Inches(gap - 0.02), Inches(0.3))
            ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT_AMBER
            ar.line.fill.background()

    add_text(s, "從今天起，你的每個分析作業都應該是",
             Inches(0.7), Inches(5.2), Inches(12), Inches(0.6),
             size=24, color=TEXT_LIGHT, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_text(s, "「可 clone、可重現、可線上看」的 repo",
             Inches(0.7), Inches(5.85), Inches(12), Inches(0.7),
             size=32, bold=True, color=ACCENT_AMBER, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    return s


def slide_assignment(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_LIGHT)
    add_accent_bar(s, color=ACCENT_AMBER)
    add_section_label(s, "ASSIGNMENT", color=ACCENT_AMBER)
    add_title(s, "本週作業：你的第一個發布專案", size=38)
    add_subtitle(s, "把今天範例改造成你自己的迷你資料分析 repo（不要重用 tw-airquality-mini）")

    items = [
        "任選一個感興趣的小資料集（Week 09 台灣 open data 或自己的 CSV）",
        "依 SOP 6 步驟建立全新 repo",
        "README 至少包含：標題、一段描述、How to run、結構、作者",
        "在 notebook 完成至少 1 個資料探索圖表",
        "啟用 GitHub Pages，把 notebook 轉 HTML 發布",
        "繳交：repo URL + Pages URL（兩個都要能打開）",
    ]
    add_bullets(s, items,
                Inches(0.9), Inches(2.45), Inches(12), Inches(4.0),
                size=22, color=TEXT_DARK,
                bullet_char="✓", bullet_color=ACCENT_AMBER)
    add_text(s, "⏰  繳交期限：Week 11 上課前",
             Inches(0.9), Inches(6.6), Inches(12), Inches(0.5),
             size=22, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_next(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_DARK)
    add_text(s, "下週預告",
             Inches(0.7), Inches(0.6), Inches(12), Inches(1.0),
             size=44, bold=True, color=ACCENT_CYAN, font=FONT_CJK)
    rows = [
        ("11", "Streamlit Cloud 部署（互動 app vs. 靜態 Pages 對比）"),
        ("12", "Taiwan open data APIs — fetch / parse / clean"),
        ("13", "Interactive dashboards with Plotly Express"),
        ("16", "Final milestone:  Live web app presentation"),
    ]
    y = 1.95
    for wk, title in rows:
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.9), Inches(y),
                                    Inches(0.95), Inches(0.95))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_CYAN
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = wk
        r.font.name = FONT_SANS; r.font.size = Pt(28); r.font.bold = True
        r.font.color.rgb = BG_DARK
        add_text(s, title,
                 Inches(2.1), Inches(y + 0.18), Inches(11), Inches(0.7),
                 size=24, color=TEXT_LIGHT, font=FONT_CJK)
        y += 1.13

    add_text(s, "Thank you  •  Questions?",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.6),
             size=24, bold=True, color=ACCENT_AMBER,
             align=PP_ALIGN.CENTER)
    return s


# ============================================================
# Compose
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Plan total page count first
    total = 47   # adjust if you add/remove slides
    p = 1
    archive_dir = "week-10-agentic_workflows_and_github/archive/extracted_images"

    slide_title(prs); p += 1
    slide_objectives(prs, p, total); p += 1
    slide_agenda(prs, p, total); p += 1

    # Section 0
    slide_divider(prs, p, total, 0,
                  "Why SOP & Public Repo",
                  "為什麼需要 SOP 與公開 repo",
                  label="SECTION"); p += 1
    slide_why_sop(prs, p, total); p += 1
    slide_three_tragedies(prs, p, total); p += 1

    # Section 1
    slide_divider(prs, p, total, 1,
                  "SOP Overview",
                  "環境設定 SOP 全景圖（6 步驟）",
                  label="SECTION"); p += 1
    slide_sop_overview(prs, p, total); p += 1
    slide_sop_table(prs, p, total); p += 1

    # Section 2
    slide_divider(prs, p, total, 2,
                  "Hands-on #1 · Folder, venv, git",
                  "實作 #1：Step 1–3",
                  label="SECTION"); p += 1
    slide_step1(prs, p, total); p += 1
    slide_step2_venv(prs, p, total); p += 1
    slide_step2_install(prs, p, total); p += 1

    # ----- Git mental-model inserts (from archived Visual Git deck) -----
    # Before students run git init/add/commit, give them the "what is git
    # actually doing?" picture: 4 areas → staging → commit.
    slide_full_image(prs, f"{archive_dir}/slide_02.png", p, total); p += 1
    slide_full_image(prs, f"{archive_dir}/slide_04.png", p, total); p += 1
    slide_full_image(prs, f"{archive_dir}/slide_05.png", p, total); p += 1

    slide_step3(prs, p, total); p += 1
    slide_section2_check(prs, p, total); p += 1

    # Cheatsheet matrix (Intent → Command → Data Flow) — useful reference
    # card students can screenshot for the rest of the course.
    slide_full_image(prs, f"{archive_dir}/slide_14.png", p, total); p += 1

    # Section 3
    slide_divider(prs, p, total, 3,
                  "Claude Code Scaffold",
                  "用 Claude Code 產生專案骨架",
                  label="SECTION"); p += 1
    slide_why_scaffold(prs, p, total); p += 1
    slide_scaffold_prompt(prs, p, total); p += 1

    # ----- Prompt-style contrast (analysis prompts) -----
    # Bridge from scaffold prompt (structured by example) to hands-on
    # analysis: show students that the same discipline applies when they
    # ask Claude to *analyze* the data, not just to scaffold the project.
    slide_prompt_two_styles(prs, p, total); p += 1
    slide_vague_prompt(prs, p, total); p += 1
    slide_structured_prompt(prs, p, total); p += 1

    # Section 4
    slide_divider(prs, p, total, 4,
                  "Hands-on #2 · Scaffold + First Analysis",
                  "實作 #2：scaffold + 第一份分析",
                  label="SECTION"); p += 1
    slide_task_a_scaffold(prs, p, total); p += 1
    slide_task_b_first_analysis(prs, p, total); p += 1

    # Break
    slide_break(prs, p, total); p += 1

    # Section 5
    slide_divider(prs, p, total, 5,
                  "GitHub Repo & README",
                  "GitHub Repo 與 README 寫法",
                  label="SECTION"); p += 1
    slide_repo_is_namecard(prs, p, total); p += 1
    slide_create_repo_two_ways(prs, p, total); p += 1
    slide_readme_template(prs, p, total); p += 1
    slide_readme_principles(prs, p, total); p += 1

    # Section 6
    slide_divider(prs, p, total, 6,
                  "Hands-on #3 · Push to GitHub",
                  "實作 #3：建遠端 repo 並 push",
                  label="SECTION"); p += 1
    slide_push_repo(prs, p, total); p += 1
    slide_improve_readme(prs, p, total); p += 1

    # Section 7
    slide_divider(prs, p, total, 7,
                  "GitHub Pages",
                  "把分析結果變成網頁",
                  label="SECTION"); p += 1
    slide_pages_concept(prs, p, total); p += 1
    slide_pages_steps(prs, p, total); p += 1

    # Section 8
    slide_divider(prs, p, total, 8,
                  "Hands-on #4 · Enable Pages",
                  "實作 #4：啟用 Pages、發布 index.html",
                  label="SECTION"); p += 1
    slide_nbconvert(prs, p, total); p += 1
    slide_pages_landing(prs, p, total); p += 1
    slide_pages_enable(prs, p, total); p += 1

    # Section 9 / Recap
    slide_recap(prs, p, total); p += 1
    slide_assignment(prs, p, total); p += 1
    slide_next(prs, p, total); p += 1

    actual = p - 1
    print(f"Built {actual} slides (footer total constant = {total})")

    import os
    target = "week-10-agentic_workflows_and_github/week-10-slides.pptx"
    out = target
    # If the canonical file is locked (PowerPoint open), fall back to a
    # side filename so the rebuild still succeeds; the user can replace
    # the original after closing PowerPoint.
    try:
        if os.path.exists(target):
            with open(target, "ab"):
                pass
    except PermissionError:
        out = "week-10-agentic_workflows_and_github/week-10-slides-NEW.pptx"
        print(f"[locked] {target} is open in another app — writing to {out} instead")
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
