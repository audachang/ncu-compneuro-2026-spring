"""Build Week 11 lecture slides (Streamlit × Cognitive Aging Dashboard).

Visual style: ACL@NCU clean academic — white background, navy + teal accent
identity colors, generous whitespace, restrained geometric decoration.

Output: week-11-slides.pptx (16:9, ~38-42 slides for a 150-min class)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# ACL@NCU palette — clean white + identity colors
# ============================================================
BG_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE   = RGBColor(0xF7, 0xF8, 0xFA)
BG_SECTION    = RGBColor(0x14, 0x32, 0x5C)   # deep navy (NCU-ish)
BG_BREAK      = RGBColor(0xFB, 0xEA, 0xC0)
TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED    = RGBColor(0x5F, 0x6B, 0x83)
ACCENT_PRIMARY = RGBColor(0x14, 0x32, 0x5C)  # navy
ACCENT_TEAL    = RGBColor(0x0D, 0x9B, 0x9B)  # teal — ACL identity
ACCENT_AMBER   = RGBColor(0xE8, 0xA1, 0x2A)
ACCENT_RED     = RGBColor(0xD3, 0x4F, 0x4F)
ACCENT_GREEN   = RGBColor(0x2E, 0x8B, 0x57)
CODE_BG        = RGBColor(0x1E, 0x29, 0x3B)
CODE_COMMENT   = RGBColor(0x8B, 0x9D, 0xB8)
HAIRLINE       = RGBColor(0xE2, 0xE6, 0xEC)

# Slide size: 16:9 wide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"


# ============================================================
# Primitives
# ============================================================
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_top_band(slide, color=ACCENT_TEAL, height=Inches(0.14)):
    """Thin top accent band — ACL identity element."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=24, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_section_label(slide, text, color=ACCENT_TEAL):
    """Small uppercase eyebrow label above slide title."""
    add_text(slide, text,
             Inches(0.7), Inches(0.55), Inches(11), Inches(0.35),
             size=14, bold=True, color=color, font=FONT_SANS)


def add_title(slide, title_zh, *, size=36, top=0.95, color=TEXT_DARK):
    add_text(slide, title_zh,
             Inches(0.7), Inches(top), Inches(12), Inches(0.85),
             size=size, bold=True, color=color, font=FONT_CJK)


def add_subtitle(slide, text, *, top=1.6, color=TEXT_MUTED):
    add_text(slide, text,
             Inches(0.7), Inches(top), Inches(12), Inches(0.5),
             size=18, color=color, font=FONT_CJK)


def add_bullets(slide, items, left, top, width, height, *,
                size=22, color=TEXT_DARK, line_spacing=1.25,
                bullet_color=ACCENT_TEAL, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    use_font = font or FONT_CJK
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = "▸  "
        b.font.name = FONT_SANS
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = use_font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i > 0:
            p.space_before = Pt(size * 0.45)
    return tb


def _color_for_line(line, lang):
    s = line.lstrip()
    if not s:
        return None
    if lang in ("bash", "yaml", "python", "markdown") and s.startswith("#"):
        return CODE_COMMENT
    if s.startswith(">"):
        return ACCENT_AMBER
    if s.startswith("$"):
        return ACCENT_GREEN
    return None


def add_code(slide, code, left, top, width, height, *,
             size=18, lang="python", line_spacing=1.15, padding=0.22):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(padding); tf.margin_bottom = Inches(padding)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        c = _color_for_line(line, lang) or TEXT_LIGHT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(size)
        r.font.color.rgb = c
    return box


def add_footer(slide, page, total):
    add_text(slide, "Week 11 — Streamlit × Cognitive Aging  ·  ACL@NCU",
             Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
             size=11, color=TEXT_MUTED)
    add_text(slide, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, *, fill=BG_OFFWHITE,
             accent=None, accent_w=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = HAIRLINE
    card.line.width = Pt(0.5)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(accent_w), height)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return card


# ============================================================
# Slide builders
# ============================================================
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))

    # Side accent block
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(2.5),
                               Inches(0.18), Inches(2.6))
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT_TEAL
    block.line.fill.background()

    add_text(s, "WEEK 11",
             Inches(1.1), Inches(2.4), Inches(11), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, "Web App Development",
             Inches(1.1), Inches(2.95), Inches(12), Inches(0.95),
             size=46, bold=True, color=TEXT_DARK)
    add_text(s, "with Streamlit",
             Inches(1.1), Inches(3.85), Inches(12), Inches(0.85),
             size=42, bold=True, color=ACCENT_TEAL)
    add_text(s, "從 Jupyter 到互動式 Dashboard  —  用認知老化資料",
             Inches(1.1), Inches(4.85), Inches(12), Inches(0.6),
             size=22, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s, "NS5116 · Programming & AI Applications in Behavioral Science",
             Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
             size=14, color=TEXT_MUTED)
    add_text(s, "ACL@NCU  ·  Spring 2026  ·  2026-05-07  ·  150 min",
             Inches(1.1), Inches(6.4), Inches(11), Inches(0.4),
             size=13, color=TEXT_MUTED)
    return s


def slide_objectives(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE)
    add_top_band(s)
    add_section_label(s, "LEARNING OBJECTIVES")
    add_title(s, "本週結束後，你將能夠⋯")
    items = [
        "解釋 Streamlit 的 rerun-on-interaction 執行模型",
        "建立含 title / metric / dataframe 的多區塊 app",
        "加入 selectbox / slider / multiselect 等互動式 widgets",
        "整合 Matplotlib 與 Streamlit 內建 chart",
        "用 columns / sidebar / tabs 規劃 dashboard 版面",
        "套用 @st.cache_data 提升效能",
        "把 app 部署到 Streamlit Community Cloud",
    ]
    add_bullets(s, items,
                Inches(0.9), Inches(2.0), Inches(12), Inches(5),
                size=22, color=TEXT_DARK)
    add_footer(s, page, total)
    return s


def slide_agenda(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE)
    add_top_band(s)
    add_section_label(s, "TODAY'S 150 MINUTES")
    add_title(s, "今天的時間表")

    rows = [
        ("0", "Why Streamlit? (Jupyter → web app)",        "10 min", "Lecture"),
        ("1", "Mental model + st.write + magic",            "15 min", "Lecture"),
        ("2", "認識 cognitive aging dataset",                "10 min", "Lecture"),
        ("3", "Hands-on #1：第一個 app",                    "15 min", "Practice"),
        ("4", "Widgets：selectbox / slider / multiselect",  "15 min", "Lecture"),
        ("5", "Hands-on #2：加入篩選器",                     "15 min", "Practice"),
        ("—", "休息",                                        "10 min", "—"),
        ("6", "Charts：matplotlib + 內建 chart",             "15 min", "Lecture"),
        ("7", "Layout：columns / sidebar / tabs",            "10 min", "Lecture"),
        ("8", "@st.cache_data 與效能",                       "10 min", "Lecture"),
        ("9", "Hands-on #3：完整 dashboard",                 "15 min", "Practice"),
        ("10", "Deploy 到 Streamlit Cloud",                  "10 min", "Demo"),
        ("11", "Recap、Q&A、作業",                           "5 min",  "Lecture"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 4,
        Inches(0.7), Inches(1.95),
        Inches(11.9), Inches(4.95),
    )
    t = table_shape.table
    t.columns[0].width = Inches(0.7)
    t.columns[1].width = Inches(7.0)
    t.columns[2].width = Inches(1.5)
    t.columns[3].width = Inches(2.7)

    headers = ["#", "內容", "時間", "形式"]
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT

    for ri, row in enumerate(rows, start=1):
        is_break = row[1] == "休息"
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            if is_break:
                cell.fill.fore_color.rgb = BG_BREAK
            else:
                cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_CJK
            r.font.size = Pt(13)
            r.font.color.rgb = TEXT_DARK
    add_footer(s, page, total)
    return s


def slide_divider(prs, page, total, num, label_en, title_zh):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_SECTION)
    add_text(s, f"PART {num}",
             Inches(0.7), Inches(2.6), Inches(12), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, label_en,
             Inches(0.7), Inches(3.2), Inches(12.5), Inches(0.95),
             size=44, bold=True, color=TEXT_LIGHT)
    add_text(s, title_zh,
             Inches(0.7), Inches(4.25), Inches(12.5), Inches(0.6),
             size=22, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)
    return s


def slide_break(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_BREAK)
    add_text(s, "BREAK",
             Inches(0), Inches(2.7), Inches(13.333), Inches(1.3),
             size=72, bold=True, color=ACCENT_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, "10 分鐘休息",
             Inches(0), Inches(4.1), Inches(13.333), Inches(0.7),
             size=28, color=TEXT_DARK, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    return s


# ----- Section 0: Why Streamlit -----
def slide_why_streamlit(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 0 · MOTIVATION")
    add_title(s, "Why Streamlit?")

    # 3 problem cards
    cards = [
        ("01", "Jupyter 不能分享",
         "寄一份 .ipynb 給合作者 — 對方多半打不開、跑不動，最後只看到截圖。"),
        ("02", "PI 想自己探索",
         "想拉拉看 slider 試試不同年齡組？沒人想為了改一個參數開 Jupyter。"),
        ("03", "Final presentation",
         "Week 16 你要展示作品給觀眾看 — 一個 URL 比 GitHub repo 直觀 100 倍。"),
    ]
    x0, y0 = Inches(0.8), Inches(2.3)
    w, h = Inches(3.95), Inches(3.6)
    gap = Inches(0.15)
    for i, (icon, title, body) in enumerate(cards):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=ACCENT_TEAL)
        add_text(s, icon, left + Inches(0.25), y0 + Inches(0.25),
                 Inches(1.5), Inches(0.7), size=28, bold=True,
                 color=ACCENT_TEAL)
        add_text(s, title, left + Inches(0.18), y0 + Inches(0.95),
                 w - Inches(0.4), Inches(0.55),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, left + Inches(0.18), y0 + Inches(1.55),
                 w - Inches(0.4), h - Inches(1.7),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "Streamlit 把同一份 Python 分析變成網頁 — 全程不用 HTML / CSS / JS。",
             Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
             size=18, italic=True, color=ACCENT_PRIMARY, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_footer(s, page, total)
    return s


def slide_compare_frameworks(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 0 · MOTIVATION")
    add_title(s, "Streamlit vs. Flask vs. Dash")

    rows = [
        ("學習曲線",
         "平 — 純 Python", "陡 — HTML/CSS/JS/routing", "中 — 需懂 callback"),
        ("寫法",
         "top-to-bottom script", "route decorator + template", "callback graph"),
        ("適合場景",
         "資料分析 dashboard / demo", "通用 web app / API server", "複雜互動圖表"),
        ("部署",
         "Streamlit Cloud (免費)", "自架 / PaaS", "自架 / Dash Enterprise"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 4,
        Inches(0.7), Inches(2.0),
        Inches(11.9), Inches(3.6),
    )
    t = table_shape.table
    t.columns[0].width = Inches(2.0)
    t.columns[1].width = Inches(3.4)
    t.columns[2].width = Inches(3.3)
    t.columns[3].width = Inches(3.2)

    headers = ["", "Streamlit", "Flask", "Dash (Plotly)"]
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_TEAL if ci == 1 else ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_SANS; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_CJK
            r.font.size = Pt(13)
            r.font.color.rgb = ACCENT_TEAL if ci == 1 and ri > 0 else TEXT_DARK
            r.font.bold = (ci == 1)

    add_card(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.95),
             accent=ACCENT_TEAL)
    add_text(s,
             "結論：對研究者最低摩擦 — 你已經會 Python + pandas，再學 5–10 個 st.* 函式就能做 web app。",
             Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.55),
             size=15, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_streamlit_hello(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 0 · MOTIVATION")
    add_title(s, "先看官方 demo：streamlit hello")

    add_code(s, "$ streamlit hello",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.7),
             size=22, lang="bash")

    add_text(s, "啟動四個內建範例，先掌握 Streamlit 的能力範圍：",
             Inches(0.7), Inches(2.85), Inches(12), Inches(0.5),
             size=16, color=TEXT_MUTED, font=FONT_CJK)

    demos = [
        ("01", "Animated chart",
         "即時更新的折線圖 + progress bar — slider 控制 freq"),
        ("02", "Mapping demo",
         "互動式地圖視覺化（Uber 紐約叫車資料）"),
        ("03", "DataFrame demo",
         "用 multiselect 篩選國家、即時更新表格與圖"),
        ("04", "Self-driving car",
         "上傳影像 → 即時跑 object detection model"),
    ]
    y = 3.5
    x_left, w_card = Inches(0.7), Inches(5.85)
    gap = Inches(0.15)
    for i, (num, title, desc) in enumerate(demos):
        col = i % 2
        row = i // 2
        left = x_left + col * (w_card + gap)
        top = Inches(y + row * 1.45)
        add_card(s, left, top, w_card, Inches(1.3), accent=ACCENT_TEAL)
        add_text(s, num, left + Inches(0.2), top + Inches(0.15),
                 Inches(0.7), Inches(0.4),
                 size=18, bold=True, color=ACCENT_TEAL)
        add_text(s, title, left + Inches(0.85), top + Inches(0.13),
                 Inches(4.8), Inches(0.45),
                 size=16, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, desc, left + Inches(0.85), top + Inches(0.6),
                 Inches(4.9), Inches(0.65),
                 size=12, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_st_write(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 1 · MENTAL MODEL")
    add_title(s, "st.write — 一個函式應付所有顯示")

    code = """st.write("# 這是 markdown 標題")           # markdown
st.write("一段純文字。")                       # plain text
st.write({"a": 1, "b": [2, 3]})               # dict → JSON tree
st.write(pd.DataFrame({"x": [1, 2, 3]}))      # interactive table
st.write(fig)                                 # matplotlib figure
st.write(plotly_fig)                          # plotly chart"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.7),
             size=15)

    add_card(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(1.85),
             accent=ACCENT_TEAL)
    add_text(s, "為什麼方便？",
             Inches(0.95), Inches(5.1), Inches(11.5), Inches(0.5),
             size=17, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s,
             "st.write 自動偵測物件型別 → DataFrame、figure、dict、字串都不用個別記 API。"
             "Prototyping 時很省事；正式 dashboard 仍建議顯式呼叫 st.dataframe / st.pyplot 比較易讀。",
             Inches(0.95), Inches(5.6), Inches(11.5), Inches(1.2),
             size=14, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_magic_commands(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 1 · MENTAL MODEL")
    add_title(s, "Magic Commands — 連 st.write 都省了")

    code = """import streamlit as st
import pandas as pd

df = pd.DataFrame({"age": [25, 60], "rt": [320, 410]})

"# Streamlit Magic"        # 字串 → markdown
df                         # dataframe → 互動表格
fig                        # matplotlib figure → 圖"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.85),
             size=16)

    add_text(s, "Streamlit 把任何「孤立的表達式」自動視為 st.write(...)：",
             Inches(0.7), Inches(5.05), Inches(12), Inches(0.5),
             size=16, color=TEXT_DARK, font=FONT_CJK)
    items = [
        "Prototyping 寫得像 Jupyter — 最快看到結果",
        "正式 dashboard 用 st.dataframe / st.pyplot — 程式較易讀",
    ]
    add_bullets(s, items, Inches(0.9), Inches(5.55), Inches(12), Inches(1.3),
                size=15, line_spacing=1.25)
    add_footer(s, page, total)
    return s


def slide_media_widgets(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 4 · WIDGETS")
    add_title(s, "媒體 widget — 在 dashboard 顯示 stimuli")

    code = """# 行為實驗 dashboard 常見需求：旁邊放當前 trial 的 stimulus
st.image("stim/gabor_45deg.png",
         caption="Gabor patch (45°)", width=240)

st.audio("stim/tone_1khz.wav")

st.video("stim/biological_motion.mp4")"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.4),
             size=16)

    add_card(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.2),
             accent=ACCENT_TEAL)
    add_text(s, "三個 API 接受多種輸入：",
             Inches(0.95), Inches(4.75), Inches(11), Inches(0.5),
             size=17, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = [
        "File path（本機檔案）— 最常見",
        "URL（外部 host 的圖 / 音訊）",
        "bytes / numpy array — 動態生成的 stimulus 也能直接餵",
        "適合作 trial-by-trial 行為資料 dashboard 的 stimulus 預覽",
    ]
    add_bullets(s, items, Inches(0.95), Inches(5.25), Inches(11.5), Inches(1.5),
                size=14, line_spacing=1.2, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_status_messages(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 4 · WIDGETS")
    add_title(s, "Status messages — 視覺化的使用者回饋")

    code = """st.success("Loaded 400 participants successfully.")
st.warning("12 trials had RT < 200 ms — excluded as outliers.")
st.error("File missing required column: 'reaction_time_ms'")
st.info("Tip: Slide age range below 30 to see young adult subsample.")

with st.spinner("Running RSA computation..."):
    rsm = compute_rsm(df)
st.success("Done!")

bar = st.progress(0)
for i in range(100):
    bar.progress(i + 1)

st.balloons()       # 任務完成的小彩蛋"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(3.85),
             size=14)

    add_card(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.85),
             accent=ACCENT_TEAL)
    add_text(s,
             "把 success / warning / error / info 當作 dashboard 與使用者溝通的回饋管道 — 比靜默更新或丟例外好。",
             Inches(0.95), Inches(6.15), Inches(11.5), Inches(0.55),
             size=14, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_text_family(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 1 · MENTAL MODEL")
    add_title(s, "完整 text 顯示家族")

    rows = [
        ("st.title", "頁面主標題", "32-pt 級"),
        ("st.header", "section heading", "24-pt 級"),
        ("st.subheader", "subsection heading", "20-pt 級"),
        ("st.markdown", "完整 markdown 含 LaTeX/HTML", "** _ # 等都支援"),
        ("st.caption", "灰色小字註腳", "圖下方 caption / 註解"),
        ("st.code", "語法 highlight 程式碼塊", "language='python'"),
        ("st.latex", "LaTeX 數學式渲染", r"\sum_i x_i^2"),
        ("st.write", "萬能 — 自動偵測型別", "fallback default"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 3,
        Inches(0.7), Inches(2.0),
        Inches(11.9), Inches(4.6),
    )
    t = table_shape.table
    t.columns[0].width = Inches(2.7)
    t.columns[1].width = Inches(5.3)
    t.columns[2].width = Inches(3.9)
    for ci, h in enumerate(["API", "用途", "備註 / 範例"]):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_MONO if ci == 0 else FONT_CJK
            r.font.size = Pt(13)
            r.font.color.rgb = ACCENT_TEAL if ci == 0 else TEXT_DARK
            r.font.bold = (ci == 0)
    add_text(s, "Prototyping 用 st.write 即可；正式 dashboard 用對的 API 可控制視覺層級。",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             size=14, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_cache_resource(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 8 · CACHING")
    add_title(s, "@st.cache_data vs. @st.cache_resource")

    rows = [
        ("@st.cache_data",
         "可序列化資料",
         "DataFrame, str, dict, ndarray",
         "CSV / API / 計算結果"),
        ("@st.cache_resource",
         "不可序列化的資源物件",
         "DB connection, ML model, tokenizer",
         "joblib.load(...) sklearn model"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 4,
        Inches(0.7), Inches(2.0),
        Inches(11.9), Inches(2.4),
    )
    t = table_shape.table
    t.columns[0].width = Inches(3.0)
    t.columns[1].width = Inches(2.7)
    t.columns[2].width = Inches(3.4)
    t.columns[3].width = Inches(2.8)
    for ci, h in enumerate(["Decorator", "適用", "型別範例", "情境"]):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_MONO if ci in (0, 2) else FONT_CJK
            r.font.size = Pt(12)
            r.font.color.rgb = TEXT_DARK
            r.font.bold = (ci == 0)

    code = """@st.cache_resource
def load_model(path):
    import joblib
    return joblib.load(path)

model = load_model("models/cognitive_age_predictor.pkl")"""
    add_code(s, code, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.7),
             size=14)

    add_text(s,
             "判斷原則：你想要拿到「同一個物件實例」（cache_resource）還是「相同內容的副本」（cache_data）？",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             size=14, italic=True, color=ACCENT_PRIMARY, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_multipage(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 9 · MULTI-PAGE")
    add_title(s, "Multi-page app — dashboard → ML app")

    code = """mode = st.sidebar.selectbox("Page",
    ["Overview", "EDA", "Predict"])

if mode == "Overview":
    show_overview(df)
elif mode == "EDA":
    show_eda(df)
else:
    show_prediction(df, model)"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(7.5), Inches(3.05),
             size=14)

    # Right column — page architecture
    add_text(s, "三個 page 的角色分工",
             Inches(8.4), Inches(2.0), Inches(4.5), Inches(0.4),
             size=16, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = [
        "Overview：總結 metric + 主要圖表",
        "EDA：互動式探索（filters + scatter）",
        "Predict：使用者輸入 → model 預測",
    ]
    add_bullets(s, items, Inches(8.4), Inches(2.5), Inches(4.5), Inches(2),
                size=14, font=FONT_CJK)

    add_card(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.7),
             accent=ACCENT_AMBER)
    add_text(s, "進階做法",
             Inches(0.95), Inches(5.25), Inches(11), Inches(0.5),
             size=16, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_text(s,
             "Streamlit 內建多頁機制：在 repo 內建 pages/ 資料夾，每個 .py 檔自動成一頁。"
             "對複雜 app 推薦 pages/，對 prototype 用上面 selectbox 即可。",
             Inches(0.95), Inches(5.7), Inches(11.5), Inches(1.0),
             size=13, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 1: Mental model -----
def slide_mental_model(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 1 · MENTAL MODEL")
    add_title(s, "每次互動 = 整支 app.py 從頭跑到尾")

    code = """import streamlit as st

count = st.slider("Pick a number", 0, 100, 50)
st.write(f"You picked {count}")"""
    add_code(s, code, Inches(0.7), Inches(2.05), Inches(7.5), Inches(1.95),
             size=20)

    # Right column — explanation
    add_text(s, "拖動 slider 時，Streamlit 會：",
             Inches(8.5), Inches(2.05), Inches(4.5), Inches(0.5),
             size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
    items = [
        "重新執行整支 app.py",
        "st.slider() 回傳新值",
        "st.write() 渲染新文字",
    ]
    add_bullets(s, items,
                Inches(8.5), Inches(2.65), Inches(4.5), Inches(2),
                size=16, line_spacing=1.3, font=FONT_CJK)

    add_card(s, Inches(0.7), Inches(4.4), Inches(12), Inches(2.2),
             accent=ACCENT_AMBER)
    add_text(s, "重點",
             Inches(0.95), Inches(4.55), Inches(11), Inches(0.5),
             size=18, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_text(s,
             "沒有 callback、沒有 event handler、沒有 onClick —— 你寫程式的方式跟 Jupyter 一樣（由上到下），"
             "卻得到一個互動式網頁。每個 widget 呼叫立刻回傳當前值。",
             Inches(0.95), Inches(5.1), Inches(11.5), Inches(1.4),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_install_run(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 1 · MENTAL MODEL")
    add_title(s, "安裝與執行")

    add_text(s, "1. 安裝", Inches(0.7), Inches(2.0), Inches(11), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_code(s, "$ pip install streamlit",
             Inches(0.7), Inches(2.45), Inches(12), Inches(0.7),
             size=20, lang="bash")

    add_text(s, "2. 執行", Inches(0.7), Inches(3.4), Inches(11), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_code(s, "$ streamlit run app.py",
             Inches(0.7), Inches(3.85), Inches(12), Inches(0.7),
             size=20, lang="bash")

    add_text(s, "3. 瀏覽器自動開啟",
             Inches(0.7), Inches(4.85), Inches(11), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, "→ http://localhost:8501",
             Inches(0.95), Inches(5.3), Inches(11), Inches(0.5),
             size=20, color=TEXT_DARK, font=FONT_MONO)
    add_text(s, "存檔後右上角會出現 *Rerun* 與 *Always rerun* 按鈕，立即看到改動結果。",
             Inches(0.95), Inches(5.85), Inches(12), Inches(0.5),
             size=15, color=TEXT_MUTED, italic=True, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 2: Dataset -----
def slide_dataset_intro(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 2 · DATASET")
    add_title(s, "Cognitive Aging Dataset (n=400)")

    add_text(s, "今天的範例資料 — 模仿 Cambridge Brain Sciences / Taiwan Biobank 認知測驗結構：",
             Inches(0.7), Inches(1.95), Inches(12), Inches(0.5),
             size=16, color=TEXT_MUTED, font=FONT_CJK)

    rows = [
        ("subject_id",             "S001 – S400"),
        ("age",                    "20 – 80 (years)"),
        ("sex",                    "F / M"),
        ("education",              "Years of formal education (9–22)"),
        ("group",                  "young / middle / older"),
        ("reaction_time_ms",       "Simple RT (lower = faster)"),
        ("working_memory_span",    "n-back / digit span (2–9)"),
        ("processing_speed",       "Digit-symbol substitution"),
        ("moca_score",             "Montreal Cognitive Assessment (0–30)"),
        ("stroop_interference_ms", "Incongruent − congruent RT"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 2,
        Inches(0.7), Inches(2.55),
        Inches(11.9), Inches(4.1),
    )
    t = table_shape.table
    t.columns[0].width = Inches(3.8)
    t.columns[1].width = Inches(8.1)

    for ci, h in enumerate(["Column", "Description"]):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_SANS; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, (col, desc) in enumerate(rows, start=1):
        for ci, val in enumerate([col, desc]):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.02)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_MONO if ci == 0 else FONT_SANS
            r.font.size = Pt(12)
            r.font.color.rgb = TEXT_DARK
    add_footer(s, page, total)
    return s


def slide_dataset_pattern(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 2 · DATASET")
    add_title(s, "資料生成符合認知老化文獻的典型模式")

    rows = [
        ("Reaction time",          "↑ ~2 ms / year — Salthouse (1996)"),
        ("Working memory span",    "Peaks ~25, gradual decline after 30"),
        ("Processing speed",       "Linear decline across lifespan"),
        ("MoCA score",             "Ceiling effect; mild decline late life"),
        ("Stroop interference",    "↑ with age — increased control demands"),
        ("Education effect",       "Mild boost on WM & MoCA (cognitive reserve proxy)"),
    ]
    y = 2.0
    for label, desc in rows:
        add_card(s, Inches(0.7), Inches(y), Inches(11.9), Inches(0.7),
                 accent=ACCENT_TEAL)
        add_text(s, label, Inches(0.95), Inches(y + 0.16),
                 Inches(4), Inches(0.4),
                 size=16, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, desc, Inches(5.1), Inches(y + 0.18),
                 Inches(7.5), Inches(0.4),
                 size=15, color=TEXT_MUTED)
        y += 0.78
    add_footer(s, page, total)
    return s


# ----- Section 3: Hands-on #1 -----
def slide_handson_1(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s, ACCENT_AMBER)
    add_section_label(s, "HANDS-ON #1 · 15 MIN", color=ACCENT_AMBER)
    add_title(s, "你的第一個 Streamlit App")

    add_text(s, "目標：顯示資料前 10 列、總受試者數、平均年齡。",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             size=17, color=TEXT_DARK, font=FONT_CJK)

    code = """# practice_step1.py
import streamlit as st
import pandas as pd

st.title("Cognitive Aging Dashboard")
st.write("First look at the dataset.")

df = pd.read_csv("data/cognitive_aging_taiwan.csv")

st.metric("Participants", len(df))
st.metric("Mean age", f"{df['age'].mean():.1f} years")

st.dataframe(df.head(10))"""
    add_code(s, code, Inches(0.7), Inches(2.6), Inches(8.0), Inches(3.7),
             size=14)

    # Checklist
    add_text(s, "檢查清單",
             Inches(8.95), Inches(2.6), Inches(4), Inches(0.4),
             size=16, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = ["看到 title", "看到兩個 metric", "看到互動式 dataframe"]
    add_bullets(s, items, Inches(8.95), Inches(3.1), Inches(4), Inches(2),
                size=14, font=FONT_CJK)

    add_text(s, "$ streamlit run practice_step1.py",
             Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
             size=15, color=ACCENT_GREEN, font=FONT_MONO)
    add_footer(s, page, total)
    return s


# ----- Section 4: Widgets -----
def slide_widgets_overview(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 4 · WIDGETS")
    add_title(s, "Widgets — 每個都回傳一個值")

    widgets = [
        ("st.selectbox",   "單選下拉",        "F"),
        ("st.slider",      "範圍 / 數值",      "(20, 80)"),
        ("st.multiselect", "多選清單",        "['F', 'M']"),
        ("st.checkbox",    "Boolean 開關",    "True"),
        ("st.radio",       "單選按鈕組",       "'overview'"),
        ("st.text_input",  "文字輸入",        "'S001'"),
        ("st.date_input",  "日期 / 範圍",      "(date1, date2)"),
        ("st.file_uploader", "檔案上傳",      "UploadedFile or None"),
    ]
    table_shape = s.shapes.add_table(
        len(widgets) + 1, 3,
        Inches(0.7), Inches(2.0),
        Inches(11.9), Inches(4.6),
    )
    t = table_shape.table
    t.columns[0].width = Inches(3.5)
    t.columns[1].width = Inches(4.2)
    t.columns[2].width = Inches(4.2)
    for ci, h in enumerate(["Widget", "用途", "回傳範例"]):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(widgets, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_MONO if ci != 1 else FONT_CJK
            r.font.size = Pt(13)
            r.font.color.rgb = TEXT_DARK
    add_footer(s, page, total)
    return s


def slide_widgets_code(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 4 · WIDGETS")
    add_title(s, "用 widgets 篩選資料", size=34)

    code = """# 1. Selectbox — 選一個認知測驗
measure = st.selectbox("Cognitive measure",
    options=["reaction_time_ms", "moca_score", "working_memory_span"])

# 2. Range slider — 年齡範圍
age_min, age_max = st.slider("Age range", 20, 80, (20, 80))

# 3. Multiselect — 性別
sex_choices = st.multiselect("Sex", options=["F", "M"], default=["F", "M"])

# 4. 用 boolean mask 篩選
mask = df["age"].between(age_min, age_max) & df["sex"].isin(sex_choices)
df_filtered = df[mask]"""
    add_code(s, code, Inches(0.7), Inches(2.05), Inches(11.9), Inches(4.4),
             size=14)
    add_text(s, "重點：每個 widget 回傳值 → 存到變數 → 後續用變數篩資料。沒有 callback。",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             size=15, color=ACCENT_PRIMARY, italic=True, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 5: Hands-on #2 -----
def slide_handson_2(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s, ACCENT_AMBER)
    add_section_label(s, "HANDS-ON #2 · 15 MIN", color=ACCENT_AMBER)
    add_title(s, "加入 age + sex 篩選器")

    items = [
        "在 step 1 基礎上加 st.slider 設定年齡範圍 (20–80)",
        "加 st.multiselect 選擇性別",
        "用 boolean mask 篩選 df",
        "用 st.columns(2) 並排顯示兩個 metric (n, mean RT)",
        "顯示篩選後的 df.head(20)",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.05), Inches(11.5), Inches(2.5),
                size=20)

    add_text(s, "提示",
             Inches(0.7), Inches(4.8), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    code = """age_min, age_max = st.slider("Age", 20, 80, (20, 80))
mask = df_all["age"].between(age_min, age_max) & df_all["sex"].isin(sex_choices)
df = df_all[mask]"""
    add_code(s, code, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.4),
             size=14)
    add_footer(s, page, total)
    return s


# ----- Section 6: Charts -----
def slide_charts_builtin(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 6 · CHARTS")
    add_title(s, "Streamlit 內建 charts — 一行搞定")

    code = """# Bar chart：三個年齡組的平均 RT
group_means = df.groupby("group", observed=True)["reaction_time_ms"].mean()
st.bar_chart(group_means)

# Line chart：MoCA 隨年齡變化
st.line_chart(df.set_index("age")["moca_score"])

# Area / scatter / map — 全部都是一行 API
st.scatter_chart(df, x="age", y="reaction_time_ms", color="sex")"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(3.0),
             size=16)

    add_card(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.3),
             accent=ACCENT_GREEN)
    add_text(s, "優點：何時用內建 chart — 快速原型、不需精細控制色彩 / 標籤。",
             Inches(0.95), Inches(5.45), Inches(11.5), Inches(0.5),
             size=16, color=TEXT_DARK, font=FONT_CJK)
    add_text(s, "限制：很難客製 — 加迴歸線、改 legend、調 axis 都要改用 matplotlib。",
             Inches(0.95), Inches(5.95), Inches(11.5), Inches(0.5),
             size=15, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_charts_matplotlib(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 6 · CHARTS")
    add_title(s, "Matplotlib — 完整控制 + st.pyplot(fig)")

    code = """import matplotlib.pyplot as plt
import numpy as np

fig, ax = plt.subplots(figsize=(8, 4))
ax.scatter(df["age"], df["reaction_time_ms"],
           s=20, alpha=0.5, c="#1F6FB4")

slope, intercept = np.polyfit(df["age"], df["reaction_time_ms"], 1)
xs = np.array([df["age"].min(), df["age"].max()])
ax.plot(xs, slope*xs + intercept, "k--",
        label=f"slope={slope:.2f} ms/yr")

ax.set_xlabel("Age"); ax.set_ylabel("Reaction Time (ms)")
ax.legend()

st.pyplot(fig)        # NOTE: 不要用 plt.show()"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.6),
             size=14)
    add_text(s, "常見錯誤：在 Streamlit 內呼叫 plt.show() — 什麼都不會發生。",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             size=15, italic=True, color=ACCENT_RED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 7: Layout -----
def slide_layout_columns(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 7 · LAYOUT")
    add_title(s, "st.columns — 左右排版")

    code = """left, right = st.columns([2, 1])    # 2:1 寬度比

with left:
    st.subheader("Scatter")
    st.pyplot(fig)

with right:
    st.subheader("Stats")
    st.dataframe(df.describe())"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(7.5), Inches(2.6),
             size=16)

    # Visual mockup (right half)
    mock_x, mock_y = Inches(8.6), Inches(2.0)
    add_card(s, mock_x, mock_y, Inches(4.0), Inches(2.6))
    add_card(s, mock_x + Inches(0.15), mock_y + Inches(0.15),
             Inches(2.5), Inches(2.3),
             fill=BG_WHITE, accent=ACCENT_TEAL)
    add_text(s, "Scatter (2x)",
             mock_x + Inches(0.3), mock_y + Inches(1.05),
             Inches(2.3), Inches(0.4),
             size=14, color=TEXT_MUTED, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_card(s, mock_x + Inches(2.75), mock_y + Inches(0.15),
             Inches(1.1), Inches(2.3),
             fill=BG_WHITE, accent=ACCENT_TEAL)
    add_text(s, "Stats",
             mock_x + Inches(2.75), mock_y + Inches(1.05),
             Inches(1.1), Inches(0.4),
             size=12, color=TEXT_MUTED,
             align=PP_ALIGN.CENTER)

    # Tabs preview
    add_text(s, "其他 layout：",
             Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = [
        "st.sidebar — 左側固定面板，常用於 filters",
        "st.tabs([...]) — 多分頁切換 view",
        "st.expander('...') — 折疊區塊（隱藏細節）",
        "st.container() — 邏輯分組",
    ]
    add_bullets(s, items, Inches(0.9), Inches(5.5), Inches(12), Inches(1.4),
                size=16)
    add_footer(s, page, total)
    return s


def slide_layout_sidebar(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 7 · LAYOUT")
    add_title(s, "st.sidebar — 固定篩選器位置")

    code = """with st.sidebar:
    st.header("Filters")
    age_range = st.slider("Age", 20, 80, (20, 80))
    sex       = st.multiselect("Sex", ["F", "M"], default=["F", "M"])
    measure   = st.selectbox("Measure", list(MEASURES.keys()))
    show_reg  = st.checkbox("Show regression line", value=True)"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.4),
             size=16)

    add_card(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.9),
             accent=ACCENT_TEAL)
    add_text(s, "為什麼把 filters 放 sidebar？",
             Inches(0.95), Inches(4.85), Inches(11.5), Inches(0.5),
             size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
    items = [
        "主畫面留給內容（charts、tables）",
        "篩選器永遠看得到，不會被內容擠掉",
        "sidebar widgets 與一般 widgets API 完全相同",
    ]
    add_bullets(s, items, Inches(0.95), Inches(5.4), Inches(11.5), Inches(1.2),
                size=15, line_spacing=1.2)
    add_footer(s, page, total)
    return s


def slide_layout_tabs(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 7 · LAYOUT")
    add_title(s, "st.tabs — 分頁切換不同 view")

    code = """tab1, tab2, tab3 = st.tabs(["Scatter", "Distribution", "Raw data"])

with tab1:
    st.pyplot(fig_scatter)

with tab2:
    st.pyplot(fig_hist)

with tab3:
    st.dataframe(df)
    st.download_button("⬇️ Download CSV",
                       data=df.to_csv(index=False),
                       file_name="cognitive_aging.csv")"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.4),
             size=16)
    add_text(s, "用 emoji 開頭讓 tab label 更易辨識；download_button 直接給使用者下載資料。",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
             size=14, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 8: Caching -----
def slide_caching_problem(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 8 · CACHING")
    add_title(s, "問題：每次互動都重讀 CSV")

    add_text(s, "Streamlit 的「每次重跑整個 script」很方便，但⋯⋯",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             size=18, color=TEXT_DARK, font=FONT_CJK)

    add_card(s, Inches(0.7), Inches(2.65), Inches(11.9), Inches(2.5),
             accent=ACCENT_RED)
    add_text(s, "沒 cache 的情況",
             Inches(0.95), Inches(2.8), Inches(11), Inches(0.5),
             size=17, bold=True, color=ACCENT_RED, font=FONT_CJK)
    items = [
        "每次拖 slider → 整支 app.py rerun → pd.read_csv 重新讀",
        "資料 100 MB？每次互動等 2 秒",
        "API 呼叫？每次重打外部 server，可能還會被 rate-limit",
    ]
    add_bullets(s, items, Inches(0.95), Inches(3.35), Inches(11.5), Inches(1.7),
                size=15, line_spacing=1.25, bullet_color=ACCENT_RED)

    add_card(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.4),
             accent=ACCENT_GREEN)
    add_text(s, "解法：@st.cache_data — 同樣輸入只跑一次",
             Inches(0.95), Inches(5.5), Inches(11), Inches(0.5),
             size=18, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_text(s, "Streamlit 用「函式名稱 + 參數值」當 cache key。第二次同樣參數 → 立刻回傳上次結果。",
             Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.5),
             size=14, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_caching_code(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 8 · CACHING")
    add_title(s, "@st.cache_data — 一個 decorator 解決")

    code = """@st.cache_data
def load_data(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    df["group"] = pd.Categorical(
        df["group"],
        categories=["young", "middle", "older"],
        ordered=True,
    )
    return df

df = load_data("data/cognitive_aging_taiwan.csv")"""
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.8),
             size=16)

    add_text(s, "何時要 cache？",
             Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = [
        "讀檔案、API request",
        "同樣輸入永遠回同樣結果且耗時 > 50 ms",
    ]
    add_bullets(s, items, Inches(0.9), Inches(5.45), Inches(12), Inches(0.85),
                size=16)

    add_text(s, "@st.cache_resource：cache thread / DB connection 等不可序列化物件。",
             Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 9: Hands-on #3 -----
def slide_handson_3(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s, ACCENT_AMBER)
    add_section_label(s, "HANDS-ON #3 · 15 MIN", color=ACCENT_AMBER)
    add_title(s, "完整 Cognitive Aging Dashboard")

    add_text(s, "整合今天所有觀念，建立完整 dashboard：",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    items = [
        "用 @st.cache_data 包住 load_data()",
        "把所有 filters 移到 st.sidebar",
        "加 measure selectbox（5 個認知測驗任選）",
        "主畫面用 st.tabs([Scatter, Distribution, Raw data]) 分三頁",
        "Scatter：散佈圖 + 線性迴歸線",
        "Distribution：young / middle / older 三組 histogram",
        "Raw data：dataframe + st.download_button",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.65), Inches(12), Inches(3.5),
                size=17)

    add_card(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.65),
             fill=BG_OFFWHITE, accent=ACCENT_GREEN)
    add_text(s, "對照解答：app/app.py — 先嘗試自己寫，再開來比對。",
             Inches(0.95), Inches(6.4), Inches(11), Inches(0.4),
             size=15, color=TEXT_DARK, font=FONT_CJK)
    add_footer(s, page, total)
    return s


# ----- Section 10: Deploy -----
def slide_deploy_overview(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 10 · DEPLOY")
    add_title(s, "把你的 app 變成一個公開 URL")

    steps = [
        ("1", "準備 repo",
         "app.py + requirements.txt + data/ 資料夾"),
        ("2", "Push 到 GitHub",
         "git push origin main（沿用 Week 10 的工作流程）"),
        ("3", "登入 share.streamlit.io",
         "用 GitHub 帳號 sign in"),
        ("4", "New app → Deploy",
         "選 repo / branch / app.py 路徑 → 等 1–2 分鐘"),
        ("5", "拿到永久公開 URL",
         "https://your-app.streamlit.app — 修改後 push 自動 redeploy"),
    ]
    y = 2.0
    for num, title, desc in steps:
        add_card(s, Inches(0.7), Inches(y), Inches(11.9), Inches(0.85),
                 accent=ACCENT_TEAL)
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.95), Inches(y + 0.18),
                                    Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT_TEAL
        circle.line.fill.background()
        add_text(s, num, Inches(0.95), Inches(y + 0.20),
                 Inches(0.5), Inches(0.45),
                 size=18, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER)
        add_text(s, title, Inches(1.65), Inches(y + 0.13),
                 Inches(4), Inches(0.4),
                 size=16, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, desc, Inches(1.65), Inches(y + 0.46),
                 Inches(10.5), Inches(0.4),
                 size=13, color=TEXT_MUTED, font=FONT_CJK)
        y += 0.95
    add_footer(s, page, total)
    return s


def slide_deploy_requirements(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "PART 10 · DEPLOY")
    add_title(s, "requirements.txt — 別忘了！")

    add_text(s, "Streamlit Cloud 用這個檔案在 server 上裝套件：",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             size=17, color=TEXT_DARK, font=FONT_CJK)

    code = """# requirements.txt
streamlit>=1.33
pandas>=2.0
numpy>=1.26
matplotlib>=3.8"""
    add_code(s, code, Inches(0.7), Inches(2.65), Inches(7.5), Inches(2.0),
             size=18, lang="bash")

    add_card(s, Inches(8.4), Inches(2.65), Inches(4.2), Inches(2.0),
             accent=ACCENT_AMBER)
    add_text(s, "產生方法",
             Inches(8.65), Inches(2.85), Inches(4), Inches(0.4),
             size=15, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_text(s, "$ pip freeze > requirements.txt",
             Inches(8.65), Inches(3.3), Inches(4), Inches(0.4),
             size=12, color=ACCENT_GREEN, font=FONT_MONO)
    add_text(s, "或手寫只列直接 import 的套件版本範圍。",
             Inches(8.65), Inches(3.85), Inches(3.7), Inches(0.7),
             size=12, color=TEXT_MUTED, font=FONT_CJK)

    # Common pitfalls
    add_text(s, "常見錯誤",
             Inches(0.7), Inches(5.0), Inches(12), Inches(0.4),
             size=16, bold=True, color=ACCENT_RED, font=FONT_CJK)
    items = [
        "忘了加 requirements.txt → Cloud build 失敗",
        "用 conda 套件名（例如 'pillow'）但 Cloud 跑 pip → 套件名不一定一致",
        "data 太大放 GitHub（>100 MB）→ 改用 hosted 資料或 git-lfs",
    ]
    add_bullets(s, items, Inches(0.9), Inches(5.45), Inches(12), Inches(1.4),
                size=14, bullet_color=ACCENT_RED)
    add_footer(s, page, total)
    return s


# ----- Recap / Pitfalls / Homework -----
def slide_pitfalls(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "RECAP")
    add_title(s, "常見錯誤 & 正確做法")

    rows = [
        ("執行模型",        "期待 callback / event handler",        "接受 every-interaction-rerun"),
        ("顯示 DataFrame",  "print(df)、st.write(df.head())",       "st.dataframe(df)"),
        ("顯示 matplotlib", "plt.show()",                            "st.pyplot(fig)"),
        ("載入大資料",       "每次都 pd.read_csv",                   "@st.cache_data"),
        ("Widget 共用變數", "兩個 widget 都 = age",                  "獨立變數名"),
        ("部署",            "忘記 requirements.txt",                "先寫好 requirements"),
    ]
    table_shape = s.shapes.add_table(
        len(rows) + 1, 3,
        Inches(0.7), Inches(2.05),
        Inches(11.9), Inches(4.4),
    )
    t = table_shape.table
    t.columns[0].width = Inches(2.5)
    t.columns[1].width = Inches(4.7)
    t.columns[2].width = Inches(4.7)
    for ci, h in enumerate(["主題", "錯誤", "正確"]):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_PRIMARY
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = TEXT_LIGHT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.04)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT_MONO if ci > 0 and "(" not in val else FONT_CJK
            r.font.size = Pt(13)
            r.font.color.rgb = ACCENT_RED if ci == 1 else (
                ACCENT_GREEN if ci == 2 else TEXT_DARK)
    add_footer(s, page, total)
    return s


def slide_homework(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s, ACCENT_AMBER)
    add_section_label(s, "HOMEWORK", color=ACCENT_AMBER)
    add_title(s, "建立並部署你的 Cognitive Aging Dashboard", size=32)

    add_text(s, "最低要求：",
             Inches(0.7), Inches(2.0), Inches(12), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items = [
        "用 @st.cache_data 載入 cognitive_aging_taiwan.csv",
        "至少 3 種不同類型的 widget",
        "至少 2 種視覺化（scatter + histogram，或 bar + line 等）",
        "用 st.sidebar 或 st.columns 規劃版面",
        "至少 1 個 st.metric（含 delta_color）",
        "加 st.download_button 讓使用者下載篩選後的資料",
        "部署到 Streamlit Cloud 並取得公開 URL",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.5), Inches(12), Inches(3.4),
                size=16)

    add_card(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.95),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "繳交：GitHub repo URL + Streamlit Cloud URL + 100 字 reflection",
             Inches(0.95), Inches(6.05), Inches(11.5), Inches(0.4),
             size=15, bold=True, color=TEXT_DARK, font=FONT_CJK)
    add_text(s, "Rubric：功能 50% · code 品質 20% · deploy 成功 20% · reflection 10%",
             Inches(0.95), Inches(6.45), Inches(11.5), Inches(0.4),
             size=13, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_resources(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_WHITE); add_top_band(s)
    add_section_label(s, "RESOURCES & WHAT'S NEXT")
    add_title(s, "更多資源 + 下週預告")

    # Two columns
    add_text(s, "Streamlit 資源",
             Inches(0.7), Inches(2.0), Inches(6), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items_left = [
        "docs.streamlit.io — 官方文件",
        "cheat-sheet.streamlit.app — 一頁速查",
        "share.streamlit.io — 部署平台",
        "陳 YT, Streamlit 入門 (Medium, 2020)",
        "Mhadhbi, Streamlit Tutorial (DataCamp, 2026)",
    ]
    add_bullets(s, items_left, Inches(0.7), Inches(2.5), Inches(6), Inches(2.5),
                size=15, line_spacing=1.3)

    add_text(s, "What's Next",
             Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    items_right = [
        "Week 12：Open Data API（PubMed / data.gov.tw）",
        "Week 13：Plotly 互動式圖表 + storytelling",
        "Week 14：Claude API — 給 app 加 AI feature",
        "Week 16：Final presentation — 展示你部署的 app",
    ]
    add_bullets(s, items_right, Inches(7.0), Inches(2.5), Inches(6), Inches(2.5),
                size=15, line_spacing=1.3, font=FONT_CJK)

    add_card(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.3),
             accent=ACCENT_TEAL)
    add_text(s, "本週的 dashboard 是你 final project 的起點",
             Inches(0.95), Inches(5.65), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_PRIMARY, font=FONT_CJK)
    add_text(s, "之後幾週會在這個基礎上加 API、加 Plotly、加 AI feature — Week 16 直接展示部署版。",
             Inches(0.95), Inches(6.15), Inches(12), Inches(0.5),
             size=14, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_thanks(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_SECTION)
    add_text(s, "Questions?",
             Inches(0), Inches(2.8), Inches(13.333), Inches(1.4),
             size=64, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text(s, "謝謝大家",
             Inches(0), Inches(4.4), Inches(13.333), Inches(0.6),
             size=24, color=ACCENT_TEAL, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_text(s, "ACL@NCU · Action & Cognition Laboratory",
             Inches(0), Inches(6.5), Inches(13.333), Inches(0.4),
             size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)
    return s


# ============================================================
# Main
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Plan slides — count first to set total page numbers
    plan = [
        ("title",      slide_title),
        ("objectives", slide_objectives),
        ("agenda",     slide_agenda),

        ("d0", lambda p, page, total: slide_divider(
            p, page, total, "0", "WHY STREAMLIT?", "從 Jupyter 到 Web App")),
        ("why",        slide_why_streamlit),
        ("cmp",        slide_compare_frameworks),
        ("hello",      slide_streamlit_hello),

        ("d1", lambda p, page, total: slide_divider(
            p, page, total, "1", "MENTAL MODEL", "Streamlit 的執行模型")),
        ("mm",         slide_mental_model),
        ("install",    slide_install_run),
        ("write",      slide_st_write),
        ("textfam",    slide_text_family),
        ("magic",      slide_magic_commands),

        ("d2", lambda p, page, total: slide_divider(
            p, page, total, "2", "DATASET", "Cognitive Aging Dataset")),
        ("ds",         slide_dataset_intro),
        ("dsp",        slide_dataset_pattern),

        ("h1",         slide_handson_1),

        ("d4", lambda p, page, total: slide_divider(
            p, page, total, "4", "WIDGETS", "互動式控制元件")),
        ("wo",         slide_widgets_overview),
        ("wc",         slide_widgets_code),
        ("media",      slide_media_widgets),
        ("status",     slide_status_messages),

        ("h2",         slide_handson_2),

        ("brk",        slide_break),

        ("d6", lambda p, page, total: slide_divider(
            p, page, total, "6", "CHARTS", "視覺化你的資料")),
        ("cb",         slide_charts_builtin),
        ("cm",         slide_charts_matplotlib),

        ("d7", lambda p, page, total: slide_divider(
            p, page, total, "7", "LAYOUT", "規劃 dashboard 版面")),
        ("lc",         slide_layout_columns),
        ("ls",         slide_layout_sidebar),
        ("lt",         slide_layout_tabs),

        ("d8", lambda p, page, total: slide_divider(
            p, page, total, "8", "CACHING", "@st.cache_data")),
        ("cap",        slide_caching_problem),
        ("cac",        slide_caching_code),
        ("cres",       slide_cache_resource),

        ("h3",         slide_handson_3),
        ("mpage",      slide_multipage),

        ("d10", lambda p, page, total: slide_divider(
            p, page, total, "10", "DEPLOY", "Streamlit Community Cloud")),
        ("dpo",        slide_deploy_overview),
        ("dpr",        slide_deploy_requirements),

        ("d11", lambda p, page, total: slide_divider(
            p, page, total, "11", "RECAP & Q&A", "重點回顧 / 作業")),
        ("pit",        slide_pitfalls),
        ("hw",         slide_homework),
        ("res",        slide_resources),
        ("end",        slide_thanks),
    ]
    total = len(plan)

    for i, (key, builder) in enumerate(plan, start=1):
        if key == "title":
            builder(prs)
        else:
            builder(prs, i, total)

    out = "week-11-slides.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")
    return out


if __name__ == "__main__":
    build()
