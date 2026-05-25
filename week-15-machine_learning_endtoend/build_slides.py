"""Build Week 15 lecture slides — End-to-End Machine Learning.

Anchor: Géron Ch.2 California Housing notebook, used to introduce
the major ML algorithm categories (linear / instance / tree / ensemble
/ kernel / unsupervised).

Visual style: ACL@NCU — white BG, navy + teal accents (matches week-13).

Output: week-15-slides.pptx (16:9, ~50 slides)
Run:    python build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# ACL@NCU palette (shared across weeks)
# ============================================================
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE    = RGBColor(0xF7, 0xF8, 0xFA)
BG_SECTION     = RGBColor(0x14, 0x32, 0x5C)
BG_BREAK       = RGBColor(0xFB, 0xEA, 0xC0)
BG_PRACTICE    = RGBColor(0xEC, 0xF7, 0xF6)
TEXT_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED     = RGBColor(0x5F, 0x6B, 0x83)
ACCENT_PRIMARY = RGBColor(0x14, 0x32, 0x5C)
ACCENT_TEAL    = RGBColor(0x0D, 0x9B, 0x9B)
ACCENT_AMBER   = RGBColor(0xE8, 0xA1, 0x2A)
ACCENT_RED     = RGBColor(0xD3, 0x4F, 0x4F)
ACCENT_GREEN   = RGBColor(0x2E, 0x8B, 0x57)
ACCENT_ORANGE  = RGBColor(0xF9, 0x71, 0x16)
CODE_BG        = RGBColor(0x1E, 0x29, 0x3B)
CODE_COMMENT   = RGBColor(0x8B, 0x9D, 0xB8)
HAIRLINE       = RGBColor(0xE2, 0xE6, 0xEC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"

FOOTER_STR = "Week 15 — End-to-End Machine Learning  ·  ACL@NCU"


# ============================================================
# Primitives (copied from week-13 build_slides.py for consistency)
# ============================================================
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_top_band(slide, color=ACCENT_TEAL, height=Inches(0.14)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=20, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_section_label(slide, text, color=ACCENT_TEAL):
    add_text(slide, text,
             Inches(0.7), Inches(0.55), Inches(11), Inches(0.35),
             size=14, bold=True, color=color, font=FONT_SANS)


def add_title(slide, title_zh, *, size=32, top=0.95, color=TEXT_DARK):
    add_text(slide, title_zh,
             Inches(0.7), Inches(top), Inches(12), Inches(0.85),
             size=size, bold=True, color=color, font=FONT_CJK)


def add_bullets(slide, items, left, top, width, height, *,
                size=18, color=TEXT_DARK, line_spacing=1.25,
                bullet_color=ACCENT_TEAL, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    use_font = font or FONT_CJK
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = "▸  "
        b.font.name = FONT_SANS
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = use_font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i > 0:
            p.space_before = Pt(size * 0.45)
    return tb


def _color_for_code_line(line, lang):
    s = line.lstrip()
    if not s:
        return None
    if lang in ("bash", "powershell", "yaml", "python", "markdown") and s.startswith("#"):
        return CODE_COMMENT
    if s.startswith(">"):
        return ACCENT_AMBER
    if s.startswith("$") or s.startswith("PS "):
        return ACCENT_GREEN
    return None


def add_code(slide, code, left, top, width, height, *,
             size=14, lang="python", line_spacing=1.15, padding=0.18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(padding); tf.margin_bottom = Inches(padding)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        c = _color_for_code_line(line, lang) or TEXT_LIGHT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(size)
        r.font.color.rgb = c
    return box


def add_card(slide, left, top, width, height, *, fill=BG_OFFWHITE,
             accent=None, accent_w=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = HAIRLINE
    card.line.width = Pt(0.5)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(accent_w), height)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return card


def add_callout(slide, text, left, top, width, height, *,
                size=22, color=ACCENT_PRIMARY, fill=BG_OFFWHITE,
                accent=ACCENT_TEAL, font=FONT_CJK, align=PP_ALIGN.LEFT):
    add_card(slide, left, top, width, height, fill=fill, accent=accent,
             accent_w=0.10)
    add_text(slide, text,
             left + Inches(0.4), top + Inches(0.2),
             width - Inches(0.6), height - Inches(0.4),
             size=size, bold=True, color=color, font=font,
             align=align, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page, total, *, theme="light"):
    color = TEXT_MUTED if theme == "light" else ACCENT_TEAL
    add_text(slide, FOOTER_STR,
             Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
             size=11, color=color)
    add_text(slide, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=color, align=PP_ALIGN.RIGHT)


def add_table(slide, headers, rows, left, top, width, height, *,
              header_fill=ACCENT_PRIMARY, header_color=TEXT_LIGHT,
              col_widths=None, header_size=13, body_size=12,
              first_col_bold=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_color
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT_CJK
            r.font.size = Pt(body_size)
            r.font.color.rgb = TEXT_DARK
            if first_col_bold and ci == 0:
                r.font.bold = True
    return table_shape


def add_slide_base(prs, *, bg=BG_WHITE, band=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, bg)
    if band:
        add_top_band(s)
    return s


def add_practice_badge(slide, left=0.7, top=0.55):
    add_text(slide, "🔬 HANDS-ON",
             Inches(left), Inches(top), Inches(2.5), Inches(0.35),
             size=14, bold=True, color=ACCENT_TEAL)


# ============================================================
# Slide builders
# ============================================================
def slide_title(prs):
    s = add_slide_base(prs, band=False)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(2.5),
                               Inches(0.18), Inches(2.6))
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT_TEAL
    block.line.fill.background()
    add_text(s, "WEEK 15",
             Inches(1.1), Inches(2.4), Inches(11), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, "End-to-End Machine Learning",
             Inches(1.1), Inches(2.95), Inches(12), Inches(0.95),
             size=40, bold=True, color=TEXT_DARK, font=FONT_SANS)
    add_text(s, "從資料到模型的完整 pipeline",
             Inches(1.1), Inches(3.95), Inches(12), Inches(0.85),
             size=32, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, "Anchor: Géron 2023 · Hands-on ML · Chapter 2",
             Inches(1.1), Inches(4.95), Inches(12), Inches(0.5),
             size=18, color=TEXT_MUTED, font=FONT_SANS, italic=True)
    add_text(s, "NS5116 · Programming & AI Applications in Behavioral Science",
             Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
             size=14, color=TEXT_MUTED)
    add_text(s, "ACL@NCU  ·  Spring 2026  ·  2026-06-04  ·  張智宏",
             Inches(1.1), Inches(6.4), Inches(11), Inches(0.4),
             size=13, color=TEXT_MUTED)
    return s


def slide_section(prs, label, title_zh, title_en=None):
    """Full-bleed dark section divider."""
    s = add_slide_base(prs, bg=BG_SECTION, band=False)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))
    add_text(s, label,
             Inches(0.9), Inches(2.5), Inches(11), Inches(0.6),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, title_zh,
             Inches(0.9), Inches(3.15), Inches(12), Inches(1.2),
             size=42, bold=True, color=TEXT_LIGHT, font=FONT_CJK)
    if title_en:
        add_text(s, title_en,
                 Inches(0.9), Inches(4.5), Inches(12), Inches(0.6),
                 size=20, color=TEXT_LIGHT, font=FONT_SANS, italic=True)
    return s


def slide_objectives(prs):
    s = add_slide_base(prs)
    add_section_label(s, "Learning Objectives")
    add_title(s, "今天結束後你會：")
    items = [
        "Frame — 判斷一個問題屬於 supervised / unsupervised、regression / classification",
        "Split — 用 StratifiedShuffleSplit 切出 reproducible 的 test set",
        "Pipeline — 組合 imputation + encoding + scaling 而不發生 data leakage",
        "Compare — 用 5-fold CV 比較 ≥5 種 regression algorithm 的優劣",
        "Tune — 用 GridSearch / RandomizedSearch 做 hyperparameter 調整",
        "Transfer — 把同一套 pipeline 套到 Stroop / Flanker RT 預測問題",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5),
                size=20)
    return s


def slide_agenda(prs):
    s = add_slide_base(prs)
    add_section_label(s, "Agenda — 3 hours")
    add_title(s, "今天的路線圖")
    rows = [
        ("§1", "Framing ML problems", "20 min", "supervised vs. unsupervised, regression vs. classification"),
        ("§2", "Anchor case: California Housing", "25 min", "load, split, EDA"),
        ("§3", "Preprocessing Pipeline", "30 min", "imputation, encoding, scaling, ColumnTransformer"),
        ("§4", "Algorithm Zoo (regression)", "40 min", "linear, k-NN, tree, ensemble, kernel"),
        ("§5", "Model Evaluation & Tuning", "25 min", "k-fold CV, Grid/Randomized search"),
        ("§6", "Unsupervised Tasters", "15 min", "k-means as features, IsolationForest"),
        ("§7", "Cogneuro Transfer", "25 min", "Stroop RT prediction with the same pipeline"),
    ]
    add_table(s,
              ["§", "Topic", "Time", "Subtopics"],
              rows,
              Inches(0.7), Inches(1.95), Inches(12.0), Inches(4.5),
              col_widths=[Inches(0.7), Inches(3.8), Inches(1.3), Inches(6.2)],
              header_size=14, body_size=13, first_col_bold=True)
    return s


def slide_why(prs):
    s = add_slide_base(prs)
    add_section_label(s, "為什麼今天要學 ML")
    add_title(s, "你已經會收資料、整理資料、畫圖 — 接著要從資料學出一個可預測的 model")
    add_callout(
        s,
        "「哪些 trial-level 變項最能預測 RT？\n  能不能用一個 model 在新受試者身上預測 RT？」",
        Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
        size=22, accent=ACCENT_TEAL)
    add_text(s, "這正是一個 supervised regression 問題。",
             Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.5),
             size=18, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "今天的策略：",
             Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.5),
             size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
    items = [
        "Anchor case：跑完 California Housing 案例（資料乾淨、文獻多）",
        "用每個 pipeline 階段引入一類 ML algorithm",
        "最後把同一套 pipeline 套到 cogneuro 問題上",
    ]
    add_bullets(s, items, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.7),
                size=18)
    return s


# ============================================================
# §1 Framing
# ============================================================
def slide_section1(prs):
    return slide_section(prs, "Section 1", "ML 問題的框架化", "Framing the ML problem")


def slide_four_questions(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§1 · Framing")
    add_title(s, "寫 sklearn 之前先問四個問題")
    rows = [
        ("有沒有 label？", "supervised / unsupervised / semi-supervised",
         "supervised（label = 觀測到的 RT）"),
        ("Label 的型別？", "regression / classification",
         "regression（RT 是連續值）"),
        ("資料一次給齊嗎？", "batch / online",
         "batch（一次拿到所有受試者）"),
        ("模型怎麼學？", "instance-based / model-based",
         "model-based（學出一組 coefficient）"),
    ]
    add_table(s, ["問題", "選項", "Stroop RT 的答案"], rows,
              Inches(0.7), Inches(2.0), Inches(12), Inches(3.2),
              col_widths=[Inches(2.5), Inches(4.7), Inches(4.8)],
              header_size=14, body_size=14, first_col_bold=True)
    add_callout(s,
                "選錯類型會讓所有 metric 變得無意義。\n例如把 RT (ms) 當 classification target，silently 失去 ordinal 結構。",
                Inches(0.7), Inches(5.5), Inches(12), Inches(1.1),
                size=16, accent=ACCENT_AMBER)
    return s


def slide_taxonomy_overview(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§1 · Framing")
    add_title(s, "本週涵蓋的 algorithm 類別地圖")

    # Supervised box
    add_card(s, Inches(0.7), Inches(2.0), Inches(6.0), Inches(4.5),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "Supervised", Inches(1.0), Inches(2.15),
             Inches(5.5), Inches(0.5),
             size=22, bold=True, color=ACCENT_PRIMARY)
    items_s = [
        "Linear  →  LinearRegression, Ridge",
        "Instance  →  KNeighborsRegressor",
        "Tree  →  DecisionTreeRegressor",
        "Ensemble  →  RandomForest, GradientBoosting",
        "Kernel  →  SVR (rbf)",
    ]
    add_bullets(s, items_s, Inches(1.0), Inches(2.75),
                Inches(5.5), Inches(3.5), size=16, font=FONT_SANS)

    # Unsupervised box
    add_card(s, Inches(6.95), Inches(2.0), Inches(5.7), Inches(4.5),
             fill=BG_OFFWHITE, accent=ACCENT_ORANGE)
    add_text(s, "Unsupervised", Inches(7.25), Inches(2.15),
             Inches(5.2), Inches(0.5),
             size=22, bold=True, color=ACCENT_PRIMARY)
    items_u = [
        "Clustering  →  KMeans",
        "Anomaly detection  →  IsolationForest",
        "(本週作為 feature engineering 與 outlier filter)",
    ]
    add_bullets(s, items_u, Inches(7.25), Inches(2.75),
                Inches(5.2), Inches(3.5), size=16, font=FONT_SANS)

    add_text(s, "Deep learning 屬於另一個範疇（representation learning），本週不涵蓋。",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             size=13, color=TEXT_MUTED, font=FONT_CJK, italic=True)
    return s


def slide_practice1(prs):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_practice_badge(s)
    add_title(s, "Hands-on 1 — 把這三個情境框架化")
    rows = [
        ("A", "從 EEG 30 秒片段判斷受試者是清醒還是睡著",
         "supervised · binary classification · batch · model-based"),
        ("B", "把 1000 篇 fMRI 論文依「研究主題」自動分群",
         "unsupervised · clustering · batch · model-based (k-means)"),
        ("C", "判斷某張 fMRI volume 是否為 motion artifact",
         "通常 supervised 或 unsupervised anomaly · 可 batch 或 online"),
    ]
    add_table(s, ["情境", "問題描述", "建議答案"], rows,
              Inches(0.7), Inches(2.0), Inches(12), Inches(3.5),
              col_widths=[Inches(0.8), Inches(5.5), Inches(5.7)],
              header_size=14, body_size=13, first_col_bold=True)
    add_text(s, "5 分鐘小組討論 → 一組分享一個情境的答案。",
             Inches(0.7), Inches(5.8), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    return s


# ============================================================
# §2 Anchor case
# ============================================================
def slide_section2(prs):
    return slide_section(prs, "Section 2",
                         "Anchor case — California Housing",
                         "Load · Split · EDA")


def slide_data_overview(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§2 · Anchor case")
    add_title(s, "資料：1990 加州人口普查 block-level")
    items = [
        "20,640 個 block，10 個欄位",
        "Target = median_house_value（USD）",
        "Features：經緯度、housing_median_age、總房間數、總臥房數、人口、家戶數、median_income、ocean_proximity",
        "注意：total_bedrooms 有 207 筆缺值，ocean_proximity 是文字",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(2.5),
                size=18)
    code = ('import pandas as pd\n'
            'from pathlib import Path\n'
            'import tarfile, urllib.request\n\n'
            'def load_housing_data():\n'
            '    p = Path("datasets/housing.tgz")\n'
            '    if not p.is_file():\n'
            '        Path("datasets").mkdir(exist_ok=True)\n'
            '        url = "https://github.com/ageron/data/raw/main/housing.tgz"\n'
            '        urllib.request.urlretrieve(url, p)\n'
            '        with tarfile.open(p) as t: t.extractall("datasets")\n'
            '    return pd.read_csv("datasets/housing/housing.csv")\n\n'
            'housing = load_housing_data()')
    add_code(s, code, Inches(0.7), Inches(4.7), Inches(12), Inches(2.1), size=13)
    return s


def slide_golden_rule(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§2 · Anchor case")
    add_title(s, "黃金法則：拿到資料的第一件事是切 test set")
    add_callout(s,
                "切完之後，把 test set 鎖起來，直到最後評估前都不准看。",
                Inches(0.7), Inches(2.1), Inches(12), Inches(1.1),
                size=24, accent=ACCENT_RED)
    items = [
        "如果先做 EDA 再切：你會 silently 用 test 的資訊影響 design choice",
        "這叫 data snooping bias — 等於把答案抄進考卷",
        "Reproducibility：random_state=42 確保每次得到相同 split",
    ]
    add_bullets(s, items, Inches(0.7), Inches(3.5), Inches(12), Inches(2.0),
                size=18)
    code = ('from sklearn.model_selection import train_test_split\n\n'
            'train_set, test_set = train_test_split(\n'
            '    housing, test_size=0.2, random_state=42)')
    add_code(s, code, Inches(0.7), Inches(5.7), Inches(12), Inches(1.0), size=14)
    return s


def slide_stratified(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§2 · Anchor case")
    add_title(s, "Stratified split — 當分布不均時")
    items = [
        "median_income 很重要，但分布右尾長",
        "純 random split 可能讓 test 的 high-income block 比例 ≠ 母體",
        "解法：把 income 分 5 個 bin，按 bin 比例分層抽樣",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(1.7),
                size=18)
    code = ('import numpy as np\n'
            'from sklearn.model_selection import StratifiedShuffleSplit\n\n'
            'housing["income_cat"] = pd.cut(\n'
            '    housing["median_income"],\n'
            '    bins=[0., 1.5, 3.0, 4.5, 6., np.inf],\n'
            '    labels=[1, 2, 3, 4, 5])\n\n'
            'split = StratifiedShuffleSplit(n_splits=1, test_size=0.2, random_state=42)\n'
            'for tr, te in split.split(housing, housing["income_cat"]):\n'
            '    strat_train = housing.iloc[tr]\n'
            '    strat_test  = housing.iloc[te]')
    add_code(s, code, Inches(0.7), Inches(3.8), Inches(12), Inches(2.4), size=13)
    add_callout(s,
                "Cogneuro 類比：若年輕人遠多於老人，用 age group stratify\n才能保證 test 反映你想 generalize 的母體。",
                Inches(0.7), Inches(6.3), Inches(12), Inches(0.8),
                size=14, accent=ACCENT_TEAL)
    return s


def slide_practice2(prs):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_practice_badge(s)
    add_title(s, "Hands-on 2 — Stratified split on age")
    add_text(s, "任務：給定一個模擬 RT dataset（n=200, age 20-80），先用 pd.cut 切 4 個 age bin，"
             "再用 StratifiedShuffleSplit 切 20% test。驗證 train/test 中各 bin 比例差異 < 1%。",
             Inches(0.7), Inches(1.95), Inches(12), Inches(1.4),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    code = ('import numpy as np, pandas as pd\n'
            'np.random.seed(42)\n'
            'df = pd.DataFrame({\n'
            '    "age": np.random.uniform(20, 80, 200),\n'
            '    "rt":  np.random.normal(500, 80, 200),\n'
            '})\n'
            '# 你的程式碼從這裡開始')
    add_code(s, code, Inches(0.7), Inches(3.5), Inches(12), Inches(1.9), size=14)
    add_text(s, "8 分鐘獨立完成 → 同桌互看程式碼。",
             Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    return s


# ============================================================
# §3 Pipeline
# ============================================================
def slide_section3(prs):
    return slide_section(prs, "Section 3",
                         "Data Preparation Pipeline",
                         "Impute · Encode · Scale · Compose")


def slide_pipeline_why(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "核心觀念：fit on train, transform on test")
    add_callout(s,
                "所有 preprocessing 步驟都要先在 train 上 fit，再 transform 到 test。\n"
                "違反這條 = data leakage = test RMSE 過度樂觀。",
                Inches(0.7), Inches(2.0), Inches(12), Inches(1.5),
                size=20, accent=ACCENT_RED)
    items = [
        "sklearn.pipeline.Pipeline 是強制你遵守這條規則的工具",
        "ColumnTransformer 讓不同欄位走不同的 preprocessing 路徑",
        "整個 pipeline 像一個 model — 對它 fit 一次、predict 一次",
    ]
    add_bullets(s, items, Inches(0.7), Inches(3.8), Inches(12), Inches(2.2),
                size=18)
    return s


def slide_pipeline_flow(prs):
    """Visual diagram of the preprocessing pipeline data flow."""
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "Preprocessing pipeline — 資料流圖", size=28)

    # --------------------------------------------------------------------
    # Layout reference (inches):
    #   raw block ~ y=1.6
    #   ColumnTransformer band ~ y=2.5
    #   numeric track   x=2.0-5.5, y=3.4-5.0
    #   categorical track x=7.8-11.3, y=3.4-5.0
    #   merge / X_prepared ~ y=5.3
    #   estimator ~ y=6.1
    # --------------------------------------------------------------------

    def box(left, top, width, height, text, *, fill, text_color=TEXT_LIGHT,
            font=FONT_SANS, size=12, bold=False, line=None):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line; sh.line.width = Pt(0.5)
        sh.adjustments[0] = 0.18
        tf = sh.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = text_color
        return sh

    def arrow(x1, y1, x2, y2, color=TEXT_DARK, weight=1.8):
        c = s.shapes.add_connector(2, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
        c.line.color.rgb = color
        c.line.width = Pt(weight)
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            ln = c.line._get_or_add_ln()
            tail = etree.SubElement(ln, qn("a:tailEnd"))
            tail.set("type", "triangle")
            tail.set("w", "med"); tail.set("len", "med")
        except Exception:
            pass
        return c

    # --- (1) Raw DataFrame ----------------------------------------------
    box(3.5, 1.55, 6.3, 0.75,
        "Raw DataFrame  (X with mixed types)",
        fill=ACCENT_PRIMARY, size=15, bold=True)
    add_text(s,
             "longitude  ·  latitude  ·  median_income  ·  total_bedrooms (NA)  ·  ocean_proximity (str)",
             Inches(2.6), Inches(2.32), Inches(8.1), Inches(0.32),
             size=10, color=TEXT_MUTED, font=FONT_MONO, align=PP_ALIGN.CENTER)

    # arrow down into ColumnTransformer
    arrow(6.65, 2.34, 6.65, 2.78)

    # --- (2) ColumnTransformer "splitter" --------------------------------
    box(2.5, 2.78, 8.3, 0.50,
        "ColumnTransformer  —  不同欄位走不同路徑",
        fill=ACCENT_TEAL, size=14, bold=True)

    # arrows to two tracks
    arrow(4.5, 3.28, 3.75, 3.55, color=ACCENT_TEAL)
    arrow(8.75, 3.28, 9.50, 3.55, color=ACCENT_ORANGE)

    # --- (3a) Numerical track (left) -------------------------------------
    add_card(s, Inches(1.45), Inches(3.55), Inches(4.6), Inches(2.4),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL, accent_w=0.06)
    add_text(s, "Numerical columns",
             Inches(1.65), Inches(3.62), Inches(4.3), Inches(0.32),
             size=12, bold=True, color=ACCENT_PRIMARY, font=FONT_SANS)
    add_text(s, "longitude, latitude, total_rooms,\nmedian_income, ...",
             Inches(1.65), Inches(3.92), Inches(4.3), Inches(0.55),
             size=10, color=TEXT_MUTED, font=FONT_MONO,
             align=PP_ALIGN.LEFT)

    box(1.75, 4.55, 4.0, 0.55,
        "SimpleImputer(strategy='median')",
        fill=ACCENT_TEAL, size=12, bold=True)
    add_text(s, "→ learns the median per column",
             Inches(1.75), Inches(5.12), Inches(4.0), Inches(0.28),
             size=10, color=TEXT_MUTED, font=FONT_CJK, italic=True,
             align=PP_ALIGN.CENTER)
    arrow(3.75, 5.10, 3.75, 5.40, color=ACCENT_TEAL)

    box(1.75, 5.42, 4.0, 0.50,
        "StandardScaler()",
        fill=ACCENT_TEAL, size=12, bold=True)
    add_text(s, "→ learns mean & std per column",
             Inches(1.75), Inches(5.93), Inches(4.0), Inches(0.28),
             size=10, color=TEXT_MUTED, font=FONT_CJK, italic=True,
             align=PP_ALIGN.CENTER)

    # --- (3b) Categorical track (right) ----------------------------------
    add_card(s, Inches(7.30), Inches(3.55), Inches(4.6), Inches(2.4),
             fill=BG_OFFWHITE, accent=ACCENT_ORANGE, accent_w=0.06)
    add_text(s, "Categorical columns",
             Inches(7.50), Inches(3.62), Inches(4.3), Inches(0.32),
             size=12, bold=True, color=ACCENT_PRIMARY, font=FONT_SANS)
    add_text(s, "ocean_proximity",
             Inches(7.50), Inches(3.92), Inches(4.3), Inches(0.32),
             size=10, color=TEXT_MUTED, font=FONT_MONO)

    box(7.55, 4.30, 4.1, 0.95,
        "OneHotEncoder(\n  handle_unknown='ignore')",
        fill=ACCENT_ORANGE, size=12, bold=True, font=FONT_MONO)
    add_text(s, "→ learns the set of categories",
             Inches(7.55), Inches(5.28), Inches(4.1), Inches(0.28),
             size=10, color=TEXT_MUTED, font=FONT_CJK, italic=True,
             align=PP_ALIGN.CENTER)

    # --- (4) Merge node + X_prepared -------------------------------------
    arrow(3.75, 5.95, 5.95, 6.30, color=ACCENT_TEAL)
    arrow(9.60, 5.55, 7.40, 6.30, color=ACCENT_ORANGE)

    box(5.45, 6.30, 2.45, 0.55,
        "concat → X_prepared",
        fill=ACCENT_PRIMARY, size=13, bold=True)

    # --- (5) Estimator ---------------------------------------------------
    arrow(6.67, 6.86, 6.67, 7.05, color=TEXT_DARK)
    add_text(s, "estimator.fit() / .predict()",
             Inches(5.0), Inches(7.05), Inches(3.4), Inches(0.32),
             size=12, bold=True, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.CENTER)

    # --- Side callouts: fit on train (left) / transform on test (right) --
    fit_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.10), Inches(3.55),
                                 Inches(1.25), Inches(2.10))
    fit_box.fill.solid(); fit_box.fill.fore_color.rgb = ACCENT_GREEN
    fit_box.line.fill.background()
    add_text(s, "fit_transform\n(X_train)",
             Inches(0.10), Inches(3.55), Inches(1.25), Inches(2.10),
             size=12, bold=True, color=TEXT_LIGHT, font=FONT_SANS,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "← 學 statistics",
             Inches(0.10), Inches(5.70), Inches(1.25), Inches(0.30),
             size=10, color=ACCENT_GREEN, font=FONT_CJK,
             align=PP_ALIGN.CENTER, italic=True)

    tx_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(12.00), Inches(3.55),
                                Inches(1.25), Inches(2.10))
    tx_box.fill.solid(); tx_box.fill.fore_color.rgb = ACCENT_RED
    tx_box.line.fill.background()
    add_text(s, "transform\n(X_test)",
             Inches(12.00), Inches(3.55), Inches(1.25), Inches(2.10),
             size=12, bold=True, color=TEXT_LIGHT, font=FONT_SANS,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "↓ 重用 statistics",
             Inches(12.00), Inches(5.70), Inches(1.25), Inches(0.30),
             size=10, color=ACCENT_RED, font=FONT_CJK,
             align=PP_ALIGN.CENTER, italic=True)

    return s


def slide_imputation(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "Step A — Imputation：處理 missing values")
    rows = [
        ("Drop rows", "df.dropna(subset=['total_bedrooms'])", "資料量小可接受；否則浪費"),
        ("Drop column", "df.drop('total_bedrooms', axis=1)", "整欄缺很多時"),
        ("Impute (推薦)", "SimpleImputer(strategy='median')", "不丟資料、不丟欄位"),
    ]
    add_table(s, ["策略", "Code", "何時用"], rows,
              Inches(0.7), Inches(2.0), Inches(12), Inches(2.5),
              col_widths=[Inches(2.2), Inches(5.0), Inches(4.8)],
              header_size=14, body_size=13, first_col_bold=True)
    code = ('from sklearn.impute import SimpleImputer\n'
            'imputer = SimpleImputer(strategy="median")\n'
            'X_num_imputed = imputer.fit_transform(housing_num)')
    add_code(s, code, Inches(0.7), Inches(4.9), Inches(12), Inches(1.4), size=14)
    add_text(s, "其他策略：mean, most_frequent, constant；進階可用 KNNImputer。",
             Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
             size=14, color=TEXT_MUTED, font=FONT_CJK)
    return s


def slide_encoding(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "Step B — Encoding 文字類別")
    rows = [
        ("OrdinalEncoder", "類別有自然順序 (low / med / high)",
         "❌ 對無序類別會強加錯誤順序"),
        ("OneHotEncoder", "類別無序（最安全的預設）",
         "⚠ 類別很多時產生 sparse 高維 feature"),
    ]
    add_table(s, ["方法", "適用情境", "缺點"], rows,
              Inches(0.7), Inches(2.0), Inches(12), Inches(2.0),
              col_widths=[Inches(2.5), Inches(4.5), Inches(5.0)],
              header_size=14, body_size=14, first_col_bold=True)
    code = ('from sklearn.preprocessing import OneHotEncoder\n\n'
            'cat_encoder = OneHotEncoder(handle_unknown="ignore",\n'
            '                            sparse_output=False)\n'
            'X_cat_1hot = cat_encoder.fit_transform(housing[["ocean_proximity"]])\n'
            'print(cat_encoder.categories_)\n'
            '# [array([\'<1H OCEAN\', \'INLAND\', \'ISLAND\', \'NEAR BAY\', \'NEAR OCEAN\'])]')
    add_code(s, code, Inches(0.7), Inches(4.4), Inches(12), Inches(2.2), size=13)
    return s


def slide_scaling(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "Step C — Feature Scaling")
    add_callout(s,
                "為什麼要 scaling：距離型 (k-NN, SVM, k-means) 與梯度型 (linear+GD, NN)\n"
                "對 feature 的 scale 極度敏感。Tree-based 演算法則不需要。",
                Inches(0.7), Inches(2.0), Inches(12), Inches(1.4),
                size=18, accent=ACCENT_AMBER)
    rows = [
        ("StandardScaler", "z-score：(x − μ) / σ", "mean=0, std=1；outlier 不會壓縮資料"),
        ("MinMaxScaler", "(x − min) / (max − min)", "壓到 [0, 1]；對 outlier 敏感"),
        ("RobustScaler", "用 median 與 IQR", "outlier 多時的替代方案"),
    ]
    add_table(s, ["Scaler", "公式", "特性"], rows,
              Inches(0.7), Inches(3.7), Inches(12), Inches(2.4),
              col_widths=[Inches(2.5), Inches(3.8), Inches(5.7)],
              header_size=14, body_size=14, first_col_bold=True)
    add_text(s, "對 RandomForest / GradientBoosting：scaling 沒影響，但浪費計算。",
             Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             size=14, color=TEXT_MUTED, font=FONT_CJK)
    return s


def slide_columntransformer(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§3 · Pipeline")
    add_title(s, "Step D — 把所有步驟串成 ColumnTransformer")
    code = ('from sklearn.pipeline import Pipeline\n'
            'from sklearn.compose import ColumnTransformer\n'
            'from sklearn.preprocessing import StandardScaler, OneHotEncoder\n'
            'from sklearn.impute import SimpleImputer\n\n'
            'num_attribs = ["longitude", "latitude", "housing_median_age",\n'
            '               "total_rooms", "total_bedrooms", "population",\n'
            '               "households", "median_income"]\n'
            'cat_attribs = ["ocean_proximity"]\n\n'
            'num_pipeline = Pipeline([\n'
            '    ("imputer", SimpleImputer(strategy="median")),\n'
            '    ("scaler",  StandardScaler()),\n'
            '])\n\n'
            'full_pipeline = ColumnTransformer([\n'
            '    ("num", num_pipeline, num_attribs),\n'
            '    ("cat", OneHotEncoder(handle_unknown="ignore"), cat_attribs),\n'
            '])')
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(12), Inches(4.5), size=13)
    add_callout(s,
                "對 train 用 fit_transform()，對 test 只能用 transform() — pipeline 自動處理。",
                Inches(0.7), Inches(6.55), Inches(12), Inches(0.55),
                size=14, accent=ACCENT_GREEN, font=FONT_CJK)
    return s


def slide_practice3(prs):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_practice_badge(s)
    add_title(s, "Hands-on 3 — 把 RT data 套進 pipeline")
    add_text(s,
             "情境：trial-level dataframe 有 congruency (字串)、isi (連續, 有 missing)、"
             "block_num (整數)、rt (target)。寫一個 ColumnTransformer 處理三個 feature。",
             Inches(0.7), Inches(1.95), Inches(12), Inches(1.4),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    code = ('import pandas as pd, numpy as np\n'
            'np.random.seed(42)\n'
            'df = pd.DataFrame({\n'
            '    "congruency": np.random.choice(["congruent", "incongruent"], 100),\n'
            '    "isi":        np.random.choice([500, 1000, 1500, np.nan], 100),\n'
            '    "block_num":  np.random.choice([1, 2, 3], 100),\n'
            '    "rt":         np.random.normal(500, 80, 100),\n'
            '})\n'
            '# 你的程式碼從這裡開始')
    add_code(s, code, Inches(0.7), Inches(3.5), Inches(12), Inches(2.5), size=13)
    add_text(s, "10 分鐘獨立完成。",
             Inches(0.7), Inches(6.2), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    return s


# ============================================================
# §4 Algorithm zoo
# ============================================================
def slide_section4(prs):
    return slide_section(prs, "Section 4",
                         "Algorithm Zoo — Regression",
                         "Linear · Instance · Tree · Ensemble · Kernel")


def slide_taxonomy_table(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo")
    add_title(s, "Algorithm taxonomy 速覽")
    rows = [
        ("Linear (parametric)", "LinearRegression, Ridge, Lasso",
         "target 是 feature 的線性組合", "baseline；接近線性時"),
        ("Instance-based", "KNeighborsRegressor",
         "鄰近的 sample 有相似 target", "區域結構強、資料密集"),
        ("Tree-based", "DecisionTreeRegressor",
         "feature 空間的 axis-aligned 切割", "interaction 強、可解釋"),
        ("Ensemble", "RandomForest, GradientBoosting",
         "多個 weak learner 投票/平均", "tabular data 最強 baseline"),
        ("Kernel methods", "SVR(kernel='rbf')",
         "high-dim 找 margin-maximizing hyperplane", "中等資料量、非線性"),
    ]
    add_table(s, ["類別", "代表 estimator", "歸納偏置 (inductive bias)", "何時用"],
              rows,
              Inches(0.5), Inches(2.0), Inches(12.4), Inches(4.5),
              col_widths=[Inches(2.5), Inches(3.0), Inches(4.4), Inches(2.5)],
              header_size=13, body_size=12, first_col_bold=True)
    return s


def _algo_card(slide, left, top, width, height,
               family, estimator, intuition, gotcha):
    add_card(slide, left, top, width, height, fill=BG_OFFWHITE,
             accent=ACCENT_TEAL)
    add_text(slide, family, left + Inches(0.3), top + Inches(0.15),
             width - Inches(0.4), Inches(0.5),
             size=22, bold=True, color=ACCENT_PRIMARY, font=FONT_SANS)
    add_text(slide, estimator, left + Inches(0.3), top + Inches(0.75),
             width - Inches(0.4), Inches(0.4),
             size=14, color=ACCENT_TEAL, font=FONT_MONO)
    add_text(slide, intuition, left + Inches(0.3), top + Inches(1.25),
             width - Inches(0.4), height - Inches(2.2),
             size=15, color=TEXT_DARK, font=FONT_CJK)
    add_text(slide, gotcha,
             left + Inches(0.3), top + height - Inches(0.8),
             width - Inches(0.4), Inches(0.7),
             size=13, color=ACCENT_RED, font=FONT_CJK, italic=True)


def slide_algo_linear(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo · 1/5")
    add_title(s, "Family 1：Linear models")
    _algo_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.6),
               "Linear / Ridge", "LinearRegression, Ridge",
               "假設 y = w·x + b。Ridge 加上 L2 penalty 防止 overfit。\n\n"
               "適合：baseline、變項數遠少於樣本數、需要可解釋 coefficient。",
               "⚠ 對 feature 非線性 / 互動完全無感")
    s.shapes.add_picture(
        "diagrams/algo_linear.png",
        Inches(0.7), Inches(5.75), width=Inches(5.8), height=Inches(1.30))
    code = ('from sklearn.linear_model import LinearRegression, Ridge\n\n'
            'lin   = LinearRegression()\n'
            'ridge = Ridge(alpha=1.0)\n\n'
            'lin.fit(X_train, y_train)\n'
            'print(lin.coef_)        # 看每個 feature 的權重\n'
            'print(lin.intercept_)')
    add_code(s, code, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5), size=13)
    return s


def slide_algo_knn(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo · 2/5")
    add_title(s, "Family 2：Instance-based — k-Nearest Neighbors")
    _algo_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.6),
               "k-NN", "KNeighborsRegressor",
               "對新 sample，找最近 k 個 train sample，取平均當預測值。\n\n"
               "適合：區域結構強、資料量大、不想假設函數形式。",
               "⚠ 必須先 scaling；維度詛咒 (curse of dimensionality)")
    s.shapes.add_picture(
        "diagrams/algo_knn.png",
        Inches(0.7), Inches(5.75), width=Inches(5.8), height=Inches(1.30))
    code = ('from sklearn.neighbors import KNeighborsRegressor\n\n'
            'knn = KNeighborsRegressor(n_neighbors=5)\n'
            '# 不 scaling 的話 latitude (±90) 會壓過\n'
            '# income (~3) 的距離貢獻\n\n'
            'pipe = make_pipeline(\n'
            '    StandardScaler(),\n'
            '    KNeighborsRegressor(n_neighbors=5))')
    add_code(s, code, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5), size=13)
    return s


def slide_algo_tree(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo · 3/5")
    add_title(s, "Family 3：Tree-based — Decision Tree")
    _algo_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.6),
               "Decision Tree", "DecisionTreeRegressor",
               "遞迴把 feature 空間切成 axis-aligned 區塊，每塊預測該區的 mean。\n\n"
               "適合：feature interaction 強、需要可視化 decision path。",
               "⚠ Single tree 容易 overfit；極度 high variance")
    s.shapes.add_picture(
        "diagrams/algo_tree.png",
        Inches(0.7), Inches(5.75), width=Inches(5.8), height=Inches(1.30))
    code = ('from sklearn.tree import DecisionTreeRegressor\n\n'
            'tree = DecisionTreeRegressor(\n'
            '    max_depth=None,        # 不限深度 → overfit\n'
            '    random_state=42)\n\n'
            '# 改善：限制深度\n'
            'tree2 = DecisionTreeRegressor(max_depth=8,\n'
            '                              min_samples_leaf=20)')
    add_code(s, code, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5), size=13)
    return s


def slide_algo_ensemble(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo · 4/5")
    add_title(s, "Family 4：Ensemble — Random Forest & Gradient Boosting")
    _algo_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.6),
               "Random Forest", "RandomForestRegressor",
               "Bagging：平行訓練多棵 deep tree（隨機 sample + 隨機 feature），\n"
               "再平均 → 降 variance。",
               "幾乎是 tabular data 的 default 強 baseline")
    _algo_card(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(3.6),
               "Gradient Boosting", "GradientBoostingRegressor",
               "Boosting：序列訓練多棵 shallow tree，\n"
               "每棵學前一棵的 residual → 降 bias。",
               "通常表現更好但比 RF 慢、對 hyperparameter 較敏感")
    s.shapes.add_picture(
        "diagrams/algo_ensemble.png",
        Inches(1.7), Inches(5.75), width=Inches(10.0), height=Inches(1.30))
    return s


def slide_algo_kernel(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo · 5/5")
    add_title(s, "Family 5：Kernel methods — SVR")
    _algo_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.6),
               "Support Vector Regression", "SVR(kernel='rbf')",
               "在 high-dim feature space 找一個 margin-maximizing hyperplane。\n\n"
               "適合：中等資料量 (<10k)、非線性結構、想要 sparse solution。",
               "⚠ 訓練時間 O(n²)–O(n³)，大資料慢；hyperparameter (C, γ) 重要")
    s.shapes.add_picture(
        "diagrams/algo_kernel.png",
        Inches(0.7), Inches(5.75), width=Inches(7.5), height=Inches(1.30))
    code = ('from sklearn.svm import SVR\n\n'
            'svr = SVR(kernel="rbf", C=10, gamma=0.1)\n\n'
            '# 在 housing 16k samples 上會很慢\n'
            '# → 課堂示範用 subsample 至 3000 rows')
    add_code(s, code, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5), size=14)
    return s


def slide_unified_eval(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo")
    add_title(s, "統一的訓練與評估 — 一個 for-loop 跑完五大家族")
    code = ('models = {\n'
            '    "Linear":       LinearRegression(),\n'
            '    "k-NN (k=5)":   KNeighborsRegressor(n_neighbors=5),\n'
            '    "DecisionTree": DecisionTreeRegressor(random_state=42),\n'
            '    "RandomForest": RandomForestRegressor(n_estimators=100, random_state=42),\n'
            '    "SVR (RBF)":    SVR(kernel="rbf", C=10, gamma=0.1),\n'
            '}\n\n'
            'for name, model in models.items():\n'
            '    pipe = make_pipeline(full_pipeline, model)\n'
            '    scores = -cross_val_score(pipe, X, y, cv=5,\n'
            '                              scoring="neg_root_mean_squared_error")\n'
            '    print(f"{name:15s}  RMSE = {scores.mean():.0f} ± {scores.std():.0f}")')
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(12), Inches(4.5), size=13)
    return s


def _bar_chart(slide, left, top, width, height, values, labels,
               *, colors=None, max_val=None):
    """Simple bar chart from connectors + rectangles."""
    pad_l = Inches(0.9); pad_b = Inches(0.7); pad_t = Inches(0.3); pad_r = Inches(0.2)
    plot_l = left + pad_l
    plot_r = left + width - pad_r
    plot_t = top + pad_t
    plot_b = top + height - pad_b
    plot_w = plot_r - plot_l
    plot_h = plot_b - plot_t
    # axes
    ax_y = slide.shapes.add_connector(1, plot_l, plot_t, plot_l, plot_b)
    ax_y.line.color.rgb = TEXT_MUTED; ax_y.line.width = Pt(0.75)
    ax_x = slide.shapes.add_connector(1, plot_l, plot_b, plot_r, plot_b)
    ax_x.line.color.rgb = TEXT_MUTED; ax_x.line.width = Pt(0.75)
    n = len(values)
    bar_gap = plot_w / (n * 2 + 1)
    bar_w = bar_gap
    mx = max_val if max_val else max(values) * 1.1
    if colors is None:
        colors = [ACCENT_TEAL] * n
    for i, (v, lab, c) in enumerate(zip(values, labels, colors)):
        x = plot_l + bar_gap + i * 2 * bar_gap
        h = plot_h * (v / mx)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, plot_b - h, bar_w, h)
        bar.fill.solid(); bar.fill.fore_color.rgb = c
        bar.line.fill.background()
        # value label on top
        add_text(slide, str(int(v)),
                 x - Inches(0.05), plot_b - h - Inches(0.30),
                 bar_w + Inches(0.1), Inches(0.28),
                 size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        # x-tick label
        add_text(slide, lab,
                 x - Inches(0.15), plot_b + Inches(0.10),
                 bar_w + Inches(0.30), Inches(0.50),
                 size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER,
                 font=FONT_CJK)


def slide_results_chart(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo")
    add_title(s, "California Housing — 5-fold CV RMSE (USD, 越低越好)")
    labels = ["Linear", "k-NN", "DecisionTree", "RandomForest", "GradBoost", "SVR"]
    values = [68628, 56500, 69100, 49500, 51800, 68000]
    colors = [ACCENT_TEAL, ACCENT_TEAL, ACCENT_AMBER,
              ACCENT_GREEN, ACCENT_GREEN, ACCENT_TEAL]
    _bar_chart(s, Inches(0.7), Inches(2.0), Inches(8.5), Inches(4.5),
               values, labels, colors=colors, max_val=80000)
    items = [
        "RandomForest > DecisionTree → ensemble 修掉 single tree 的 variance",
        "Linear ≈ SVR → 此資料線性結構主導",
        "k-NN 中等 → 區域結構有但被 noise 干擾",
    ]
    add_bullets(s, items, Inches(9.4), Inches(2.3), Inches(3.5), Inches(3.5),
                size=14)
    add_text(s, "RMSE 與 target 同單位（USD），比 R² 更好解讀。",
             Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
             size=14, color=TEXT_MUTED, font=FONT_CJK, italic=True)
    return s


def slide_bias_variance(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§4 · Algorithm Zoo")
    add_title(s, "Bias–Variance Tradeoff — 選 model 不是「選最強」")
    # Horizontal spectrum
    L = Inches(0.9); R = Inches(12.4); Y = Inches(3.5)
    arrow_l = s.shapes.add_connector(1, L, Y, R, Y)
    arrow_l.line.color.rgb = TEXT_MUTED; arrow_l.line.width = Pt(2)
    add_text(s, "high bias", L, Y - Inches(0.5),
             Inches(2), Inches(0.35), size=13, bold=True, color=ACCENT_RED)
    add_text(s, "high variance", R - Inches(2), Y - Inches(0.5),
             Inches(2), Inches(0.35), size=13, bold=True, color=ACCENT_RED,
             align=PP_ALIGN.RIGHT)
    add_text(s, "(underfit)", L, Y + Inches(0.15),
             Inches(2), Inches(0.3), size=12, color=TEXT_MUTED, italic=True)
    add_text(s, "(overfit)", R - Inches(2), Y + Inches(0.15),
             Inches(2), Inches(0.3), size=12, color=TEXT_MUTED, italic=True,
             align=PP_ALIGN.RIGHT)
    # Pins
    pin_labels = ["Linear", "Ridge", "RF", "Deep DT"]
    pin_x = [Inches(1.5), Inches(4.5), Inches(8.0), Inches(11.5)]
    for lab, x in zip(pin_labels, pin_x):
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x - Inches(0.10), Y - Inches(0.10),
                                 Inches(0.20), Inches(0.20))
        pin.fill.solid(); pin.fill.fore_color.rgb = ACCENT_TEAL
        pin.line.fill.background()
        add_text(s, lab, x - Inches(0.7), Y + Inches(0.45),
                 Inches(1.4), Inches(0.35),
                 size=13, bold=True, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER)
    add_callout(s,
                "Ensemble 的魅力：用「多個 high-variance learner 平均」\n"
                "  → 降 variance 而不增加 bias。",
                Inches(0.7), Inches(5.3), Inches(12), Inches(1.1),
                size=18, accent=ACCENT_GREEN)
    return s


def slide_practice4(prs):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_practice_badge(s)
    add_title(s, "Hands-on 4 — 加入 Ridge 與 GradientBoosting")
    add_text(s,
             "任務：把 Ridge(alpha=1.0) 與 GradientBoostingRegressor(n_estimators=100) "
             "加入比較表。哪一個 RMSE 較低？解釋為什麼 ensemble 通常贏 single tree。",
             Inches(0.7), Inches(1.95), Inches(12), Inches(1.4),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    code = ('from sklearn.linear_model import Ridge\n'
            'from sklearn.ensemble import GradientBoostingRegressor\n\n'
            'extra = {\n'
            '    "Ridge (α=1)":      Ridge(alpha=1.0),\n'
            '    "GradientBoosting": GradientBoostingRegressor(\n'
            '        n_estimators=100, random_state=42),\n'
            '}\n'
            'for name, model in extra.items():\n'
            '    pipe = make_pipeline(full_pipeline, model)\n'
            '    rmse = -cross_val_score(pipe, X, y, cv=5,\n'
            '              scoring="neg_root_mean_squared_error").mean()\n'
            '    print(f"{name:20s}  RMSE = {rmse:.0f}")')
    add_code(s, code, Inches(0.7), Inches(3.5), Inches(12), Inches(3.0), size=13)
    return s


# ============================================================
# §5 Model selection
# ============================================================
def slide_section5(prs):
    return slide_section(prs, "Section 5",
                         "Model Evaluation & Tuning",
                         "Cross-validation · Grid / Randomized search")


def slide_no_train_acc(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§5 · Evaluation")
    add_title(s, "為什麼不能用 train accuracy 比較 model？")
    code = ('# WRONG — 在 train set 上算 RMSE\n'
            'tree = DecisionTreeRegressor().fit(X_train, y_train)\n'
            'tree.score(X_train, y_train)   # → ~1.0 看似完美\n\n'
            '# 但...\n'
            'tree.score(X_test, y_test)     # → 0.62  其實爛透了')
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(12), Inches(2.5), size=14)
    add_callout(s,
                "DecisionTree(max_depth=None) 的 train RMSE 趨近 0 — \n"
                "它記住了所有 noise，這叫 overfitting。",
                Inches(0.7), Inches(4.9), Inches(12), Inches(1.1),
                size=18, accent=ACCENT_RED)
    add_text(s, "解法：用 cross-validation — 在「沒看過」的 fold 上評估。",
             Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
             size=18, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    return s


def slide_kfold_diagram(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§5 · Evaluation")
    add_title(s, "K-Fold Cross-Validation")
    # Diagram
    L = Inches(1.5); top0 = Inches(2.3)
    fold_w = Inches(10.5); fold_h = Inches(0.5)
    n_folds = 5
    seg_w = fold_w / n_folds
    for f in range(n_folds):
        y = top0 + f * (fold_h + Inches(0.10))
        add_text(s, f"Fold {f+1}",
                 L - Inches(1.2), y + Inches(0.10),
                 Inches(1.0), Inches(0.35),
                 size=13, bold=True, color=TEXT_DARK)
        for seg in range(n_folds):
            x = L + seg * seg_w
            is_test = seg == f
            box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, seg_w, fold_h)
            box.fill.solid()
            box.fill.fore_color.rgb = ACCENT_AMBER if is_test else ACCENT_TEAL
            box.line.color.rgb = TEXT_LIGHT; box.line.width = Pt(1)
            add_text(s, "TEST" if is_test else "train",
                     x, y + Inches(0.10),
                     seg_w, Inches(0.35),
                     size=11, bold=True, color=TEXT_LIGHT,
                     align=PP_ALIGN.CENTER)
    # Code
    code = ('from sklearn.model_selection import cross_val_score\n'
            'scores = cross_val_score(model, X, y, cv=5,\n'
            '                         scoring="neg_root_mean_squared_error")\n'
            'print(f"RMSE: {-scores.mean():.0f} ± {scores.std():.0f}")')
    add_code(s, code, Inches(0.7), Inches(5.5), Inches(12), Inches(1.5), size=14)
    return s


def slide_grid_vs_random(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§5 · Evaluation")
    add_title(s, "Grid Search vs. Randomized Search")
    rows = [
        ("GridSearchCV", "窮舉所有 (param 組合 × cv folds)",
         "< 4 個 hyperparameter 且離散值不多"),
        ("RandomizedSearchCV", "在分布上抽 n_iter 次",
         "連續 hyperparameter、或值很多"),
        ("HalvingRandomSearchCV", "Successive halving — 逐步淘汰差的組合",
         "預算極度有限"),
    ]
    add_table(s, ["方法", "策略", "適用情境"], rows,
              Inches(0.7), Inches(2.0), Inches(12), Inches(2.3),
              col_widths=[Inches(3.0), Inches(4.5), Inches(4.5)],
              header_size=14, body_size=13, first_col_bold=True)
    code = ('from sklearn.model_selection import RandomizedSearchCV\n'
            'from scipy.stats import randint\n\n'
            'param_dist = {\n'
            '    "n_estimators": randint(50, 300),\n'
            '    "max_features": randint(2, 8),\n'
            '}\n'
            'rand = RandomizedSearchCV(\n'
            '    RandomForestRegressor(random_state=42),\n'
            '    param_dist, n_iter=20, cv=5,\n'
            '    scoring="neg_root_mean_squared_error", n_jobs=-1)\n'
            'rand.fit(X_train, y_train)\n'
            'print(rand.best_params_)')
    add_code(s, code, Inches(0.7), Inches(4.5), Inches(12), Inches(2.4), size=13)
    return s


def slide_final_test(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§5 · Evaluation")
    add_title(s, "最終 test set 評估 — 只能做一次")
    add_callout(s,
                "對 test set 看了結果後又回去調 hyperparameter\n"
                "  = 把 test set 變成 train set 的一部分。",
                Inches(0.7), Inches(2.0), Inches(12), Inches(1.4),
                size=20, accent=ACCENT_RED)
    code = ('final_model = rand.best_estimator_\n'
            'X_test = strat_test.drop("median_house_value", axis=1)\n'
            'y_test = strat_test["median_house_value"]\n\n'
            'from sklearn.metrics import root_mean_squared_error\n'
            'final_rmse = root_mean_squared_error(\n'
            '    y_test, final_model.predict(X_test))\n'
            'print(f"FINAL TEST RMSE: {final_rmse:.0f}")')
    add_code(s, code, Inches(0.7), Inches(3.8), Inches(12), Inches(2.5), size=14)
    add_text(s,
             "如果 test RMSE 不滿意 — 接受它，並寫進 limitation。",
             Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
             size=18, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    return s


# ============================================================
# §6 Unsupervised
# ============================================================
def slide_section6(prs):
    return slide_section(prs, "Section 6",
                         "Unsupervised 小品",
                         "Clustering · Anomaly detection (as feature tools)")


def slide_kmeans_feature(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§6 · Unsupervised")
    add_title(s, "KMeans 當 feature engineering 工具")
    items = [
        "Géron 原 notebook：把 (lat, lon) 聚成 10 個 cluster，",
        "再算每個 sample 到 10 個 centroid 的 RBF 相似度 → 10 個新 feature",
        "Cogneuro 類比：對 fMRI ROI time series 做 k-means 找 functional clusters",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(1.8),
                size=18)
    code = ('from sklearn.cluster import KMeans\n'
            'from sklearn.metrics.pairwise import rbf_kernel\n'
            'from sklearn.base import BaseEstimator, TransformerMixin\n\n'
            'class ClusterSimilarity(BaseEstimator, TransformerMixin):\n'
            '    def __init__(self, n_clusters=10, gamma=1.0, random_state=None):\n'
            '        self.n_clusters = n_clusters\n'
            '        self.gamma = gamma\n'
            '        self.random_state = random_state\n\n'
            '    def fit(self, X, y=None):\n'
            '        self.kmeans_ = KMeans(self.n_clusters, n_init=10,\n'
            '                              random_state=self.random_state).fit(X)\n'
            '        return self\n\n'
            '    def transform(self, X):\n'
            '        return rbf_kernel(X, self.kmeans_.cluster_centers_, gamma=self.gamma)')
    add_code(s, code, Inches(0.7), Inches(4.0), Inches(12), Inches(3.0), size=12)
    return s


def slide_isolation_forest(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§6 · Unsupervised")
    add_title(s, "IsolationForest — 多變量 outlier 偵測")
    add_callout(s,
                "比 mean ± 3 SD 強：能找到「每個 feature 單看都正常，\n"
                "但組合起來很奇怪」的 multivariate outlier。",
                Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
                size=18, accent=ACCENT_AMBER)
    code = ('from sklearn.ensemble import IsolationForest\n\n'
            'iso = IsolationForest(contamination=0.05, random_state=42)\n'
            'outlier_mask = iso.fit_predict(X_train) == -1\n'
            'print(f"Flagged {outlier_mask.sum()} / {len(X_train)} outliers")')
    add_code(s, code, Inches(0.7), Inches(3.7), Inches(12), Inches(1.8), size=14)
    items = [
        "在 RT 分析中：抓出「正確但太快」或「motor preparation 異常」的 trial",
        "contamination 是「預期 outlier 比例」的先驗 — 0.01–0.05 為常用範圍",
    ]
    add_bullets(s, items, Inches(0.7), Inches(5.7), Inches(12), Inches(1.3),
                size=16)
    return s


# ============================================================
# §7 Cogneuro transfer
# ============================================================
def slide_section7(prs):
    return slide_section(prs, "Section 7",
                         "Cogneuro Transfer",
                         "Stroop RT prediction with the same pipeline")


def slide_rt_pipeline_code(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§7 · Cogneuro transfer")
    add_title(s, "把同一套 pipeline 套到 Stroop RT 預測")
    code = ('# 1. 模擬資料（200 受試者 × 30 trial）\n'
            'df = simulate_stroop(interaction=True)\n\n'
            '# 2. Stratified split by age group\n'
            'df["age_bin"] = pd.cut(df["age"], bins=[20, 35, 50, 65, 75])\n'
            'sss = StratifiedShuffleSplit(n_splits=1, test_size=0.2, random_state=42)\n'
            'for tr, te in sss.split(df, df["age_bin"]):\n'
            '    train, test = df.iloc[tr], df.iloc[te]\n\n'
            '# 3. Pipeline — exactly the same shape as housing\n'
            'prep = ColumnTransformer([\n'
            '    ("num", num_pipe, ["age", "isi", "trial_num"]),\n'
            '    ("cat", OneHotEncoder(), ["congruent"]),\n'
            '])\n\n'
            '# 4. Algorithm zoo on cogneuro data\n'
            'for name, model in [("Linear", LinearRegression()),\n'
            '                    ("RandomForest", RandomForestRegressor(n_estimators=200))]:\n'
            '    pipe = make_pipeline(prep, model)\n'
            '    rmse = -cross_val_score(pipe, X, y, cv=5,\n'
            '              scoring="neg_root_mean_squared_error").mean()\n'
            '    print(f"{name:14s}  RMSE = {rmse:.1f} ms")')
    add_code(s, code, Inches(0.7), Inches(2.0), Inches(12), Inches(4.9), size=12)
    return s


def slide_rt_results(prs):
    s = add_slide_base(prs)
    add_section_label(s, "§7 · Cogneuro transfer")
    add_title(s, "「最佳 model 取決於資料的真實結構」")
    rows = [
        ("Linear DGP (純線性)",       "39.8 ± 1.0 ms", "44.5 ± 1.1 ms", "Linear 勝"),
        ("Interaction DGP (age × congruency)", "51.8 ± 0.4 ms",
         "44.6 ± 1.2 ms", "RF 勝"),
    ]
    add_table(s, ["資料生成方式", "Linear CV RMSE", "RandomForest CV RMSE", "結論"],
              rows,
              Inches(0.5), Inches(2.0), Inches(12.4), Inches(2.0),
              col_widths=[Inches(3.6), Inches(2.7), Inches(3.1), Inches(3.0)],
              header_size=13, body_size=13, first_col_bold=True)
    add_callout(s,
                "結論：沒有「最強的 algorithm」 — 只有「最適合這份資料的 algorithm」。\n"
                "永遠用 CV 比較，永遠對結果保持懷疑。",
                Inches(0.7), Inches(4.6), Inches(12), Inches(1.4),
                size=20, accent=ACCENT_GREEN)
    add_text(s,
             "📄 code/ml/06_cogneuro_rt_pipeline.py 提供完整可重跑程式。",
             Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             size=14, color=TEXT_MUTED, font=FONT_CJK)
    return s


def slide_practice5(prs):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_practice_badge(s)
    add_title(s, "Hands-on 5（課堂收尾）— 改寫生成式")
    add_text(s,
             "任務：在合成資料的 RT 公式中，加入 congruency × age 的互動項："
             "「老年人受 incongruent trial 影響更大」。重新訓練，這次哪個 model 贏？",
             Inches(0.7), Inches(1.95), Inches(12), Inches(1.4),
             size=17, color=TEXT_DARK, font=FONT_CJK)
    code = ('# Before:\n'
            'cong_effect = 0 if congruent else 60\n\n'
            '# After:\n'
            'cong_effect = 0 if congruent else (60 + 4.0 * (age - 45))')
    add_code(s, code, Inches(0.7), Inches(3.5), Inches(12), Inches(1.7), size=14)
    add_callout(s,
                "💡 提示：Tree-based 與 ensemble 比 linear 更能 capture interaction。",
                Inches(0.7), Inches(5.5), Inches(12), Inches(0.8),
                size=16, accent=ACCENT_TEAL, font=FONT_CJK)
    return s


# ============================================================
# Recap + HW + closing
# ============================================================
def slide_recap(prs):
    s = add_slide_base(prs)
    add_section_label(s, "Recap")
    add_title(s, "今天要記住的五件事")
    items = [
        "先 frame 再 code；先切 test set 再 EDA",
        "Pipeline 是 leakage 防火牆 — fit on train, transform on test",
        "用 CV 比較 model，不用 train accuracy",
        "Test set 只能評估一次",
        "Ensemble (RandomForest, GradientBoosting) 是 tabular data 的強力 baseline",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(4.0),
                size=22)
    return s


def slide_pitfalls(prs):
    s = add_slide_base(prs)
    add_section_label(s, "Common Pitfalls")
    add_title(s, "最常見的五個錯誤")
    rows = [
        ("對 train+test 一起 fit StandardScaler", "data leakage，test RMSE 過度樂觀",
         "用 Pipeline，對 train fit、對 test transform"),
        ("反覆調參直到 test 變好", "test 變成 train 的一部分",
         "嚴格分 validation set，或用 nested CV"),
        ("OneHot 忘了 handle_unknown='ignore'", "test 出現新類別時 crash",
         "加 handle_unknown='ignore'"),
        ("Tree 還做 StandardScaler", "沒影響，但浪費計算",
         "對 RF/GBM 可省略 scaling"),
        ("用 R² 比 regression model", "對 outlier 敏感、難解讀",
         "用 RMSE（與 target 同單位）"),
    ]
    add_table(s, ["錯誤", "後果", "修正"], rows,
              Inches(0.5), Inches(2.0), Inches(12.4), Inches(4.5),
              col_widths=[Inches(4.0), Inches(4.2), Inches(4.2)],
              header_size=13, body_size=12, first_col_bold=True)
    return s


def slide_homework(prs):
    s = add_slide_base(prs)
    add_section_label(s, "Homework")
    add_title(s, "Week 15 作業 — Flanker RT prediction")
    items = [
        "情境：合成的 Flanker task dataset（250 受試者 × 80 trial）",
        "Task 1：EDA + stratified train/test split (10%)",
        "Task 2：Pipeline 設計含 boolean pass-through (15%)",
        "Task 3：至少 3 類 algorithm 用 5-fold CV 比較 (25%)",
        "Task 4：對最佳 model 做 RandomizedSearchCV (20%)",
        "Task 5：Final test RMSE — 只算一次 (10%)",
        "Task 6 + 7：Feature importance 解釋 + 一頁報告 (20%)",
        "Due：2026-06-11 23:59 上傳 eeclass · .ipynb + report.md",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(4.5),
                size=17)
    add_text(s,
             "📄 詳見 week-15-homework.md",
             Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             size=14, color=ACCENT_PRIMARY, font=FONT_MONO)
    return s


def slide_resources(prs):
    s = add_slide_base(prs)
    add_section_label(s, "References")
    add_title(s, "延伸閱讀")
    items = [
        "Géron, A. (2023). Hands-on Machine Learning, 3rd ed. O'Reilly. Ch. 2",
        "github.com/ageron/handson-ml3  ←  原 notebook",
        "scikit-learn user guide：cross-validation / preprocessing 章節",
        "Hastie, Tibshirani, Friedman (2009). The Elements of Statistical Learning. — bias-variance 數學",
        "Varoquaux & Cheplygina (2022). ML for medical imaging: methodological failures. npj Digital Medicine, 5(48). — 必讀",
    ]
    add_bullets(s, items, Inches(0.7), Inches(2.0), Inches(12), Inches(4.0),
                size=17, font=FONT_SANS)
    return s


def slide_closing(prs):
    s = add_slide_base(prs, bg=BG_SECTION, band=False)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))
    add_text(s, "Thank you",
             Inches(0.9), Inches(2.5), Inches(12), Inches(1.5),
             size=64, bold=True, color=TEXT_LIGHT, font=FONT_SANS)
    add_text(s, "下週：Final Project Workshop",
             Inches(0.9), Inches(4.0), Inches(12), Inches(0.7),
             size=24, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, "Erik Chang  ·  ACL@NCU  ·  audachang@gmail.com",
             Inches(0.9), Inches(6.4), Inches(12), Inches(0.4),
             size=14, color=TEXT_LIGHT)
    return s


# ============================================================
# Main
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build in order, then add footers with running page numbers.
    builders = [
        slide_title,
        slide_objectives,
        slide_agenda,
        slide_why,

        slide_section1,
        slide_four_questions,
        slide_taxonomy_overview,
        slide_practice1,

        slide_section2,
        slide_data_overview,
        slide_golden_rule,
        slide_stratified,
        slide_practice2,

        slide_section3,
        slide_pipeline_why,
        slide_pipeline_flow,
        slide_imputation,
        slide_encoding,
        slide_scaling,
        slide_columntransformer,
        slide_practice3,

        slide_section4,
        slide_taxonomy_table,
        slide_algo_linear,
        slide_algo_knn,
        slide_algo_tree,
        slide_algo_ensemble,
        slide_algo_kernel,
        slide_unified_eval,
        slide_results_chart,
        slide_bias_variance,
        slide_practice4,

        slide_section5,
        slide_no_train_acc,
        slide_kfold_diagram,
        slide_grid_vs_random,
        slide_final_test,

        slide_section6,
        slide_kmeans_feature,
        slide_isolation_forest,

        slide_section7,
        slide_rt_pipeline_code,
        slide_rt_results,
        slide_practice5,

        slide_recap,
        slide_pitfalls,
        slide_homework,
        slide_resources,
        slide_closing,
    ]

    # Speaker notes (optional)
    try:
        from speaker_notes import NOTES
    except ImportError:
        NOTES = {}

    total = len(builders)
    md_lines = [
        "# Week 15 — End-to-End Machine Learning · Speaker Notes",
        "",
        "> 自動生成自 `speaker_notes.py`。請編輯該檔，不要直接編輯本檔。",
        "> 對應的投影片：`week-15-slides.pptx` (49 張)。",
        "",
        "---",
        "",
    ]
    for i, fn in enumerate(builders, start=1):
        slide = fn(prs)
        # Skip footer on section dividers (dark bg) and title/closing
        is_dark = fn in (slide_title, slide_closing,
                         slide_section1, slide_section2, slide_section3,
                         slide_section4, slide_section5, slide_section6,
                         slide_section7)
        if not is_dark:
            add_footer(slide, i, total)

        # Attach speaker notes
        short_title, body = NOTES.get(i, (f"Slide {i}", ""))
        if body.strip():
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = body.strip()

        # Append to md transcript
        md_lines.append(f"## Slide {i:02d} · {short_title}")
        md_lines.append("")
        if body.strip():
            md_lines.append(body.strip())
        else:
            md_lines.append("_(no notes)_")
        md_lines.append("")
        md_lines.append("---")
        md_lines.append("")

    out = "week-15-slides.pptx"
    prs.save(out)
    print(f"✅ Wrote {out} with {total} slides.")

    md_out = "week-15-speaker-notes.md"
    with open(md_out, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"OK Wrote {md_out}")


if __name__ == "__main__":
    build()
