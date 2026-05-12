"""Build Week 12 lecture slides — Streamlit Caching & The Data Analysis Pipeline.

Visual style: ACL@NCU — reused from Week 11 (white BG, navy + teal accents,
thin teal top band on content slides, navy section dividers, amber break).

Output: week-12-slides.pptx (16:9, 42 slides for a ~170-min class).
Run:  python build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# ACL@NCU palette (mirrors Week 11)
# ============================================================
BG_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_OFFWHITE    = RGBColor(0xF7, 0xF8, 0xFA)
BG_SECTION     = RGBColor(0x14, 0x32, 0x5C)   # deep navy
BG_BREAK       = RGBColor(0xFB, 0xEA, 0xC0)
BG_PRACTICE    = RGBColor(0xEC, 0xF7, 0xF6)   # light teal tint for hands-on
TEXT_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED     = RGBColor(0x5F, 0x6B, 0x83)
ACCENT_PRIMARY = RGBColor(0x14, 0x32, 0x5C)   # navy
ACCENT_TEAL    = RGBColor(0x0D, 0x9B, 0x9B)   # teal — ACL identity
ACCENT_AMBER   = RGBColor(0xE8, 0xA1, 0x2A)
ACCENT_RED     = RGBColor(0xD3, 0x4F, 0x4F)
ACCENT_GREEN   = RGBColor(0x2E, 0x8B, 0x57)
CODE_BG        = RGBColor(0x1E, 0x29, 0x3B)
CODE_COMMENT   = RGBColor(0x8B, 0x9D, 0xB8)
HAIRLINE       = RGBColor(0xE2, 0xE6, 0xEC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SANS = "Calibri"
FONT_CJK  = "Microsoft JhengHei"
FONT_MONO = "Consolas"

DECK_TITLE = "Week 12 — Streamlit Caching & The Data Analysis Pipeline"
FOOTER_STR = "Week 12 — Caching × Data Analysis Pipeline  ·  ACL@NCU"


# ============================================================
# Primitives
# ============================================================
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_top_band(slide, color=ACCENT_TEAL, height=Inches(0.14)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, text, left, top, width, height, *,
             size=20, bold=False, color=TEXT_DARK, font=FONT_SANS,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_section_label(slide, text, color=ACCENT_TEAL):
    add_text(slide, text,
             Inches(0.7), Inches(0.55), Inches(11), Inches(0.35),
             size=14, bold=True, color=color, font=FONT_SANS)


def add_title(slide, title_zh, *, size=32, top=0.95, color=TEXT_DARK):
    add_text(slide, title_zh,
             Inches(0.7), Inches(top), Inches(12), Inches(0.85),
             size=size, bold=True, color=color, font=FONT_CJK)


def add_bullets(slide, items, left, top, width, height, *,
                size=18, color=TEXT_DARK, line_spacing=1.25,
                bullet_color=ACCENT_TEAL, font=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0); tf.margin_right = Inches(0)
    tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    use_font = font or FONT_CJK
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = "▸  "
        b.font.name = FONT_SANS
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = use_font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        if i > 0:
            p.space_before = Pt(size * 0.45)
    return tb


def _color_for_code_line(line, lang):
    s = line.lstrip()
    if not s:
        return None
    if lang in ("bash", "yaml", "python", "markdown") and s.startswith("#"):
        return CODE_COMMENT
    if s.startswith(">"):
        return ACCENT_AMBER
    if s.startswith("$"):
        return ACCENT_GREEN
    return None


def add_code(slide, code, left, top, width, height, *,
             size=14, lang="python", line_spacing=1.15, padding=0.18):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    box.adjustments[0] = 0.04
    tf = box.text_frame
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(padding); tf.margin_bottom = Inches(padding)
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        c = _color_for_code_line(line, lang) or TEXT_LIGHT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(size)
        r.font.color.rgb = c
    return box


def add_card(slide, left, top, width, height, *, fill=BG_OFFWHITE,
             accent=None, accent_w=0.06):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = HAIRLINE
    card.line.width = Pt(0.5)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left, top, Inches(accent_w), height)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
    return card


def add_callout(slide, text, left, top, width, height, *,
                size=24, color=ACCENT_PRIMARY, fill=BG_OFFWHITE,
                accent=ACCENT_TEAL, font=FONT_CJK, align=PP_ALIGN.LEFT):
    add_card(slide, left, top, width, height, fill=fill, accent=accent,
             accent_w=0.10)
    add_text(slide, text,
             left + Inches(0.4), top + Inches(0.2),
             width - Inches(0.6), height - Inches(0.4),
             size=size, bold=True, color=color, font=font,
             align=align, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page, total, *, theme="light"):
    color = TEXT_MUTED if theme == "light" else ACCENT_TEAL
    add_text(slide, FOOTER_STR,
             Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
             size=11, color=color)
    add_text(slide, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=color, align=PP_ALIGN.RIGHT)


def add_demo_ref(slide, path, *, top=6.65):
    """Small italic ribbon citing a demo file path."""
    add_text(slide, f"demo: {path}",
             Inches(0.5), Inches(top), Inches(10), Inches(0.3),
             size=11, italic=True, color=ACCENT_TEAL, font=FONT_MONO)


def add_demo_callout(slide, label, path, *, left=0.85, top=6.40,
                     width=11.6, height=0.5,
                     accent=ACCENT_AMBER):
    """Inline ▶ DEMO badge — more visible than add_demo_ref().

    label : short imperative ("Run this", "Compare with", "Try changing ...")
    path  : monospace file path / command
    """
    add_card(slide, Inches(left), Inches(top), Inches(width), Inches(height),
             fill=BG_OFFWHITE, accent=accent, accent_w=0.08)
    # Badge
    add_text(slide, "▶ DEMO",
             Inches(left + 0.18), Inches(top + 0.10),
             Inches(1.05), Inches(height - 0.15),
             size=13, bold=True, color=accent)
    # Label (CJK)
    add_text(slide, label,
             Inches(left + 1.30), Inches(top + 0.10),
             Inches(5.2), Inches(height - 0.15),
             size=13, color=TEXT_DARK, font=FONT_CJK)
    # Path / command (mono)
    add_text(slide, path,
             Inches(left + 6.55), Inches(top + 0.10),
             Inches(width - 6.7), Inches(height - 0.15),
             size=12, color=ACCENT_PRIMARY, font=FONT_MONO)


def add_table(slide, headers, rows, left, top, width, height, *,
              header_fill=ACCENT_PRIMARY, header_color=TEXT_LIGHT,
              col_widths=None, header_size=13, body_size=12,
              first_col_bold=False):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    t = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w

    # Header
    for ci, h in enumerate(headers):
        cell = t.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = h
        r.font.name = FONT_CJK; r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = header_color

    # Body
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE if ri % 2 else BG_OFFWHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT_CJK
            r.font.size = Pt(body_size)
            r.font.color.rgb = TEXT_DARK
            if first_col_bold and ci == 0:
                r.font.bold = True
    return table_shape


def add_slide_base(prs, *, bg=BG_WHITE, band=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, bg)
    if band:
        add_top_band(s)
    return s


# ============================================================
# Slide builders
# ============================================================
def slide_title(prs):
    s = add_slide_base(prs, band=False)
    add_top_band(s, ACCENT_TEAL, Inches(0.5))

    # Side accent block
    block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(2.5),
                               Inches(0.18), Inches(2.6))
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT_TEAL
    block.line.fill.background()

    add_text(s, "WEEK 12",
             Inches(1.1), Inches(2.4), Inches(11), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, "Streamlit Caching",
             Inches(1.1), Inches(2.95), Inches(12), Inches(0.95),
             size=44, bold=True, color=TEXT_DARK)
    add_text(s, "& The Data Analysis Pipeline",
             Inches(1.1), Inches(3.85), Inches(12), Inches(0.85),
             size=40, bold=True, color=ACCENT_TEAL)
    add_text(s, "從 dashboard 走向「資料是可信的」data app",
             Inches(1.1), Inches(4.85), Inches(12), Inches(0.6),
             size=22, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "NS5116 · Programming & AI Applications in Behavioral Science",
             Inches(1.1), Inches(6.0), Inches(11), Inches(0.4),
             size=14, color=TEXT_MUTED)
    add_text(s, "ACL@NCU  ·  Spring 2026  ·  2026-05-14  ·  張智宏",
             Inches(1.1), Inches(6.4), Inches(11), Inches(0.4),
             size=13, color=TEXT_MUTED)
    return s


def slide_motivation(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "TODAY'S MOTIVATION")
    add_title(s, "為什麼今天要回頭講 app.py？")

    items = [
        "上週介紹了 demo01–demo14 與整合 dashboard app.py",
        "但因為時間關係，沒有逐段拆解整合範例",
        "同時 app.py 用了一個關鍵元件我們還沒解釋：@st.cache_data",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.1),
                size=20)

    # Two-card today summary
    cards = [
        ("Part 1 · 60 min",
         "把 app.py 與 caching 講清楚"),
        ("Part 2 · 90 min",
         "把視角拉開，談一般化的資料分析流程"),
    ]
    x0, y0 = Inches(0.9), Inches(4.3)
    w, h = Inches(5.7), Inches(2.0)
    gap = Inches(0.35)
    for i, (title, body) in enumerate(cards):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=ACCENT_TEAL)
        add_text(s, title, left + Inches(0.25), y0 + Inches(0.2),
                 w - Inches(0.4), Inches(0.5),
                 size=20, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
        add_text(s, body, left + Inches(0.25), y0 + Inches(0.85),
                 w - Inches(0.4), h - Inches(1.0),
                 size=17, color=TEXT_DARK, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_objectives_1(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "LEARNING OBJECTIVES · 1 / 2")
    add_title(s, "本週前半 — 你將能夠⋯")
    items = [
        "解讀 app.py 的七大區塊與 Streamlit rerun model",
        "正確使用 @st.cache_data（ttl, max_entries, show_spinner）",
        "區分 @st.cache_data vs @st.cache_resource",
        "診斷 UnhashableParamError 等典型錯誤",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.2), Inches(11.5), Inches(4),
                size=22)
    add_footer(s, page, total)
    return s


def slide_objectives_2(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "LEARNING OBJECTIVES · 2 / 2")
    add_title(s, "本週後半 — 你將能夠⋯")
    items = [
        "描繪通用 pipeline：load → inspect → describe → fix → re-describe → analyse",
        "用 descriptive statistics 診斷資料品質",
        "以 observation-driven 方式修補資料",
        "區分 cleaning vs analysis 的邊界",
        "把流程拆成 pure functions（為 final project 鋪路）",
    ]
    add_bullets(s, items, Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.5),
                size=20)
    add_footer(s, page, total)
    return s


def slide_schedule(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "TODAY'S 170 MINUTES")
    add_title(s, "Schedule at a Glance")

    rows = [
        ("1.1",  "app.py 整體結構導覽",                "15 min"),
        ("1.1b", "Sidebar — widget → boolean mask",     "5 min"),
        ("1.2",  "Streamlit 執行模型 — 為什麼 rerun 是關鍵", "10 min"),
        ("1.3", "◆  @st.cache_data 深入解析",            "25 min"),
        ("1.4", "@st.cache_resource & 兩種快取比較",      "10 min"),
        ("—",   "Break",                                  "10 min"),
        ("2.1", "資料分析流程總覽",                       "15 min"),
        ("2.2", "Descriptive statistics 作為健康檢查",    "25 min"),
        ("2.3", "◆  Observation-driven fixing",          "30 min"),
        ("2.4", "完整走一遍 — Stroop demo",               "15 min"),
        ("2.5", "Cleaning vs analysis 邊界",              "5 min"),
        ("—",   "Recap & Homework",                       "10 min"),
    ]
    add_table(s, ["#", "主題", "時間"], rows,
              Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.7),
              col_widths=[Inches(1.0), Inches(7.5), Inches(1.8)],
              header_size=14, body_size=12)

    add_text(s, "總計 170 min（含 10 min break）",
             Inches(0.7), Inches(6.75), Inches(12), Inches(0.35),
             size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_footer(s, page, total)
    return s


def slide_divider(prs, page, total, num, label_en, title_zh, mins):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_SECTION)
    add_text(s, f"PART {num}  ·  {mins} min",
             Inches(0.7), Inches(2.6), Inches(12), Inches(0.55),
             size=22, bold=True, color=ACCENT_TEAL)
    add_text(s, label_en,
             Inches(0.7), Inches(3.2), Inches(12.5), Inches(0.95),
             size=42, bold=True, color=TEXT_LIGHT)
    add_text(s, title_zh,
             Inches(0.7), Inches(4.25), Inches(12.5), Inches(0.6),
             size=22, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)
    return s


def slide_break(prs, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_BREAK)
    add_text(s, "BREAK",
             Inches(0), Inches(2.7), Inches(13.333), Inches(1.3),
             size=72, bold=True, color=ACCENT_PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, "10 分鐘休息  ·  回來繼續 Data Analysis Pipeline",
             Inches(0), Inches(4.1), Inches(13.333), Inches(0.7),
             size=24, color=TEXT_DARK, font=FONT_CJK,
             align=PP_ALIGN.CENTER)
    add_text(s, f"{page} / {total}",
             Inches(11.3), Inches(7.05), Inches(1.5), Inches(0.35),
             size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    return s


# -----------------------------------------------------------------
# Part 1 — app.py + caching
# -----------------------------------------------------------------
def slide_app_tour(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.1 APP STRUCTURE")
    add_title(s, "app.py 整體結構導覽")

    add_text(s, "情境：Cognitive Aging Dashboard",
             Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
             size=20, bold=True, color=ACCENT_PRIMARY, font=FONT_CJK)
    add_bullets(s, [
        "n = 400，lifespan 認知測驗（RT / WM / processing speed / MoCA / Stroop）",
        "Sidebar 篩 age × sex × education",
        "四個 tabs：Age trajectory · Distributions · By group · Raw data",
    ], Inches(0.9), Inches(2.65), Inches(11.5), Inches(2.0), size=18)

    # Right column: dashboard mock — quick visual hint of layout
    add_card(s, Inches(7.0), Inches(4.85), Inches(5.6), Inches(1.95),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "Dashboard 主視覺",
             Inches(7.25), Inches(5.05), Inches(5.2), Inches(0.4),
             size=14, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    add_text(s,
        "Sidebar (左) · KPI metrics (上) · 4 tabs (主視覺) · Filtered DataFrame",
        Inches(7.25), Inches(5.5), Inches(5.2), Inches(1.3),
        size=13, color=TEXT_MUTED, font=FONT_CJK)

    add_callout(s,
        "Streamlit app 通常由七個區塊組成 — 下一頁逐一拆解。",
        Inches(0.9), Inches(4.85), Inches(5.9), Inches(1.95),
        size=18, color=ACCENT_PRIMARY)

    add_demo_callout(s, "跟著拆解 7 個區塊", "streamlit run demo/00_week11_app/app.py", accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


def slide_seven_blocks(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.1 SEVEN BLOCKS")
    add_title(s, "一個典型 Streamlit app 的七大區塊")

    blocks = [
        ("1", "Page config", "st.set_page_config()  — 必須最先呼叫"),
        ("2", "Constants",   "DATA_PATH, MEASURES"),
        ("3", "Data loading","@st.cache_data + load_data()   ◆ 今天重點"),
        ("4", "Sidebar widgets", "st.slider, st.multiselect, st.selectbox"),
        ("5", "Filtering",   "boolean mask on DataFrame"),
        ("6", "Header + KPI","st.title, st.metric × 4"),
        ("7", "Tabs",        "st.tabs(...)  ×  4   ← 主視覺"),
    ]
    y = 2.0
    for num, head, body in blocks:
        row_h = 0.58
        add_text(s, num,
                 Inches(0.85), Inches(y), Inches(0.5), Inches(row_h),
                 size=22, bold=True, color=ACCENT_TEAL)
        add_text(s, head,
                 Inches(1.5), Inches(y), Inches(2.6), Inches(row_h),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body,
                 Inches(4.3), Inches(y), Inches(8.5), Inches(row_h),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)
        y += 0.62

    add_text(s, "▲  注意：st.set_page_config() 必須是第一個 Streamlit 呼叫；否則 raise StreamlitAPIException。",
             Inches(0.85), Inches(6.7), Inches(12), Inches(0.35),
             size=13, italic=True, color=ACCENT_RED, font=FONT_CJK)

    add_demo_callout(s, "打開檔案找出每一個區塊", "demo/00_week11_app/app.py", accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


def slide_sidebar(prs, page, total):
    """Sidebar 的角色 — 把 widget 值接到 filtering."""
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.1b  SIDEBAR")
    add_title(s, "Sidebar 的角色 — 把 widget 值接到 filtering")

    # Left column: 5 widgets list (from app.py)
    add_text(s, "app.py 的 5 個 sidebar widgets",
             Inches(0.85), Inches(1.95), Inches(6.5), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    widgets = [
        ("st.slider",      "Age range  →  tuple (age_min, age_max)"),
        ("st.multiselect", "Sex        →  list[\"F\", \"M\"]"),
        ("st.slider",      "Education  →  tuple"),
        ("st.selectbox",   "Measure    →  單值 (with format_func)"),
        ("st.checkbox",    "Regression line  →  bool"),
    ]
    y = Inches(2.4)
    for code, body in widgets:
        add_text(s, code,
                 Inches(0.95), y, Inches(2.4), Inches(0.4),
                 size=13, bold=True, color=ACCENT_PRIMARY, font=FONT_MONO)
        add_text(s, body,
                 Inches(3.4), y, Inches(4.0), Inches(0.4),
                 size=13, color=TEXT_DARK, font=FONT_CJK)
        y += Inches(0.42)

    # Right column: widget → mask code
    add_text(s, "Widget → boolean mask",
             Inches(7.6), Inches(1.95), Inches(5.3), Inches(0.4),
             size=15, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    code = (
        "with st.sidebar:\n"
        "    age_min, age_max = st.slider(...)\n"
        "    sex_choices      = st.multiselect(...)\n"
        "\n"
        "mask = (df_all[\"age\"].between(age_min, age_max)\n"
        "        & df_all[\"sex\"].isin(sex_choices))\n"
        "df = df_all[mask].copy()"
    )
    add_code(s, code,
             Inches(7.6), Inches(2.4), Inches(5.3), Inches(2.55),
             size=12)

    # Bottom: 3 design disciplines
    add_text(s, "三條 sidebar 設計紀律",
             Inches(0.85), Inches(4.95), Inches(12), Inches(0.4),
             size=15, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    disciplines = [
        ("1", "with st.sidebar: 的 scope",
         "widgets 寫在 sidebar；篩選 logic 寫在主程式"),
        ("2", "default= / value= 寫死合理初值",
         "否則使用者打開 app 看到空白會困惑"),
        ("3", "format_func 分離顯示與 key",
         "selectbox 顯示中文、回傳英文 key"),
    ]
    x0 = Inches(0.85)
    w, h = Inches(4.05), Inches(1.45)
    gap = Inches(0.15)
    for i, (num, head, body) in enumerate(disciplines):
        left = x0 + i * (w + gap)
        add_card(s, left, Inches(5.4), w, h, accent=ACCENT_AMBER)
        add_text(s, num, left + Inches(0.18), Inches(5.5),
                 Inches(0.5), Inches(0.45),
                 size=18, bold=True, color=ACCENT_AMBER)
        add_text(s, head, left + Inches(0.65), Inches(5.5),
                 w - Inches(0.8), Inches(0.5),
                 size=13, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, left + Inches(0.18), Inches(6.0),
                 w - Inches(0.3), Inches(0.85),
                 size=12, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s,
        "與 caching 的連結：sidebar widget 是觸發 rerun 的最常見來源 → load_data() 必須要 cache。",
        Inches(0.85), Inches(6.92), Inches(12), Inches(0.3),
        size=12, italic=True, color=ACCENT_PRIMARY, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_practice(prs, page, total, num, title, tasks, reflect=None,
                   answer=None, demo=None):
    s = add_slide_base(prs, bg=BG_PRACTICE)
    add_section_label(s, f"HANDS-ON  ·  {num}")
    add_title(s, title)

    add_bullets(s, tasks, Inches(0.9), Inches(2.05), Inches(11.5), Inches(3.0),
                size=18, bullet_color=ACCENT_PRIMARY)

    if reflect:
        add_card(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.0),
                 fill=BG_WHITE, accent=ACCENT_AMBER)
        add_text(s, "思考",
                 Inches(1.1), Inches(5.08), Inches(2), Inches(0.4),
                 size=14, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
        add_text(s, reflect,
                 Inches(1.1), Inches(5.42), Inches(11.2), Inches(0.5),
                 size=14, color=TEXT_DARK, font=FONT_CJK)

    if answer:
        add_text(s, "預期觀察：" + answer,
                 Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.5),
                 size=13, italic=True, color=ACCENT_GREEN, font=FONT_CJK)

    if demo:
        add_demo_ref(s, demo, top=6.65)

    add_footer(s, page, total)
    return s


def slide_execution_model(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.2 EXECUTION MODEL")
    add_title(s, "Streamlit 的核心執行模型")

    add_callout(s,
        "每當使用者操作任何 widget，Streamlit 都會從頭到尾重新執行整支 app.py。",
        Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
        size=24, color=ACCENT_PRIMARY)

    # Simple textual loop diagram (three arrows)
    steps = ["Widget interaction", "Full script rerun", "New UI"]
    y = Inches(4.6)
    box_w, box_h = Inches(3.3), Inches(1.1)
    gap = Inches(0.45)
    total_w = box_w * 3 + gap * 2
    x0 = (SLIDE_W - total_w) / 2
    for i, label in enumerate(steps):
        left = x0 + i * (box_w + gap)
        add_card(s, left, y, box_w, box_h, fill=BG_OFFWHITE,
                 accent=ACCENT_TEAL)
        add_text(s, label, left, y + Inches(0.3),
                 box_w, Inches(0.5),
                 size=18, bold=True, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER, font=FONT_CJK)
        if i < 2:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       left + box_w + Inches(0.05),
                                       y + Inches(0.43),
                                       gap - Inches(0.10), Inches(0.25))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_TEAL
            arrow.line.fill.background()

    add_text(s, "↻  迴圈 — 每次互動都觸發整輪",
             Inches(0), Inches(6.0), Inches(13.333), Inches(0.5),
             size=16, italic=True, color=TEXT_MUTED, font=FONT_CJK,
             align=PP_ALIGN.CENTER)

    add_demo_callout(s, "並排對比：有 vs 無 @st.cache_data", "streamlit run demo/00_week11_app/app_no_cache.py", accent=ACCENT_RED)
    add_footer(s, page, total)
    return s


def slide_rerun_consequences(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.2 CONSEQUENCES")
    add_title(s, "Rerun 模型的兩個重大後果")

    # Two large cards
    cards = [
        (ACCENT_GREEN, "✓",
         "State management 變簡單",
         "Widget 值自動傳入下一次 rerun；不需要手動處理「更新狀態 → 重繪 UI」的邏輯。"),
        (ACCENT_RED, "▲",
         "昂貴計算被反覆執行",
         "load_data() 每 rerun 都會重讀一次 CSV。換成 API 或 fMRI volume → app 完全不能用。"),
    ]
    x0, y0 = Inches(0.85), Inches(2.0)
    w, h = Inches(5.95), Inches(3.3)
    gap = Inches(0.25)
    for i, (color, icon, head, body) in enumerate(cards):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, fill=BG_OFFWHITE, accent=color)
        add_text(s, icon, left + Inches(0.3), y0 + Inches(0.2),
                 Inches(1), Inches(0.6), size=28, bold=True, color=color)
        add_text(s, head, left + Inches(0.3), y0 + Inches(0.85),
                 w - Inches(0.5), Inches(0.6),
                 size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, left + Inches(0.3), y0 + Inches(1.55),
                 w - Inches(0.5), h - Inches(1.7),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_callout(s,
        "→ 這就是為什麼需要 @st.cache_data；它是 functools.lru_cache 的 Streamlit 強化版。",
        Inches(0.85), Inches(5.6), Inches(11.5), Inches(1.05),
        size=18, color=ACCENT_PRIMARY)

    add_footer(s, page, total)
    return s


def slide_cache_data_core(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3   ◆  @st.cache_data")
    add_title(s, "@st.cache_data — 核心範例")

    code = (
        "@st.cache_data\n"
        "def load_data(path: Path) -> pd.DataFrame:\n"
        "    df = pd.read_csv(path)\n"
        "    df[\"group\"] = pd.Categorical(\n"
        "        df[\"group\"],\n"
        "        categories=[\"young\",\"middle\",\"older\"],\n"
        "        ordered=True,\n"
        "    )\n"
        "    return df"
    )
    add_code(s, code,
             Inches(0.9), Inches(2.05), Inches(7.5), Inches(3.6),
             size=15)

    # Speed comparison badges
    badge_x, badge_y = Inches(9.0), Inches(2.2)
    add_card(s, badge_x, badge_y, Inches(3.5), Inches(1.4),
             fill=BG_OFFWHITE, accent=ACCENT_RED)
    add_text(s, "Without cache", badge_x + Inches(0.2), badge_y + Inches(0.15),
             Inches(3.1), Inches(0.4),
             size=14, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_text(s, "≈ 80 ms / rerun",
             badge_x + Inches(0.2), badge_y + Inches(0.55),
             Inches(3.1), Inches(0.7),
             size=22, bold=True, color=TEXT_DARK)

    add_card(s, badge_x, badge_y + Inches(1.6), Inches(3.5), Inches(1.4),
             fill=BG_OFFWHITE, accent=ACCENT_GREEN)
    add_text(s, "With @st.cache_data",
             badge_x + Inches(0.2), badge_y + Inches(1.75),
             Inches(3.1), Inches(0.4),
             size=14, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_text(s, "≈  0 ms (cache hit)",
             badge_x + Inches(0.2), badge_y + Inches(2.15),
             Inches(3.1), Inches(0.7),
             size=22, bold=True, color=TEXT_DARK)

    add_text(s, "一個裝飾器 — 避免 rerun 時重讀 CSV。",
             Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.5),
             size=18, italic=True, color=ACCENT_PRIMARY, font=FONT_CJK)

    add_demo_callout(s, "看 cache hit 的威力（頁內計時器）", "streamlit run demo/00_week11_app/app_with_cache.py", accent=ACCENT_GREEN)
    add_footer(s, page, total)
    return s


def slide_cache_key(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3 CACHE KEY")
    add_title(s, "Cache key 是怎麼算出來的？")

    keys = [
        ("01", "函式名稱",         "load_data"),
        ("02", "參數的 hash",      "Path 物件 → 字串"),
        ("03", "函式原始碼 hash",  "修改函式內容 → 舊 cache 自動失效"),
    ]
    x0, y0 = Inches(0.9), Inches(2.1)
    w, h = Inches(3.85), Inches(3.5)
    gap = Inches(0.18)
    for i, (num, head, body) in enumerate(keys):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=ACCENT_TEAL)
        add_text(s, num, left + Inches(0.25), y0 + Inches(0.25),
                 Inches(1.5), Inches(0.7),
                 size=28, bold=True, color=ACCENT_TEAL)
        add_text(s, head, left + Inches(0.25), y0 + Inches(1.05),
                 w - Inches(0.4), Inches(0.6),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, left + Inches(0.25), y0 + Inches(1.7),
                 w - Inches(0.4), h - Inches(1.9),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_callout(s,
        "同一個 path 第二次傳入 → 整個函式被跳過，直接回傳上次結果。",
        Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.75),
        size=16, color=ACCENT_PRIMARY)

    add_demo_callout(s, "改 load_data 內容後重啟，cache 自動失效", "demo/01_cache_data/cache_demo.py", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_cache_copy(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3 COPY SEMANTICS")
    add_title(s, "為什麼回傳的是 copy 而不是 reference？")

    add_bullets(s, [
        "底層做 serialize → deserialize → 等同 deep copy",
        "保證：使用者下游怎麼亂改 DataFrame 都不會污染 cache",
        "代價：序列化開銷；> 100 MB 的物件 → 改用 @st.cache_resource",
    ], Inches(0.9), Inches(2.05), Inches(11.5), Inches(2.0), size=18)

    code = (
        "df  = load_data(p)\n"
        "df[\"age\"] = -999          # 你愛怎麼改怎麼改\n"
        "df2 = load_data(p)\n"
        "print(df2[\"age\"].head())  # 仍然是原始值"
    )
    add_code(s, code,
             Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.2),
             size=16)

    add_demo_callout(s, "看 deep copy vs reference 的實際差異", "demo/02_cache_vs_resource/compare.py", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_cache_params(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3 KEY PARAMETERS")
    add_title(s, "@st.cache_data 重要參數")

    rows = [
        ("ttl",            "time-to-live（秒）",       "API 資料每小時刷新"),
        ("max_entries",    "LRU 上限",                 "函式被多種 input 呼叫"),
        ("show_spinner",   "第一次計算的提示文字",     "計算 > 1 秒，給使用者回饋"),
        ("persist=\"disk\"", "cache 寫到磁碟",         "重啟 app 也保留"),
    ]
    add_table(s, ["參數", "用途", "何時用"], rows,
              Inches(1.0), Inches(2.1), Inches(11.3), Inches(3.6),
              col_widths=[Inches(2.5), Inches(4.0), Inches(4.8)],
              header_size=14, body_size=14, first_col_bold=True)

    code = (
        "@st.cache_data(ttl=3600, max_entries=10,\n"
        "               show_spinner=\"Loading...\", persist=\"disk\")\n"
        "def load_data(path): ..."
    )
    add_code(s, code,
             Inches(1.0), Inches(5.40), Inches(11.3), Inches(0.85),
             size=12)

    add_demo_callout(s, "改 TTL=None → 10 觀察 spinner 行為", "demo/01_cache_data/cache_demo.py", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_pitfall_1(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3 PITFALL #1")
    add_title(s, "▲  常見錯誤 1 — Widget 寫在 cached function 內")

    bad = (
        "@st.cache_data\n"
        "def bad_load():\n"
        "    n = st.slider(\"rows\", 10, 100)\n"
        "    # slider 改動後仍回舊結果\n"
        "    return df.sample(n)"
    )
    good = (
        "n  = st.slider(\"rows\", 10, 100)\n"
        "df = good_load(n)\n"
        "# widget 在外面讀\n\n"
        "@st.cache_data\n"
        "def good_load(n): ..."
    )

    # Bad card
    add_card(s, Inches(0.85), Inches(2.05), Inches(5.95), Inches(0.55),
             fill=ACCENT_RED, accent=ACCENT_RED)
    add_text(s, "✗  錯誤示範",
             Inches(1.0), Inches(2.12), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=TEXT_LIGHT, font=FONT_CJK)
    add_code(s, bad,
             Inches(0.85), Inches(2.65), Inches(5.95), Inches(2.6),
             size=14)

    # Good card
    add_card(s, Inches(6.95), Inches(2.05), Inches(5.95), Inches(0.55),
             fill=ACCENT_GREEN, accent=ACCENT_GREEN)
    add_text(s, "✓  正確寫法",
             Inches(7.10), Inches(2.12), Inches(5.5), Inches(0.4),
             size=15, bold=True, color=TEXT_LIGHT, font=FONT_CJK)
    add_code(s, good,
             Inches(6.95), Inches(2.65), Inches(5.95), Inches(2.6),
             size=14)

    add_callout(s,
        "原因：cache key 不包含 widget 的當前值 — slider 改動後仍回傳舊結果。",
        Inches(0.85), Inches(5.55), Inches(12.05), Inches(1.0),
        size=16, color=ACCENT_PRIMARY)

    add_demo_ref(s, "demo/01_cache_data/pitfalls.py")
    add_footer(s, page, total)
    return s


def slide_pitfall_23(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.3 PITFALL #2 & #3")
    add_title(s, "▲  常見錯誤 2 & 3")

    # Pitfall 2
    add_card(s, Inches(0.85), Inches(2.05), Inches(5.95), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_RED)
    add_text(s, "錯誤 2  ·  UnhashableParamError",
             Inches(1.0), Inches(2.2), Inches(5.7), Inches(0.45),
             size=16, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_text(s,
        "傳了 dict / DataFrame 進 cached function。",
        Inches(1.0), Inches(2.7), Inches(5.7), Inches(0.5),
        size=14, color=TEXT_MUTED, font=FONT_CJK)
    add_code(s,
        "@st.cache_data\n"
        "def f(df, _params):   # _ → 跳過 hash\n"
        "    ...",
        Inches(1.0), Inches(3.25), Inches(5.65), Inches(1.5),
        size=13)
    add_text(s,
        "底線參數不參與 key — 你必須自己保證它不影響輸出。",
        Inches(1.0), Inches(4.95), Inches(5.7), Inches(1.4),
        size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)

    # Pitfall 3
    add_card(s, Inches(6.95), Inches(2.05), Inches(5.95), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_RED)
    add_text(s, "錯誤 3  ·  Side effect 被略過",
             Inches(7.1), Inches(2.2), Inches(5.7), Inches(0.45),
             size=16, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_text(s,
        "cached function 第二次以同 input 呼叫時被「跳過」。",
        Inches(7.1), Inches(2.7), Inches(5.7), Inches(0.5),
        size=14, color=TEXT_MUTED, font=FONT_CJK)
    add_code(s,
        "@st.cache_data\n"
        "def fetch(url):\n"
        "    print(f\"Fetching {url}\")  # 第二次不會印！\n"
        "    return requests.get(url).json()",
        Inches(7.1), Inches(3.25), Inches(5.65), Inches(2.0),
        size=13)
    add_text(s,
        "不要把 logging / 寫檔等 side effect 放在 cached function 內。",
        Inches(7.1), Inches(5.4), Inches(5.7), Inches(1.0),
        size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)

    add_demo_callout(s, "三個 pitfalls 並列示範", "streamlit run demo/01_cache_data/pitfalls.py", accent=ACCENT_RED)
    add_footer(s, page, total)
    return s


def slide_cache_vs_resource(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.4   ◆  CACHE COMPARISON")
    add_title(s, "@st.cache_data vs @st.cache_resource")

    rows = [
        ("回傳值處理",   "Serialize → deep copy",      "回傳同一個 reference（singleton）"),
        ("適用場景",     "DataFrame · array · JSON",   "DB conn · ML model · LLM client"),
        ("跨 user session", "各自快取",                 "共享同一個物件"),
        ("副作用風險",   "低（每人都是 copy）",        "高 — 多人改同物件互相影響"),
    ]
    add_table(s, ["", "@st.cache_data", "@st.cache_resource"], rows,
              Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.5),
              col_widths=[Inches(2.8), Inches(4.35), Inches(4.35)],
              header_size=13, body_size=13, first_col_bold=True)

    code = (
        "@st.cache_resource\n"
        "def get_anthropic_client():\n"
        "    from anthropic import Anthropic\n"
        "    return Anthropic(api_key=st.secrets[\"ANTHROPIC_API_KEY\"])"
    )
    add_code(s, code,
             Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.7),
             size=14)
    add_text(s, "範例：Week 14 的 Anthropic client 應該用 cache_resource。",
             Inches(0.9), Inches(6.62), Inches(12), Inches(0.4),
             size=13, italic=True, color=TEXT_MUTED, font=FONT_CJK)

    add_demo_ref(s, "demo/02_cache_vs_resource/compare.py", top=6.70)
    add_footer(s, page, total)
    return s


def slide_decision_rule(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 1 · 1.4 DECISION RULE")
    add_title(s, "Decision Rule")

    add_callout(s,
        "要的是「資料」？  →  @st.cache_data",
        Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.4),
        size=24, color=ACCENT_PRIMARY)
    add_callout(s,
        "要的是「資源 / 連線 / 物件」？  →  @st.cache_resource",
        Inches(0.9), Inches(4.3), Inches(11.5), Inches(1.4),
        size=22, color=ACCENT_PRIMARY, accent=ACCENT_AMBER)

    add_footer(s, page, total)
    return s


# -----------------------------------------------------------------
# Part 2 — Data analysis pipeline
# -----------------------------------------------------------------
def slide_why_pipeline(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.1 MOTIVATION")
    add_title(s, "為什麼今天先講通用流程？")

    add_bullets(s, [
        "資料可能來自 CSV、API、PsychoPy log、SQL、fMRI volume ⋯⋯",
        "raw data → analysable data 永遠隔著相同的中間幾步",
        "跳過 → 結論與資料的雜訊綁在一起，未來 reviewer 會挑剔",
    ], Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.5), size=18)

    # Four source icons → one pipeline
    sources = [("CSV", ACCENT_TEAL), ("API", ACCENT_AMBER),
               ("Log", ACCENT_GREEN), ("Brain", ACCENT_RED)]
    y = Inches(4.85)
    box_w, box_h = Inches(1.7), Inches(1.0)
    x_start = Inches(0.9)
    for i, (label, color) in enumerate(sources):
        left = x_start + i * (box_w + Inches(0.18))
        add_card(s, left, y, box_w, box_h, accent=color)
        add_text(s, label, left, y + Inches(0.25),
                 box_w, Inches(0.5),
                 size=18, bold=True, color=color,
                 align=PP_ALIGN.CENTER, font=FONT_CJK)

    # Arrow + pipeline box
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(8.55), y + Inches(0.35),
                               Inches(0.8), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_PRIMARY
    arrow.line.fill.background()

    add_card(s, Inches(9.55), y, Inches(2.8), Inches(1.0),
             fill=ACCENT_PRIMARY)
    add_text(s, "通用 pipeline",
             Inches(9.55), y + Inches(0.25),
             Inches(2.8), Inches(0.5),
             size=18, bold=True, color=TEXT_LIGHT,
             align=PP_ALIGN.CENTER, font=FONT_CJK)

    add_demo_callout(s, "跑一遍完整 pipeline", "python demo/03_pipeline/pipeline.py", accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


def slide_pipeline_diagram(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.1   ◆  PIPELINE")
    add_title(s, "完整 Pipeline — 六個步驟")

    steps = [
        ("1", "load",         "I/O only, no transform",         ACCENT_TEAL,  False),
        ("2", "inspect",      "shape · dtypes · head",          ACCENT_TEAL,  False),
        ("3", "describe",     "descriptive statistics（診斷）", ACCENT_PRIMARY, True),
        ("4", "fix",          "observation-driven",             ACCENT_PRIMARY, True),
        ("5", "re-describe",  "驗證 fix 沒引入新問題",          ACCENT_TEAL,  False),
        ("6", "analyse",      "回答研究問題",                   ACCENT_TEAL,  False),
    ]
    y0 = Inches(1.95)
    row_h = Inches(0.72)
    for i, (num, name, body, color, star) in enumerate(steps):
        y = y0 + i * row_h
        # Number badge
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(1.1), y + Inches(0.08),
                                    Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(s, num, Inches(1.1), y + Inches(0.13),
                 Inches(0.55), Inches(0.5),
                 size=16, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER)
        # Name
        add_text(s, name, Inches(1.95), y + Inches(0.12),
                 Inches(3.2), Inches(0.5),
                 size=20, bold=True, color=TEXT_DARK)
        # Body
        add_text(s, body, Inches(5.2), y + Inches(0.18),
                 Inches(6.5), Inches(0.5),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)
        # Star
        if star:
            add_text(s, "◆  今天重點",
                     Inches(11.4), y + Inches(0.18),
                     Inches(1.7), Inches(0.5),
                     size=13, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
        # Connector line (except last)
        if i < len(steps) - 1:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(1.36), y + Inches(0.63),
                                      Inches(0.04), Inches(0.18))
            line.fill.solid(); line.fill.fore_color.rgb = HAIRLINE
            line.line.fill.background()

    add_demo_callout(s, "pipeline.py 6 步對應這張圖", "demo/03_pipeline/pipeline.py", accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


def slide_two_traps(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.1 COMMON TRAPS")
    add_title(s, "兩個常見誤區")

    traps = [
        ("✗", "跳過 describe 直接畫圖",
         "你只看到「眼睛能掃到的前幾行」。第 32 row 的 RT = 99999 永遠不會被發現。"),
        ("✗", "先 df.dropna() 再說",
         "沒看缺值樣態（隨機？系統？某 condition 集中？） → 容易引入 selection bias。"),
    ]
    x0, y0 = Inches(0.85), Inches(2.05)
    w, h = Inches(5.95), Inches(4.4)
    gap = Inches(0.25)
    for i, (icon, head, body) in enumerate(traps):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, fill=BG_OFFWHITE, accent=ACCENT_RED)
        add_text(s, icon, left + Inches(0.3), y0 + Inches(0.2),
                 Inches(1), Inches(0.7), size=32, bold=True, color=ACCENT_RED)
        add_text(s, head, left + Inches(0.3), y0 + Inches(1.0),
                 w - Inches(0.5), Inches(0.6),
                 size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body, left + Inches(0.3), y0 + Inches(1.85),
                 w - Inches(0.5), h - Inches(2.0),
                 size=16, color=TEXT_MUTED, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_descriptive_stats(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.2   ◆  DIAGNOSTIC")
    add_title(s, "Descriptive Statistics 作為「健康檢查」")

    add_bullets(s, [
        "重新定位：不是「論文 Table 1」，而是資料品質的診斷儀器",
        "五件每次都做的事 →",
    ], Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.2), size=18)

    code = (
        "print(df.shape)\n"
        "print(df.dtypes)\n"
        "print(df.isnull().sum())\n"
        "print(df.describe(include=\"all\"))\n"
        "for col in df.select_dtypes(\"object\"):\n"
        "    print(df[col].value_counts(dropna=False).head(10))"
    )
    add_code(s, code,
             Inches(0.9), Inches(3.5), Inches(11.5), Inches(2.5),
             size=15)

    add_callout(s,
        "關鍵：每個指令不只是看數字 — 跑完後問「這合不合理？」",
        Inches(0.9), Inches(5.80), Inches(11.5), Inches(0.55),
        size=16, color=ACCENT_PRIMARY)

    add_demo_callout(s, "用這 5 個指令掃過 messy data", "demo/data/messy_stroop.csv", accent=ACCENT_TEAL)
    add_footer(s, page, total)
    return s


def slide_diagnose_table(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.2   ◆  CHEAT SHEET")
    add_title(s, "從 summary 讀出問題 — 速查表")

    rows = [
        ("某欄 count 較少",                        "大量缺值"),
        ("dtype 是 object，但內容像數字",          "混入 \"NA\" / \"-\" / 空字串"),
        ("min / max 物理上不可能",                 "sentinel value 或編碼錯誤"),
        ("mean ≠ median",                          "分佈偏斜或極端 outlier"),
        ("std 異常大 / ≈ 0",                       "outlier 或近常數"),
        ("value_counts 出現 \"NA\",\"-\",\"\",999", "缺值偽裝"),
        ("level 重複（M / male / Male）",          "編碼不一致"),
        ("某 condition trial 數遠少",              "不平衡設計或 logging 失敗"),
    ]
    add_table(s, ["觀察到的現象", "可能的問題"], rows,
              Inches(1.4), Inches(2.0), Inches(10.5), Inches(4.30),
              col_widths=[Inches(5.3), Inches(5.2)],
              header_size=13, body_size=12, first_col_bold=True)

    add_demo_callout(s, "在 messy_stroop.csv 上練習找問題", "demo/data/messy_stroop.csv", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_viz_diagnostic(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.2 VISUAL DIAGNOSTIC")
    add_title(s, "視覺化也是 descriptive statistics")

    plots = [
        ("df.hist()",            "每個 numeric 欄位的分佈", ACCENT_TEAL),
        ("df.plot.box()",        "outlier 一目了然",        ACCENT_AMBER),
        ("df.isnull().sum()\n  .plot.bar()", "缺值集中在哪幾欄", ACCENT_RED),
    ]
    x0, y0 = Inches(0.9), Inches(2.1)
    w, h = Inches(3.95), Inches(3.5)
    gap = Inches(0.2)
    for i, (code, body, color) in enumerate(plots):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=color)
        # multiline-aware
        tb = s.shapes.add_textbox(left + Inches(0.25), y0 + Inches(0.4),
                                   w - Inches(0.4), Inches(0.9))
        tf = tb.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.margin_top = Inches(0); tf.margin_bottom = Inches(0)
        tf.word_wrap = True
        for li, ln in enumerate(code.split("\n")):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = ln
            r.font.name = FONT_MONO
            r.font.size = Pt(15)
            r.font.bold = True
            r.font.color.rgb = color
        # Mock plot area
        plot_area = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       left + Inches(0.4), y0 + Inches(1.3),
                                       w - Inches(0.8), Inches(1.4))
        plot_area.fill.solid(); plot_area.fill.fore_color.rgb = BG_OFFWHITE
        plot_area.line.color.rgb = HAIRLINE
        plot_area.line.width = Pt(0.5)
        add_text(s, "[ plot ]", left + Inches(0.4), y0 + Inches(1.7),
                 w - Inches(0.8), Inches(0.6),
                 size=18, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
                 italic=True)
        add_text(s, body,
                 left + Inches(0.25), y0 + Inches(2.85),
                 w - Inches(0.4), Inches(0.5),
                 size=14, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s, "數字摘要會 average-out 結構性問題；圖才看得到。",
             Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.5),
             size=17, italic=True, color=ACCENT_PRIMARY, font=FONT_CJK,
             align=PP_ALIGN.CENTER)

    add_footer(s, page, total)
    return s


def slide_obs_driven_principle(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.3   ◆  CORE PRINCIPLE")
    add_title(s, "Observation-Driven Fixing — 核心原則")

    add_callout(s,
        "每一個 cleaning 動作都應該對應到 §2.2 的具體觀察，且你能說出它的代價。",
        Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6),
        size=22, color=ACCENT_PRIMARY)

    # Three-step arrow: 觀察 → 動作 → 代價
    steps = [
        ("觀察",   "從 descriptive statistics", ACCENT_TEAL),
        ("動作",   "對應的 fix（coerce / drop / fill ...）", ACCENT_PRIMARY),
        ("代價",   "你必須能說出這個決定犧牲了什麼", ACCENT_AMBER),
    ]
    y = Inches(4.85)
    box_w, box_h = Inches(3.5), Inches(1.4)
    gap = Inches(0.3)
    total_w = box_w * 3 + gap * 2
    x0 = (SLIDE_W - total_w) / 2
    for i, (head, body, color) in enumerate(steps):
        left = x0 + i * (box_w + gap)
        add_card(s, left, y, box_w, box_h, accent=color)
        add_text(s, head, left + Inches(0.25), y + Inches(0.2),
                 box_w - Inches(0.4), Inches(0.45),
                 size=20, bold=True, color=color, font=FONT_CJK)
        add_text(s, body, left + Inches(0.25), y + Inches(0.7),
                 box_w - Inches(0.4), box_h - Inches(0.8),
                 size=14, color=TEXT_DARK, font=FONT_CJK)
        if i < 2:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       left + box_w + Inches(0.0),
                                       y + Inches(0.55),
                                       gap, Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()

    add_demo_callout(s, "clean() docstring 寫出對應關係", "demo/03_pipeline/pipeline.py::clean", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_obs_action_cost(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.3   ◆  MAPPING TABLE")
    add_title(s, "觀察 → 動作 → 代價 對應表")

    rows = [
        ("object dtype 像數字",       "pd.to_numeric(errors=\"coerce\")", "無法解析 → NaN 悄悄發生"),
        ("sentinel 偽裝缺值",          "replace(...) → to_numeric",         "漏掉 sentinel 就污染分析"),
        ("缺值少且隨機",               "dropna(subset=...)",                "損失 n"),
        ("缺值多 / 系統性",            "fillna(策略)",                      "填補本身是一個假設"),
        ("物理不可能值",               "between(...) filter",               "範圍錯 → 剔除真實 outlier"),
        ("categorical 不一致",         "str.lower() + replace",             "合併不該合的 level"),
        ("time 非 datetime",           "pd.to_datetime(errors=\"coerce\")", "format 不對整欄變 NaT"),
    ]
    add_table(s, ["觀察", "動作", "代價"], rows,
              Inches(0.7), Inches(2.0), Inches(11.95), Inches(4.7),
              col_widths=[Inches(3.6), Inches(4.35), Inches(4.0)],
              header_size=14, body_size=12, first_col_bold=True)

    add_footer(s, page, total)
    return s


def slide_cleaning_discipline(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.3 DISCIPLINE")
    add_title(s, "三條 cleaning 紀律")

    items = [
        ("01", "每步留 log",         "print(f\"dropna: {before} → {len(df)}\")"),
        ("02", "改完再跑一次 describe()", "沒驗證 = 沒做"),
        ("03", "不要 chain method",   "一行一動作，error 才追得到"),
    ]
    y0 = Inches(2.1)
    for i, (num, head, body) in enumerate(items):
        y = y0 + i * Inches(0.72)
        add_text(s, num,
                 Inches(0.9), y, Inches(0.6), Inches(0.6),
                 size=24, bold=True, color=ACCENT_TEAL)
        add_text(s, head,
                 Inches(1.7), y + Inches(0.05),
                 Inches(4.0), Inches(0.5),
                 size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body,
                 Inches(5.9), y + Inches(0.1),
                 Inches(6.8), Inches(0.5),
                 size=15, color=TEXT_MUTED, font=FONT_MONO)

    # Chain vs step-by-step
    bad = (
        "df = (df.replace(-999, np.nan)\n"
        "        .dropna()\n"
        "        .query(\"rt_ms.between(150, 3000)\")\n"
        "        .reset_index(drop=True))"
    )
    good = (
        "df = df.replace(-999, np.nan)\n"
        "print(df.isnull().sum())\n"
        "df = df.dropna(subset=[\"age\"])\n"
        "df = df[df[\"rt_ms\"].between(150, 3000)]"
    )
    add_text(s, "✗ Chain method",
             Inches(0.9), Inches(4.6), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_code(s, bad,
             Inches(0.9), Inches(5.0), Inches(5.95), Inches(1.7),
             size=12)
    add_text(s, "✓ Step-by-step",
             Inches(7.05), Inches(4.6), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=ACCENT_GREEN, font=FONT_CJK)
    add_code(s, good,
             Inches(7.05), Inches(5.0), Inches(5.85), Inches(1.7),
             size=12)

    add_demo_callout(s, "每步都有 print log", "demo/03_pipeline/pipeline.py", accent=ACCENT_AMBER)
    add_footer(s, page, total)
    return s


def slide_error_modes(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.3 ERROR MODES")
    add_title(s, "▲  兩種錯誤模式")

    # A
    add_card(s, Inches(0.85), Inches(2.05), Inches(5.95), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_RED)
    add_text(s, "錯誤 A  ·  無腦 dropna()",
             Inches(1.0), Inches(2.2), Inches(5.7), Inches(0.5),
             size=18, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_code(s,
        "df = df.dropna()  # 我不在乎 bias",
        Inches(1.0), Inches(2.85), Inches(5.65), Inches(0.85),
        size=13)
    add_text(s,
        "若 age 缺值集中在年長者（他們較不填寫）→ 系統性低估 cognitive aging 效應。",
        Inches(1.0), Inches(3.95), Inches(5.7), Inches(2.0),
        size=14, color=TEXT_MUTED, font=FONT_CJK)

    # B
    add_card(s, Inches(6.95), Inches(2.05), Inches(5.95), Inches(4.4),
             fill=BG_OFFWHITE, accent=ACCENT_RED)
    add_text(s, "錯誤 B  ·  把 analysis 偽裝成 cleaning",
             Inches(7.1), Inches(2.2), Inches(5.7), Inches(0.5),
             size=18, bold=True, color=ACCENT_RED, font=FONT_CJK)
    add_code(s,
        "df = df[df[\"rt_ms\"] <\n"
        "        df[\"rt_ms\"].mean() + 3*df[\"rt_ms\"].std()]",
        Inches(7.1), Inches(2.85), Inches(5.65), Inches(1.3),
        size=13)
    add_text(s,
        "這條 outlier rule 直接影響 condition mean → 屬於 analysis decision，"
        "應放在 analyse() 並暴露閾值。",
        Inches(7.1), Inches(4.4), Inches(5.7), Inches(2.0),
        size=14, color=TEXT_MUTED, font=FONT_CJK)

    add_footer(s, page, total)
    return s


def slide_end_to_end_demo(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.4   ◆  END-TO-END DEMO")
    add_title(s, "完整走一遍 — Stroop-like RT 流程")

    code = (
        "# 1. load — 生成 messy raw data\n"
        "raw = pd.DataFrame({...})\n\n"
        "# 2. inspect — shape / dtypes\n"
        "# 3. describe — 找出問題\n"
        "raw.describe(include=\"all\")\n\n"
        "# 4. fix\n"
        "def clean_stroop(df):\n"
        "    df = df.copy()\n"
        "    df[\"rt_ms\"]     = pd.to_numeric(df[\"rt_ms\"], errors=\"coerce\")\n"
        "    df               = df[df[\"rt_ms\"].between(150, 3000)]\n"
        "    df[\"age\"]       = df[\"age\"].replace({-999: np.nan})\n"
        "    df[\"condition\"] = (df[\"condition\"].str.lower()\n"
        "                          .replace({\"incong\": \"incongruent\"}))\n"
        "    return df\n\n"
        "# 5. re-describe → 6. analyse\n"
        "clean.groupby(\"condition\")[\"rt_ms\"].agg([\"mean\",\"std\",\"count\"])"
    )
    add_code(s, code,
             Inches(0.6), Inches(1.95), Inches(7.6), Inches(4.9),
             size=12)

    # Observations panel
    add_card(s, Inches(8.4), Inches(1.95), Inches(4.5), Inches(4.9),
             fill=BG_OFFWHITE, accent=ACCENT_TEAL)
    add_text(s, "預期觀察",
             Inches(8.6), Inches(2.1), Inches(4.2), Inches(0.5),
             size=18, bold=True, color=ACCENT_TEAL, font=FONT_CJK)
    obs = [
        "condition 從看似 4 個 → 真正 2 個",
        "mean RT 從被 99999 拉高的數千 ms → 約 500 ms",
        "n 損失需在 report 聲明",
        "若要看 Stroop effect，生成資料需讓 incongruent +50–80 ms",
    ]
    add_bullets(s, obs,
                Inches(8.6), Inches(2.7), Inches(4.2), Inches(4.0),
                size=13, bullet_color=ACCENT_TEAL)

    add_demo_ref(s, "demo/03_pipeline/pipeline.py")
    add_footer(s, page, total)
    return s


def slide_cleaning_vs_analysis(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.5 BOUNDARY")
    add_title(s, "Cleaning vs Analysis 的邊界")

    code = (
        "def load_raw(path):              # I/O only\n"
        "    return pd.read_csv(path)\n\n"
        "def clean(df):                   # 不可爭議的修補\n"
        "    ...\n\n"
        "def describe(df):                # 健康檢查報表\n"
        "    ...\n\n"
        "def analyse(df, *, outlier_sd=3.0):   # 可爭議閾值 → parameter\n"
        "    ..."
    )
    add_code(s, code,
             Inches(0.9), Inches(2.0), Inches(7.8), Inches(4.6),
             size=14)

    add_card(s, Inches(8.9), Inches(2.0), Inches(3.95), Inches(4.6),
             fill=BG_OFFWHITE, accent=ACCENT_AMBER)
    add_text(s, "為什麼這樣分？",
             Inches(9.05), Inches(2.15), Inches(3.7), Inches(0.45),
             size=16, bold=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_bullets(s, [
        "clean() 可用 pytest 直接驗",
        "analyse() 的參數寫進 paper",
        "兩者混在一起 → 結論依賴沒被檢視的清理決定",
        "對 reproducibility / pre-registration 至關重要",
    ], Inches(9.05), Inches(2.7), Inches(3.7), Inches(4.0),
       size=13, bullet_color=ACCENT_AMBER)

    add_demo_ref(s, "demo/03_pipeline/pipeline.py + test_clean.py")
    add_footer(s, page, total)
    return s


def slide_litmus_question(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "PART 2 · 2.5 LITMUS TEST")
    add_title(s, "一個檢驗問題")

    add_callout(s,
        "「如果換另一位研究者，這個操作會做同一個決定嗎？」",
        Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6),
        size=24, color=ACCENT_PRIMARY)

    # Two branches
    branches = [
        ("✓ 是", "→  clean()  · 不可爭議",   ACCENT_GREEN),
        ("▲ 否", "→  analyse() · 把決定寫進 report", ACCENT_AMBER),
    ]
    y = Inches(4.7)
    for i, (head, body, color) in enumerate(branches):
        left = Inches(0.85) + i * Inches(6.2)
        add_card(s, left, y, Inches(5.9), Inches(1.5), accent=color)
        add_text(s, head, left + Inches(0.3), y + Inches(0.2),
                 Inches(2.5), Inches(0.5),
                 size=22, bold=True, color=color, font=FONT_CJK)
        add_text(s, body, left + Inches(0.3), y + Inches(0.8),
                 Inches(5.4), Inches(0.6),
                 size=17, color=TEXT_DARK, font=FONT_CJK)

    add_footer(s, page, total)
    return s


# -----------------------------------------------------------------
# Wrap-up
# -----------------------------------------------------------------
def slide_recap(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "RECAP")
    add_title(s, "一句話總結")

    add_callout(s,
        "@st.cache_data 讓 dashboard 跑得動；",
        Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.1),
        size=22, color=ACCENT_PRIMARY)
    add_callout(s,
        "descriptive statistics 讓 dashboard 背後的資料站得住。",
        Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.1),
        size=22, color=ACCENT_PRIMARY, accent=ACCENT_AMBER)
    add_callout(s,
        "任何 cleaning 動作都應對應到一個可觀察的現象，並能說出代價。",
        Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.1),
        size=22, color=ACCENT_PRIMARY, accent=ACCENT_GREEN)

    add_footer(s, page, total)
    return s


def slide_cheat_sheet(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "RECAP · CHEAT SHEET")
    add_title(s, "常見錯誤 cheat sheet")

    rows = [
        ("Slider 動了 dashboard 沒變",         "Widget 寫在 @st.cache_data 內 → 拿到函式外"),
        ("UnhashableParamError",                "在參數前加底線 _ 跳過 hash"),
        ("describe() 看不到關鍵問題",           "先 pd.to_numeric(errors=\"coerce\") 再 describe"),
        ("value_counts() 沒顯示缺值",           "永遠用 value_counts(dropna=False)"),
        ("結論隨 cleaning 一改就變",            "把可爭議閾值提到 analyse() 並暴露參數"),
        ("to_datetime 整欄變 NaT",              "format 不對 — 用 errors=\"coerce\" 並印幾筆原值"),
        ("st.set_page_config 報錯",             "把它移到所有 st.* 之前"),
    ]
    add_table(s, ["症狀", "修法"], rows,
              Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.7),
              col_widths=[Inches(4.5), Inches(7.4)],
              header_size=14, body_size=12, first_col_bold=True)

    add_footer(s, page, total)
    return s


def slide_homework(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "HOMEWORK · DUE 2026-05-20")
    add_title(s, "用 pipeline 處理 messy Stroop-like dataset")

    items = [
        ("1", "pipeline.py",
         "load_raw / describe / clean / analyse — 四個 pure functions"),
        ("2", "report.ipynb",
         "前後對照 + 觀察→修補對應 + 為何 outlier 放 analyse"),
        ("3", "tests/test_clean.py",
         "至少一個 pytest — 餵 sentinel 或 \"NA\" 驗證輸出"),
    ]
    y0 = Inches(2.0)
    for i, (num, head, body) in enumerate(items):
        y = y0 + i * Inches(1.1)
        add_text(s, num, Inches(0.9), y + Inches(0.1),
                 Inches(0.6), Inches(0.6),
                 size=28, bold=True, color=ACCENT_TEAL)
        add_text(s, head, Inches(1.7), y + Inches(0.1),
                 Inches(4.0), Inches(0.5),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_MONO)
        add_text(s, body, Inches(5.9), y + Inches(0.15),
                 Inches(7.0), Inches(1.0),
                 size=15, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s, "詳見講義 §Homework 的 Rubric。",
             Inches(0.9), Inches(5.45), Inches(11.5), Inches(0.4),
             size=14, italic=True, color=TEXT_MUTED, font=FONT_CJK)
    add_text(s, "可參考：demo/03_pipeline/ 的 pipeline.py + test_clean.py",
             Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.4),
             size=13, italic=True, color=ACCENT_TEAL, font=FONT_MONO)

    add_footer(s, page, total)
    return s


def slide_whats_next(prs, page, total):
    s = add_slide_base(prs)
    add_section_label(s, "WHAT COMES NEXT")
    add_title(s, "後續週次")

    nexts = [
        ("13", "Plotly Express", "Interactive dashboards & 資料敘事"),
        ("14", "Anthropic SDK",  "在 Streamlit app 中呼叫 Claude API"),
        ("15", "Workshop",       "Final project 同儕互評 & UI polish"),
        ("16", "Final",          "Live app 簡報"),
    ]
    x0, y0 = Inches(0.85), Inches(2.2)
    w, h = Inches(3.0), Inches(2.5)
    gap = Inches(0.15)
    for i, (num, head, body) in enumerate(nexts):
        left = x0 + i * (w + gap)
        add_card(s, left, y0, w, h, accent=ACCENT_TEAL)
        add_text(s, f"WEEK {num}",
                 left + Inches(0.2), y0 + Inches(0.2),
                 w - Inches(0.4), Inches(0.4),
                 size=13, bold=True, color=ACCENT_TEAL)
        add_text(s, head,
                 left + Inches(0.2), y0 + Inches(0.65),
                 w - Inches(0.4), Inches(0.6),
                 size=20, bold=True, color=TEXT_DARK, font=FONT_CJK)
        add_text(s, body,
                 left + Inches(0.2), y0 + Inches(1.35),
                 w - Inches(0.4), h - Inches(1.5),
                 size=13, color=TEXT_MUTED, font=FONT_CJK)

    add_text(s, "註：Open data API 將在 final project workshop 按需引入。",
             Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.5),
             size=15, italic=True, color=ACCENT_AMBER, font=FONT_CJK)
    add_text(s, "Thank you  ·  Questions?",
             Inches(0), Inches(6.05), Inches(13.333), Inches(0.7),
             size=24, bold=True, color=ACCENT_PRIMARY, font=FONT_CJK,
             align=PP_ALIGN.CENTER)

    add_footer(s, page, total)
    return s


# ============================================================
# Main — assemble deck
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 43

    # Section 0 — Opening
    slide_title(prs)                                               # 1
    slide_motivation(prs, 2, TOTAL)                                # 2
    slide_objectives_1(prs, 3, TOTAL)                              # 3
    slide_objectives_2(prs, 4, TOTAL)                              # 4
    slide_schedule(prs, 5, TOTAL)                                  # 5

    # Section 1 — Part 1
    slide_divider(prs, 6, TOTAL, 1,
                  "Streamlit App & Caching",
                  "app.py × @st.cache_data", 60)                   # 6
    slide_app_tour(prs, 7, TOTAL)                                  # 7
    slide_seven_blocks(prs, 8, TOTAL)                              # 8
    slide_sidebar(prs, 9, TOTAL)                                   # 9
    slide_practice(prs, 10, TOTAL, "Hands-on 1",
                   "把 app 跑起來、改一個參數",
                   tasks=[
                       "streamlit run app.py",
                       "把 \"Reaction Time (ms)\" 改成 \"Simple RT (ms)\"，存檔",
                       "觀察「Source file changed. Rerun.」提示",
                   ],
                   reflect="為什麼整支程式重新執行？對 load_data() 開銷意味著什麼？",
                   demo="demo/00_week11_app/app.py")  # 10
    slide_execution_model(prs, 11, TOTAL)                          # 11
    slide_rerun_consequences(prs, 12, TOTAL)                       # 12
    slide_cache_data_core(prs, 13, TOTAL)                          # 13
    slide_cache_key(prs, 14, TOTAL)                                # 14
    slide_cache_copy(prs, 15, TOTAL)                               # 15
    slide_cache_params(prs, 16, TOTAL)                             # 16
    slide_pitfall_1(prs, 17, TOTAL)                                # 17
    slide_pitfall_23(prs, 18, TOTAL)                               # 18
    slide_practice(prs, 19, TOTAL, "Hands-on 2",
                   "觀察 cache 行為",
                   tasks=[
                       "在 load_data() 內加 print(\"DEBUG: reading CSV\")",
                       "情境 A：有 @st.cache_data → 只印一次",
                       "情境 B：拿掉裝飾器 → 每次 slider 都印",
                       "情境 C：@st.cache_data(ttl=10) → 10 秒後再觸發重讀",
                   ],
                   demo="demo/01_cache_data/cache_demo.py")        # 19
    slide_cache_vs_resource(prs, 20, TOTAL)                        # 20
    slide_decision_rule(prs, 21, TOTAL)                            # 21

    # Break
    slide_break(prs, 22, TOTAL)                                    # 22

    # Section 2 — Part 2
    slide_divider(prs, 23, TOTAL, 2,
                  "The Data Analysis Pipeline",
                  "descriptive statistics × observation-driven fixing",
                  90)                                              # 23
    slide_why_pipeline(prs, 24, TOTAL)                             # 24
    slide_pipeline_diagram(prs, 25, TOTAL)                         # 25
    slide_two_traps(prs, 26, TOTAL)                                # 26
    slide_practice(prs, 27, TOTAL, "Hands-on 3",
                   "替 app.py 標出 pipeline 位置",
                   tasks=[
                       "1) 哪一行對應 load？",
                       "2) 哪一段對應 analyse / visualise？",
                       "3) 缺少哪幾步？對合成資料 vs 真實資料合理嗎？",
                   ],
                   answer="load_data() 在第 56 行；tabs 內畫圖；缺 inspect / describe / fix / re-describe")  # 27
    slide_descriptive_stats(prs, 28, TOTAL)                        # 28
    slide_diagnose_table(prs, 29, TOTAL)                           # 29
    slide_viz_diagnostic(prs, 30, TOTAL)                           # 30
    slide_practice(prs, 31, TOTAL, "Hands-on 4",
                   "從 summary 讀出三個問題",
                   tasks=[
                       "情境：n=200 Stroop-like dataset（刻意 messy）",
                       "工具：info() / describe(include=\"all\") / isnull().sum() / value_counts(dropna=False)",
                       "任務：列出至少三個資料品質問題",
                   ],
                   answer="rt_ms object dtype；max=99999 sentinel；condition 大小寫不一致；age 含 -999",
                   demo="demo/data/messy_stroop.csv")  # 31
    slide_obs_driven_principle(prs, 32, TOTAL)                     # 32
    slide_obs_action_cost(prs, 33, TOTAL)                          # 33
    slide_cleaning_discipline(prs, 34, TOTAL)                      # 34
    slide_error_modes(prs, 35, TOTAL)                              # 35
    slide_practice(prs, 36, TOTAL, "Hands-on 5",
                   "把觀察寫成 clean_stroop(df)",
                   tasks=[
                       "根據 Practice 4 觀察到的問題寫 cleaning function",
                       "註解必須寫出「觀察 → 動作 → 代價」",
                       "cleaning 前後跑 describe(include=\"all\") 對照",
                   ],
                   answer="4 步：to_numeric / between(150,3000) / replace(-999) / lower+replace",
                   demo="demo/03_pipeline/pipeline.py::clean")  # 36
    slide_end_to_end_demo(prs, 37, TOTAL)                          # 37
    slide_cleaning_vs_analysis(prs, 38, TOTAL)                     # 38
    slide_litmus_question(prs, 39, TOTAL)                          # 39

    # Section 3 — Wrap-up
    slide_recap(prs, 40, TOTAL)                                    # 40
    slide_cheat_sheet(prs, 41, TOTAL)                              # 41
    slide_homework(prs, 42, TOTAL)                                 # 42
    slide_whats_next(prs, 43, TOTAL)                               # 43

    out = "week-12-slides.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
